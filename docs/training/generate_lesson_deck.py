#!/usr/bin/env python3
"""Generate the 120-minute VulnHunterX *lesson* deck as a .pptx file.

Run from anywhere:  python docs/training/generate_lesson_deck.py
Requires:           pip install python-pptx

This is the teaching deck that accompanies docs/training/LESSON.md — a broader,
more foundational talk than the product-intro 60-min deck in docs/presentation/:
SAST vs DAST, how CodeQL/Semgrep scan, AST + control/data flow (CodeQL vs
tree-sitter), LLM-verification pros/cons, then VulnHunterX architecture, CLI, and
results/limitations (static + LLM only; fuzzing out of scope).

The visual toolkit (theme, slide templates, code/flow helpers) is copied from
docs/presentation/generate_deck.py — that module builds its own deck at import
time, so we duplicate the helpers here instead of importing it. Rule/question
counts are recomputed from config/ at build time so the deck never drifts.

Section numbers map to LESSON.md sections; each content slide carries speaker
notes (the 120-minute talk track).
"""
from __future__ import annotations

import glob
import os
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# --------------------------------------------------------------------------- #
# Paths
# --------------------------------------------------------------------------- #
HERE = os.path.dirname(os.path.abspath(__file__))      # docs/training/
REPO = os.path.dirname(os.path.dirname(HERE))          # repo root
CONFIG = os.path.join(REPO, "config")
OUT_PATH = os.path.join(HERE, "VulnHunterX-120min.pptx")

# --------------------------------------------------------------------------- #
# Theme  — modern slate / indigo / cyan with orange brand accent
# --------------------------------------------------------------------------- #
# Dark theme — slate/indigo/cyan with orange brand accent.
# NOTE: the slide body references these names, so recoloring here recolors the
# whole deck. Title + section dividers use WHITE/GHOST/AMBER/ORANGE (below) and
# are unaffected.
DARK = RGBColor(0x0A, 0x0F, 0x1C)      # slide background base (darkest)
DARK2 = RGBColor(0x0E, 0x1A, 0x30)     # gradient end
GHOST = RGBColor(0x1C, 0x2C, 0x4A)     # watermark numerals on dividers
INK = RGBColor(0xD4, 0xDE, 0xEA)       # body text (light)
NAVY_TX = RGBColor(0xF2, 0xF6, 0xFB)   # headings (near-white)
GREY = RGBColor(0x9D, 0xB0, 0xC6)      # muted light sub-text
MUTE = RGBColor(0x90, 0x9C, 0xAB)      # footers / captions
BORDER = RGBColor(0x2C, 0x3E, 0x5E)    # subtle panel border on dark
BG = RGBColor(0x0A, 0x0F, 0x1C)        # content background
CARD = RGBColor(0x18, 0x23, 0x3D)      # cards / table body / flow nodes / panels
TRACK = RGBColor(0x1E, 0x2C, 0x4A)     # progress track

ORANGE = RGBColor(0xF9, 0x73, 0x16)
INDIGO = RGBColor(0x81, 0x84, 0xF7)    # flow accent (lifted for dark contrast)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
EMERALD = RGBColor(0x34, 0xD3, 0x99)
ROSE = RGBColor(0xFB, 0x71, 0x85)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
SLATE = RGBColor(0x7C, 0x8B, 0xA5)     # connectors / arrows (lighter for dark)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HEADER_BG = RGBColor(0x26, 0x35, 0x52)  # table header
ROW_ALT = RGBColor(0x13, 0x1D, 0x33)    # table banding

CODE_BG = RGBColor(0x0F, 0x1A, 0x2E)
CODE_FG = RGBColor(0xCC, 0xD8, 0xE6)
CODE_COMMENT = RGBColor(0x6E, 0x8B, 0xA3)
CODE_CMD = RGBColor(0x4F, 0xC9, 0xA8)
CODE_FLAG = RGBColor(0xF5, 0x9E, 0x0B)

BODY_FONT = "Segoe UI"
MONO_FONT = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

LEFT, CENTER, RIGHT = PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.RIGHT
TOP, MID = MSO_ANCHOR.TOP, MSO_ANCHOR.MIDDLE
RR, RECT, OVAL = MSO_SHAPE.ROUNDED_RECTANGLE, MSO_SHAPE.RECTANGLE, MSO_SHAPE.OVAL

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# --------------------------------------------------------------------------- #
# Build-time counters (computed from config/ so the deck never drifts)
# --------------------------------------------------------------------------- #
def count_codeql_queries() -> dict[str, int]:
    out: dict[str, int] = {}
    base = os.path.join(CONFIG, "codeql-custom")
    if not os.path.isdir(base):
        return out
    for lang in sorted(os.listdir(base)):
        d = os.path.join(base, lang)
        if not os.path.isdir(d):
            continue
        n = len(glob.glob(os.path.join(d, "**", "*.ql"), recursive=True))
        if n:
            out[lang] = n
    return out


def count_semgrep_rules() -> dict[str, int]:
    out: dict[str, int] = {}
    for path in sorted(glob.glob(os.path.join(CONFIG, "semgrep-custom", "*.yaml"))):
        lang = os.path.splitext(os.path.basename(path))[0]
        with open(path, encoding="utf-8") as fh:
            n = sum(1 for line in fh if re.match(r"\s*-\s+id:", line))
        out[lang] = n
    return out


def count_guided_questions() -> dict[str, int]:
    out: dict[str, int] = {}
    for path in sorted(glob.glob(os.path.join(CONFIG, "prompts", "*questions*.yaml"))):
        name = os.path.basename(path)
        with open(path, encoding="utf-8") as fh:
            n = sum(1 for line in fh if re.match(r"^[^\s#][^:]*:\s*$", line.rstrip()))
        out[name] = n
    return out


CODEQL = count_codeql_queries()
SEMGREP = count_semgrep_rules()
QUESTIONS = count_guided_questions()
CODEQL_TOTAL = sum(CODEQL.values())
SEMGREP_TOTAL = sum(SEMGREP.values())
QUESTIONS_TOTAL = sum(QUESTIONS.values())
CODEQL_CPP = CODEQL.get("cpp", 0) + CODEQL.get("c", 0)
SEMGREP_CPP = SEMGREP.get("cpp", 0) + SEMGREP.get("c", 0)


# --------------------------------------------------------------------------- #
# Low-level helpers
# --------------------------------------------------------------------------- #
def _bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def _set(run, size, color, bold=False, italic=False, font=BODY_FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def shape(slide, mso, l, t, w, h, fill=None, line=None, line_w=1.0):
    sp = slide.shapes.add_shape(mso, int(l), int(t), int(w), int(h))
    if mso == RR:
        # tame the default rounded-rectangle radius (subtle, modern corners)
        try:
            sp.adjustments[0] = 0.045
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def text_in(sp, paras, align=CENTER, anchor=MID, wrap=True,
            ml=0.10, mt=0.05, mr=0.10, mb=0.05):
    tf = sp.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(ml)
    tf.margin_right = Inches(mr)
    tf.margin_top = Inches(mt)
    tf.margin_bottom = Inches(mb)
    first = True
    for para in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(2)
        for txt, size, color, bold, italic in para:
            r = p.add_run()
            r.text = txt
            _set(r, size, color, bold=bold, italic=italic)
    return sp


def textbox(slide, l, t, w, h, paras, align=LEFT, anchor=TOP, wrap=True):
    tb = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
    return text_in(tb, paras, align=align, anchor=anchor, wrap=wrap,
                   ml=0, mt=0, mr=0, mb=0)


def soft_shadow(sp, alpha=22000, blur=80000, dist=38000):
    spPr = sp._element.spPr
    old = spPr.find(qn("a:effectLst"))
    if old is not None:
        spPr.remove(old)
    el = spPr.makeelement(qn("a:effectLst"), {})
    sh = el.makeelement(qn("a:outerShdw"),
                        {"blurRad": str(blur), "dist": str(dist),
                         "dir": "5400000", "rotWithShape": "0"})
    clr = sh.makeelement(qn("a:srgbClr"), {"val": "0F172A"})
    clr.append(clr.makeelement(qn("a:alpha"), {"val": str(alpha)}))
    sh.append(clr)
    el.append(sh)
    spPr.append(el)
    return sp


def fill_alpha(sp, pct):
    fill = sp._element.spPr.find(qn("a:solidFill"))
    clr = fill.find(qn("a:srgbClr"))
    clr.append(clr.makeelement(qn("a:alpha"), {"val": str(int(pct * 1000))}))
    return sp


def grad(sp, c1, c2, angle=45):
    sp.fill.gradient()
    gs = sp.fill.gradient_stops
    gs[0].position, gs[0].color.rgb = 0.0, c1
    gs[1].position, gs[1].color.rgb = 1.0, c2
    try:
        sp.fill.gradient_angle = angle
    except Exception:
        pass
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def connect(slide, x1, y1, x2, y2, color, width=2.0, elbow=False, arrow=True):
    typ = MSO_CONNECTOR.ELBOW if elbow else MSO_CONNECTOR.STRAIGHT
    cn = slide.shapes.add_connector(typ, int(x1), int(y1), int(x2), int(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    cn.shadow.inherit = False
    if arrow:
        ln = cn.line._get_or_add_ln()
        ln.append(ln.makeelement(qn("a:tailEnd"),
                                 {"type": "triangle", "w": "med", "len": "med"}))
    return cn


def pill(slide, l, t, w, h, label, fill, color=WHITE, size=11):
    p = shape(slide, RR, l, t, w, h, fill=fill)
    text_in(p, [[(label, size, color, True, False)]])
    return p


# --------------------------------------------------------------------------- #
# Slide templates
# --------------------------------------------------------------------------- #
PART_TOTAL = 7


def title_slide(footer):
    s = prs.slides.add_slide(BLANK)
    grad(shape(s, RECT, 0, 0, SLIDE_W, SLIDE_H, fill=DARK), DARK, DARK2, angle=60)
    fill_alpha(shape(s, OVAL, Inches(9.4), Inches(-1.8), Inches(6.4), Inches(6.4),
                     fill=INDIGO), 14)
    fill_alpha(shape(s, OVAL, Inches(10.8), Inches(3.6), Inches(4.6), Inches(4.6),
                     fill=CYAN), 12)
    text_in(soft_shadow(shape(s, RR, Inches(0.9), Inches(0.85), Inches(0.95),
                              Inches(0.95), fill=ORANGE)),
            [[("VX", 30, WHITE, True, False)]])
    textbox(s, Inches(2.0), Inches(1.02), Inches(7), Inches(0.7),
            [[("VulnHunterX", 22, WHITE, True, False)]])
    shape(s, RECT, Inches(0.92), Inches(2.7), Inches(0.55), Inches(0.12), fill=ORANGE)
    textbox(s, Inches(1.62), Inches(2.46), Inches(11), Inches(0.5),
            [[("90-MINUTE LESSON  ·  SELF-SCAN YOUR SOURCE CODE", 15, AMBER, True, False)]])
    textbox(s, Inches(0.85), Inches(3.0), Inches(11.8), Inches(1.5),
            [[("Find real bugs.", 54, WHITE, True, False)],
             [("Skip the false alarms.", 54, CYAN, True, False)]])
    textbox(s, Inches(0.92), Inches(5.0), Inches(11.4), Inches(0.6),
            [[("Static analysis finds candidates — an LLM decides which are real.",
               20, MUTE, False, False)]])
    chips = [(f"{CODEQL_TOTAL}", "custom CodeQL queries"),
             (f"{SEMGREP_TOTAL}", "custom Semgrep rules"),
             ("8", "languages  ·  5 profiles")]
    x = Inches(0.9)
    for big, small in chips:
        c = soft_shadow(shape(s, RR, x, Inches(5.75), Inches(3.55), Inches(0.7),
                              fill=GHOST))
        text_in(c, [[(big + "  ", 18, ORANGE, True, False),
                     (small, 12.5, WHITE, False, False)]], align=CENTER)
        x += Inches(3.75)
    textbox(s, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.5),
            [[(footer, 13, MUTE, False, False)]])
    return s


def section_slide(part_idx, kicker, title):
    s = prs.slides.add_slide(BLANK)
    grad(shape(s, RECT, 0, 0, SLIDE_W, SLIDE_H, fill=DARK), DARK, DARK2, angle=30)
    textbox(s, Inches(7.4), Inches(1.05), Inches(5.55), Inches(4.9),
            [[(str(part_idx).zfill(2), 250, GHOST, True, False)]], align=RIGHT,
            anchor=MID, wrap=False)
    shape(s, RECT, Inches(0.95), Inches(2.92), Inches(0.55), Inches(0.12), fill=ORANGE)
    textbox(s, Inches(1.65), Inches(2.68), Inches(9), Inches(0.55),
            [[(kicker.upper(), 18, AMBER, True, False)]])
    textbox(s, Inches(0.9), Inches(3.22), Inches(11), Inches(1.8),
            [[(title, 40, WHITE, True, False)]])
    shape(s, RECT, Inches(0.95), Inches(4.75), Inches(3.2), Inches(0.05), fill=ORANGE)
    x = Inches(0.95)
    for i in range(1, PART_TOTAL + 1):
        col = ORANGE if i == part_idx else GHOST
        w = Inches(0.5) if i == part_idx else Inches(0.22)
        shape(s, RR, x, Inches(6.7), w, Inches(0.16), fill=col)
        x += w + Inches(0.16)
    return s


def content_slide(title):
    s = prs.slides.add_slide(BLANK)
    _bg(s, BG)
    grad(shape(s, RECT, 0, 0, SLIDE_W, SLIDE_H, fill=DARK), DARK, DARK2, angle=70)
    shape(s, RECT, 0, 0, Inches(0.14), SLIDE_H, fill=ORANGE)
    textbox(s, Inches(0.55), Inches(0.34), Inches(11.0), Inches(0.85),
            [[(title, 26, NAVY_TX, True, False)]])
    shape(s, RECT, Inches(0.6), Inches(1.14), Inches(1.5), Inches(0.045), fill=ORANGE)
    textbox(s, Inches(9.7), Inches(0.42), Inches(3.05), Inches(0.4),
            [[("VulnHunterX", 12, MUTE, True, False)]], align=RIGHT)
    textbox(s, Inches(12.0), Inches(7.02), Inches(1.0), Inches(0.4),
            [[(str(len(prs.slides)), 11, MUTE, False, False)]], align=RIGHT)
    return s


def bullets_slide(title, bullets, notes="", top=Inches(1.5), left=Inches(0.7),
                  width=Inches(12.0), height=Inches(5.4)):
    s = content_slide(title)
    tb = s.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for text, level, color, bold in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(9)
        size = 21 - 3 * level
        if text:
            if level == 0 and bold:
                m = p.add_run()
                m.text = "▸  "
                _set(m, size, ORANGE, bold=True)
            elif level >= 1:
                m = p.add_run()
                m.text = "–  "
                _set(m, size, MUTE, bold=True)
        r = p.add_run()
        r.text = text
        _set(r, size, color or INK, bold=bold)
    if notes:
        _notes(s, notes)
    return s


def lay_cards(s, items, cols, top=1.55, card_h=1.12, gap=0.25, left=0.6,
              total_w=12.2, head_size=15.5, body_size=12.5):
    gap_e = Inches(gap)
    card_w = int((Inches(total_w) - gap_e * (cols - 1)) / cols)
    ch = Inches(card_h)
    for i, it in enumerate(items):
        r, c = divmod(i, cols)
        x = int(Inches(left) + c * (card_w + gap_e))
        y = int(Inches(top) + r * (ch + gap_e))
        accent = it.get("accent", ORANGE)
        soft_shadow(shape(s, RR, x, y, card_w, ch, fill=CARD, line=BORDER))
        shape(s, RR, x, y, Inches(0.11), ch, fill=accent)
        paras = [[(it["head"], head_size, NAVY_TX, True, False)]]
        if it.get("body"):
            paras.append([(it["body"], body_size, GREY, False, False)])
        body_box = shape(s, RECT, x + Inches(0.22), y, card_w - Inches(0.34), ch)
        body_box.fill.background()
        body_box.line.fill.background()
        text_in(body_box, paras, align=LEFT, anchor=MID, ml=0.06, mr=0.06)
    return s


def card_grid_slide(title, items, cols, notes="", **kw):
    s = content_slide(title)
    lay_cards(s, items, cols, **kw)
    if notes:
        _notes(s, notes)
    return s


def table_slide(title, headers, rows, col_widths, notes="", caption="",
                head_size=13, body_size=12.5):
    """Native PowerPoint table. col_widths in inches; should sum to ~11.95."""
    s = content_slide(title)
    top = Inches(1.55)
    if caption:
        textbox(s, Inches(0.7), Inches(1.24), Inches(12), Inches(0.4),
                [[(caption, 14, MUTE, False, True)]])
        top = Inches(1.78)
    nrows, ncols = len(rows) + 1, len(headers)
    left = Inches(0.7)
    width = Inches(sum(col_widths))
    height = Inches(0.42 * nrows)
    gframe = s.shapes.add_table(nrows, ncols, int(left), int(top),
                                int(width), int(height))
    table = gframe.table
    table.first_row = False
    table.horz_banding = False
    for j, w in enumerate(col_widths):
        table.columns[j].width = Inches(w)

    def style_cell(cell, text, size, color, bold, fill, align=LEFT):
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)
        cell.vertical_anchor = MID
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        _set(r, size, color, bold=bold)

    for j, h in enumerate(headers):
        style_cell(table.cell(0, j), h, head_size, WHITE, True, HEADER_BG)
    for i, row in enumerate(rows):
        fill = CARD if i % 2 == 0 else ROW_ALT
        for j, val in enumerate(row):
            bold = (j == 0)
            color = NAVY_TX if bold else INK
            style_cell(table.cell(i + 1, j), val, body_size, color, bold, fill)
    if notes:
        _notes(s, notes)
    return s


def _emit_code_line(p, line):
    if not line.strip():
        r = p.add_run()
        r.text = " "
        _set(r, 15.5, CODE_FG, font=MONO_FONT)
        return
    if line.lstrip().startswith("#"):
        r = p.add_run()
        r.text = line
        _set(r, 15.5, CODE_COMMENT, italic=True, font=MONO_FONT)
        return
    code_part, comment_part = line, ""
    idx = line.find("  #")
    if idx != -1:
        code_part, comment_part = line[:idx], line[idx:]
    for tok in re.split(r"(\s+)", code_part):
        if tok == "":
            continue
        r = p.add_run()
        r.text = tok
        if tok.startswith("--") or (len(tok) > 1 and tok[0] == "-" and tok[1].isalpha()):
            _set(r, 15.5, CODE_FLAG, font=MONO_FONT)
        elif tok in ("vuln-hunter-x", "python", "git", "uv", "cp", "cat", "source"):
            _set(r, 15.5, CODE_CMD, bold=True, font=MONO_FONT)
        else:
            _set(r, 15.5, CODE_FG, font=MONO_FONT)
    if comment_part:
        r = p.add_run()
        r.text = comment_part
        _set(r, 15.5, CODE_COMMENT, italic=True, font=MONO_FONT)


def code_slide(title, lines, notes="", caption="", size=15.5):
    s = content_slide(title)
    top = Inches(1.5)
    if caption:
        textbox(s, Inches(0.7), Inches(1.26), Inches(12), Inches(0.5),
                [[(caption, 15, MUTE, False, True)]])
        top = Inches(1.78)
    box = soft_shadow(shape(s, RR, Inches(0.7), top, Inches(11.95),
                            int(SLIDE_H - top - Inches(0.55)), fill=CODE_BG,
                            line=BORDER, line_w=1.0))
    for i, col in enumerate((ROSE, AMBER, EMERALD)):
        shape(s, OVAL, Inches(0.95 + i * 0.28), top + Inches(0.22),
              Inches(0.16), Inches(0.16), fill=col)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = TOP
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.62)
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(3)
        # size override vs the default 15.5 in _emit_code_line
        if not line.strip():
            run = p.add_run(); run.text = " "; _set(run, size, CODE_FG, font=MONO_FONT)
        elif line.lstrip().startswith("#"):
            run = p.add_run(); run.text = line
            _set(run, size, CODE_COMMENT, italic=True, font=MONO_FONT)
        else:
            _emit_code_line(p, line)
    if notes:
        _notes(s, notes)
    return s


# --------------------------------------------------------------------------- #
# Flow / diagram helpers
# --------------------------------------------------------------------------- #
def flow_node(s, l, t, w, h, head, sub, accent, head_size=13.5, sub_size=10.5):
    soft_shadow(shape(s, RR, int(l), int(t), int(w), int(h), fill=CARD, line=BORDER))
    shape(s, RR, int(l), int(t), int(w), Inches(0.12), fill=accent)
    box = shape(s, RECT, int(l), int(t + Inches(0.16)), int(w), int(h - Inches(0.16)))
    box.fill.background()
    box.line.fill.background()
    paras = [[(head, head_size, NAVY_TX, True, False)]]
    if sub:
        for ln in str(sub).split("\n"):
            paras.append([(ln, sub_size, GREY, False, False)])
    text_in(box, paras, align=CENTER, anchor=MID, ml=0.04, mr=0.04)
    return box


def flow_chain(s, nodes, top, left=0.7, node_w=2.2, node_h=1.1, gap=0.55,
               accent=INDIGO, vertical=False, head_size=13.5, sub_size=10.5):
    """nodes: list of (head, sub). Draws nodes with arrows between them."""
    xs, ys = [], []
    L, T = Inches(left), Inches(top)
    nw, nh, g = Inches(node_w), Inches(node_h), Inches(gap)
    for i, (head, sub) in enumerate(nodes):
        if vertical:
            x, y = int(L), int(T + i * (nh + g))
        else:
            x, y = int(L + i * (nw + g)), int(T)
        flow_node(s, x, y, nw, nh, head, sub, accent,
                  head_size=head_size, sub_size=sub_size)
        xs.append(x)
        ys.append(y)
    for i in range(len(nodes) - 1):
        if vertical:
            connect(s, int(xs[i] + nw / 2), int(ys[i] + nh),
                    int(xs[i] + nw / 2), int(ys[i + 1]), accent, width=2.5)
        else:
            connect(s, int(xs[i] + nw), int(ys[i] + nh / 2),
                    int(xs[i + 1]), int(ys[i] + nh / 2), accent, width=2.5)
    return xs, ys


def taint_diagram(s, top=2.55):
    """SOURCE → SANITIZER → SINK with a red 'no sanitizer ⇒ vuln' bypass."""
    nodes = [("SOURCE", "untrusted input\nargv · request · file", ROSE),
             ("SANITIZER", "validate / escape\nbounds-check", EMERALD),
             ("SINK", "dangerous op\nsystem() · memcpy()", AMBER)]
    nw, nh, gap = Inches(3.0), Inches(1.3), Inches(1.15)
    L, cy = Inches(1.0), Inches(top)
    xs = []
    for i, (head, sub, acc) in enumerate(nodes):
        x = int(L + i * (nw + gap))
        flow_node(s, x, int(cy), nw, nh, head, sub, acc, head_size=17, sub_size=11.5)
        xs.append(x)
    for i in range(2):
        connect(s, int(xs[i] + nw), int(cy + nh / 2),
                int(xs[i + 1]), int(cy + nh / 2), SLATE, width=3)
    # red bypass arrow beneath the row: source → sink, skipping the sanitizer
    by = int(cy + nh + Inches(0.55))
    connect(s, int(xs[0] + nw / 2), by, int(xs[2] + nw / 2), by, ROSE, width=3)
    textbox(s, int(xs[0]), int(by + Inches(0.12)), int(nw * 3 + gap * 2), Inches(0.5),
            [[("missing / inadequate sanitizer  ⇒  SOURCE reaches SINK  =  vulnerability",
               14, ROSE, True, False)]], align=CENTER)
    textbox(s, int(xs[1] - Inches(0.2)), int(cy - Inches(0.5)),
            int(nw + Inches(0.4)), Inches(0.4),
            [[("adequate check here  ⇒  safe", 12, EMERALD, True, True)]], align=CENTER)


def ast_tree(s, top=2.0):
    """Hardcoded AST for: x = a + b * 2."""
    def node(cx, cy, label, accent, w=1.15, h=0.62):
        x, y = int(Inches(cx) - Inches(w) / 2), int(Inches(cy))
        soft_shadow(shape(s, RR, x, y, Inches(w), Inches(h), fill=CARD, line=accent, line_w=1.6))
        text_in(s.shapes[-1], [[(label, 15, NAVY_TX, True, False)]])
        return (Inches(cx), Inches(cy), Inches(h))

    def edge(a, b):
        connect(s, int(a[0]), int(a[1] + a[2]), int(b[0]), int(b[1]), SLATE,
                width=2.0, arrow=False)

    assign = node(4.3, top, "=  (assign)", INDIGO)
    xvar = node(2.3, top + 1.2, "x", CYAN)
    plus = node(6.3, top + 1.2, "+", INDIGO)
    avar = node(4.9, top + 2.4, "a", CYAN)
    mul = node(7.7, top + 2.4, "*", INDIGO)
    bvar = node(6.7, top + 3.6, "b", CYAN)
    two = node(8.7, top + 3.6, "2", CYAN)
    for a, b in [(assign, xvar), (assign, plus), (plus, avar), (plus, mul),
                 (mul, bvar), (mul, two)]:
        edge(a, b)
    textbox(s, Inches(9.6), Inches(top + 0.4), Inches(3.4), Inches(3.0),
            [[("Source:", 14, MUTE, True, False)],
             [("x = a + b * 2", 16, ORANGE, True, False)],
             [("", 8, INK, False, False)],
             [("Precedence is baked", 12.5, GREY, False, False)],
             [("into the tree: * binds", 12.5, GREY, False, False)],
             [("tighter than +, so it", 12.5, GREY, False, False)],
             [("sits deeper. Tools match", 12.5, GREY, False, False)],
             [("on this shape, not text.", 12.5, GREY, False, False)]], align=LEFT)


def two_col_compare(s, left_title, left_items, right_title, right_items,
                    left_accent=CYAN, right_accent=INDIGO, top=1.6, foot=None):
    panel_w, gap, h = Inches(5.95), Inches(0.3), Inches(4.55)
    L1 = Inches(0.7)
    L2 = int(L1 + panel_w + gap)
    for Lx, title, items, acc in [(int(L1), left_title, left_items, left_accent),
                                  (L2, right_title, right_items, right_accent)]:
        soft_shadow(shape(s, RR, Lx, int(Inches(top)), int(panel_w), int(h),
                          fill=CARD, line=BORDER))
        hdr = shape(s, RR, Lx, int(Inches(top)), int(panel_w), Inches(0.62), fill=acc)
        text_in(hdr, [[(title, 16, WHITE, True, False)]])
        tb = s.shapes.add_textbox(int(Lx + Inches(0.24)),
                                  int(Inches(top) + Inches(0.82)),
                                  int(panel_w - Inches(0.48)),
                                  int(h - Inches(1.0)))
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for it in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(7)
            mono = it.startswith("»")
            txt = it[1:].strip() if mono else it
            if not mono:
                m = p.add_run()
                m.text = "•  "
                _set(m, 14, acc, bold=True)
            r = p.add_run()
            r.text = txt
            _set(r, 12.5 if mono else 14, GREY if mono else INK,
                 italic=mono, font=MONO_FONT if mono else BODY_FONT)
    if foot:
        textbox(s, Inches(0.7), int(Inches(top) + h + Inches(0.12)), Inches(12),
                Inches(0.5), [[(foot, 13.5, MUTE, False, True)]], align=CENTER)


# --------------------------------------------------------------------------- #
# SLIDES
# --------------------------------------------------------------------------- #

# 0 — Title + agenda --------------------------------------------------------- #
title_slide("120-min lesson + 30-min workshop  ·  for developers  ·  SAST + LLM triage")

bullets_slide(
    "Agenda — seven parts",
    [("Source-code security scanning: SAST & DAST", 0, INK, True),
     ("How CodeQL & Semgrep actually find bugs", 0, INK, True),
     ("AST, control flow & data flow (CodeQL vs tree-sitter)", 0, INK, True),
     ("LLM vulnerability verification — pros & cons", 0, INK, True),
     ("VulnHunterX architecture & stages (no fuzzing)", 0, INK, True),
     ("Using it: CLI + a worked example", 0, INK, True),
     ("Results & limitations  →  then the hands-on workshop", 0, INK, True)],
    notes="Seven parts. The first four are vendor-neutral fundamentals you can "
          "take to any tool. The last three are VulnHunterX-specific and lead "
          "into the workshop, where you scan a real vulnerable C++ repo with a "
          "free Ollama Cloud model.")

# 1 — SAST & DAST ------------------------------------------------------------ #
section_slide(1, "Part 1", "Source-code scanning: SAST & DAST")

bullets_slide(
    "Why scan source code at all?",
    [("Vulnerabilities are cheapest to fix BEFORE they ship — shift left.", 0, INK, True),
     ("Manual review can't scale to every commit, repo, and language.", 0, INK, True),
     ("Scanners give breadth; humans give judgement — we want both.", 0, INK, True),
     ("Core tension: catch everything ⇒ cry wolf constantly.", 0, ROSE, True)],
    notes="A bug caught in review costs minutes; in production after a breach, "
          "orders of magnitude more. But developers ignore tools that bury one "
          "real bug under fifty false alarms. VulnHunterX is a direct response "
          "to that alert-fatigue problem.")

table_slide(
    "SAST vs DAST in one picture",
    ["", "SAST (static)", "DAST (dynamic)"],
    [["Sees", "Code / bytecode at rest", "A running application"],
     ["Runs the app?", "No", "Yes"],
     ["Coverage", "Every path, even unreached", "Only paths exercised at runtime"],
     ["Finds", "Injection, memory bugs, config, secrets", "Auth flaws, real exploitability"],
     ["Weakness", "False positives", "False negatives (misses unreached code)"],
     ["In the SDLC", "Commit / PR / IDE", "Staging / QA / pre-prod"]],
    col_widths=[1.9, 5.0, 5.05],
    notes="SAST reads code like a pedantic reviewer — sees everything, can't "
          "prove reachability, over-reports. DAST throws real traffic at a live "
          "app — exploitable but only covers what it hit. Complementary. "
          "VulnHunterX is a SAST tool (its optional fuzzing stage adds DAST-like "
          "runtime confirmation for C/C++ — out of scope today).")

bullets_slide(
    "The SAST false-positive problem",
    [("A raw SAST run can emit hundreds to thousands of findings.", 0, INK, True),
     ("Many are unreachable, already mitigated, or the wrong CWE class.", 0, INK, True),
     ("Hand triage is slow and demoralising → tools get abandoned.", 0, INK, True),
     ("The opportunity: automate the triage REASONING, not just detection.", 0, EMERALD, True),
     ("Remember: detection finds candidates, verification decides.", 1, GREY, False)],
    notes="Detection is mostly solved — CodeQL and Semgrep find candidates well. "
          "The expensive, unsolved part is deciding which are real. VulnHunterX "
          "inserts an LLM as a first-pass triager so humans review only what "
          "survived scrutiny. Hold this idea for the rest of the talk.")

# 2 — CodeQL & Semgrep ------------------------------------------------------- #
section_slide(2, "Part 2", "How CodeQL & Semgrep scan")

_s = content_slide("Two philosophies of static analysis")
two_col_compare(
    _s,
    "Pattern matching  ·  Semgrep",
    ["“Does this code LOOK LIKE a known bad shape?”",
     "Parses each file to an AST, matches a written pattern.",
     "File-local · fast · no build required.",
     "Pipeline:",
     "» source file → AST → pattern match → finding"],
    "Semantic / data-flow  ·  CodeQL",
    ["“Can tainted data actually REACH a dangerous sink?”",
     "Builds a queryable model of the whole program.",
     "Cross-function · type-aware · needs a build.",
     "Pipeline:",
     "» build → relational DB → taint query → finding"],
    left_accent=CYAN, right_accent=INDIGO,
    foot="Different questions, not rivals — VulnHunterX runs both and reconciles; "
         "both emit findings as SARIF.")
_notes(_s,
       "The key reframing: these aren't two competing tools, they're two "
       "different QUESTIONS. Pattern matching asks 'does this look wrong?' — fast "
       "and syntactic. Semantic analysis asks 'can bad data actually get here?' — "
       "deep and cross-function. The old slide had a confusing third pillar; the "
       "shared SARIF output is just the bridge, shown in the footer. VulnHunterX "
       "runs both and reconciles their verdicts.")

bullets_slide(
    "CodeQL: code becomes a database",
    [("Build the code under a tracer → a relational DB of the program.", 0, INK, True),
     ("(functions, expressions, types, control & data flow)", 1, GREY, False),
     ("Query it with .ql queries — an SQL-like, OO language.", 0, INK, True),
     ("Built-in suites: security-extended, security-and-quality.", 0, INK, True),
     ("Output: SARIF with ruleId, location, severity, data-flow path.", 0, INK, True),
     ("Needs to observe the build for compiled langs (C/C++/Java/Go/C#).", 0, ROSE, True)],
    notes="Key fact: CodeQL needs to observe the build for compiled languages — "
          "a broken build collapses coverage (remember for limitations). Once "
          "the DB exists, a query like 'flow from argv to system()' runs across "
          "function boundaries. That cross-function taint is what pattern tools "
          "can't do.")

_s = content_slide("CodeQL: the extraction → query pipeline")
flow_chain(
    _s,
    [("Source code", "your repo"),
     ("Build tracer", "observes compilation"),
     ("Relational DB", "AST + CFG + DFG\n+ types"),
     (".ql query", "security suite\n+ custom rules"),
     ("SARIF", "ruleId · location\n· data-flow path")],
    top=2.7, left=0.55, node_w=2.15, node_h=1.35, gap=0.42, accent=INDIGO,
    head_size=14, sub_size=10.5)
textbox(_s, Inches(0.7), Inches(4.9), Inches(12), Inches(1.4),
        [[("Two phases: ", 15, NAVY_TX, True, False),
          ("EXTRACTION builds the database once (the slow part — needs the "
           "build for compiled languages); then ANY NUMBER of queries run over "
           "it cheaply.", 15, INK, False, False)],
         [("Break the build → the database is incomplete → memory-safety recall "
           "collapses. (Keep builds green.)", 14, ROSE, True, False)]])
_notes(_s,
       "Walk left to right. The expensive, fragile step is the build tracer + DB "
       "extraction — it must observe compilation for C/C++/Java/Go/C#. Once the "
       "DB exists, queries are cheap and you can run many. This is why a broken "
       "build is catastrophic for CodeQL: no DB, no semantic findings. The "
       "workshop repos were chosen precisely because they build reliably.")

code_slide(
    "Anatomy of a .ql query (taint tracking)",
    ["import cpp",
     "import semmle.code.cpp.dataflow.TaintTracking",
     "",
     "// A source: command-line input (argv)",
     "// A sink:   the command string passed to system()",
     "from DataFlow::Node source, DataFlow::Node sink",
     "where",
     "    source.asExpr() instanceof ArgvAccess and          // tainted input",
     "    sink.asExpr() = aSystemCall().getArgument(0) and    // dangerous op",
     "    TaintTracking::flowPath(source, sink)               // a real path?",
     "select sink, \"argv reaches system() with no sanitizer\"  // CWE-78"],
    caption="Plain-English: “find a data-flow path from argv to system()’s "
            "argument, with no sanitizer in between.”",
    notes="You don't need to write CodeQL to use VulnHunterX, but seeing one "
          "query demystifies it. Note the three ingredients from the AST section: "
          "a SOURCE, a SINK, and flowPath() asking whether tainted data actually "
          "connects them. That's taint tracking expressed as a query. Custom "
          "rules under --profile full are exactly files like this.")

code_slide(
    "Semgrep: patterns over syntax",
    ["# A rule is YAML; the pattern looks like the target language",
     "rules:",
     "  - id: dangerous-system-call",
     "    patterns:",
     "      - pattern: system($CMD)",
     "    message: \"Command injection risk if $CMD is tainted\"",
     "    metadata:",
     "      cwe: [\"CWE-78\"]        # used later to route guided questions",
     "",
     "# No build required — reads source files directly. Fast, CI-friendly.",
     "# Registry packs: p/security-audit, p/gosec, p/owasp-top-ten",
     "# In-repo packs: config/semgrep-custom/<lang>.yaml (offline, --profile full)"],
    notes="metadata.cwe matters later — VulnHunterX routes a finding to guided "
          "questions by it. Offline caveat: p/... registry packs need "
          "semgrep.dev. For reliable offline coverage, the in-repo "
          "config/semgrep-custom/<lang>.yaml rules load under --profile full.")

table_slide(
    "How pattern matching works",
    ["Construct", "Means", "Example"],
    [["$X, $CMD", "Metavariable — matches any expression", "system($CMD)"],
     ["...", "Ellipsis — matches any sequence (args, stmts)", "foo(..., $X)"],
     ["pattern", "Must match", "exec($C)"],
     ["pattern-not", "Excludes a safe form (cuts false positives)", "exec(\"static\")"],
     ["pattern-either", "Any of several patterns (OR)", "MD5(...) | SHA1(...)"],
     ["metavariable-pattern", "Constrain what a $VAR may be", "$URL ~ user input"]],
    col_widths=[3.1, 5.4, 3.45],
    caption="Matching is on the AST, not raw text — whitespace, comments and "
            "formatting don't matter.",
    notes="The power and the limit in one slide. Metavariables + ellipsis make a "
          "pattern generalise across formatting; pattern-not is how you suppress "
          "the safe form to cut false positives. But it's still file-local and "
          "syntactic — it can't follow tainted data across functions the way "
          "CodeQL's flowPath() does. That gap is exactly why LLM verification "
          "earns its keep.")

bullets_slide(
    "SARIF: the common currency",
    [("Static Analysis Results Interchange Format — JSON.", 0, INK, True),
     ("Each finding: ruleId, file+line, severity, message, dataflow, CWE tags.", 0, INK, True),
     ("VulnHunterX discovers ALL *.sarif files and parses them to one [Finding] list.", 0, INK, True),
     ("Tool-agnostic downstream → enables cross-tool reconciliation.", 0, EMERALD, True)],
    notes="Because everything funnels through SARIF, verification doesn't care "
          "which tool produced a finding. It also lets VulnHunterX reconcile when "
          "two tools flag the same line — majority vote for sibling rules, safe "
          "FP→Needs-More-Data downgrade across different CWE classes.")

code_slide(
    "SARIF anatomy — the fields the parser reads",
    ["{ \"runs\": [ {",
     "  \"tool\": { \"driver\": { \"name\": \"CodeQL\" } },",
     "  \"results\": [ {",
     "      \"ruleId\": \"cpp/use-after-free\",          # → routes guided questions",
     "      \"level\": \"error\",                         # severity",
     "      \"message\": { \"text\": \"freed pointer reused\" },",
     "      \"locations\": [ { ... \"uri\": \"imgRead.c\",  # file",
     "                         \"startLine\": 67 } ],     # line",
     "      \"codeFlows\": [ ... ],                        # the data-flow path",
     "      \"properties\": { \"cwe\": [\"CWE-416\"] }      # → CWE fallback routing",
     "  } ] } ] }"],
    caption="SarifParser discovers every *.sarif in the output dir and flattens "
            "results into one [Finding] list.",
    notes="Don't read it line by line — point out the five fields that matter "
          "downstream: ruleId (exact question routing), level, location "
          "(file+line the LLM anchors on), codeFlows (the data-flow path CodeQL "
          "found), and properties.cwe (the fallback routing key for Semgrep). "
          "This is the literal hand-off from detection to verification.")

table_slide(
    "Strengths & weaknesses, side by side",
    ["", "CodeQL", "Semgrep"],
    [["Setup", "Needs build (compiled langs)", "None — reads source"],
     ["Depth", "Cross-function taint, CFG/DFG", "File-local pattern / AST"],
     ["Speed", "Slower (build + DB)", "Fast"],
     ["Best at", "Injection, memory-safety flow", "Config, secrets, syntactic bugs"],
     ["Offline", "Yes (queries are local)", "Custom yes; p/ packs need network"]],
    col_widths=[1.7, 5.1, 5.15],
    notes="Honest summary: use both. --tool both runs CodeQL + Semgrep; --profile full "
          "adds in-repo custom rules for both engines. The redundancy is "
          "deliberate — different engines miss different bugs, and agreement is "
          "a useful triage signal.")

# 3 — AST, control flow & data flow ----------------------------------------- #
section_slide(3, "Part 3", "AST, control flow & data flow")

bullets_slide(
    "What is an AST?",
    [("Abstract Syntax Tree — the parse tree of code as structured nodes.", 0, INK, True),
     ("x = a + b * 2  →  assignment ▸ (+) ▸ (*)  … structure, not text.", 1, GREY, False),
     ("Scanners match on the tree, not characters.", 0, INK, True),
     ("Robust to whitespace, formatting, and variable names.", 1, GREY, False),
     ("system(cmd) ≡ system( cmd ) ≡ system(/*x*/cmd)  — same tree.", 1, GREY, False)],
    notes="Why not regex the source? Because those three forms are the same "
          "program but different text. An AST normalises that away. Both Semgrep "
          "and tree-sitter work on ASTs; CodeQL builds a richer model on top.")

_s = content_slide("An AST, drawn")
ast_tree(_s, top=1.9)
_notes(_s,
       "Point at the tree. The single line x = a + b * 2 becomes a tree whose "
       "ROOT is the assignment. Operator precedence isn't a parsing rule the tool "
       "re-derives every time — it's baked into the shape: * sits below +. A "
       "scanner pattern like '$X = $A + $B' matches the top two levels regardless "
       "of spacing or names. This is the structure everything else builds on.")

bullets_slide(
    "From AST to control flow & data flow",
    [("Control-Flow Graph (CFG): the order statements CAN execute.", 0, INK, True),
     ("Data-Flow Graph / taint: how values move between variables.", 0, INK, True),
     ("Security questions are flow questions:", 0, NAVY_TX, True),
     ("SOURCE (untrusted input)  →  SANITIZER  →  SINK (dangerous op)", 1, INDIGO, True),
     ("A vuln = source → sink flow with NO adequate sanitizer.", 0, ROSE, True)],
    notes="This source → sanitizer → sink model is the single most useful idea "
          "in the talk. Taint tracking = 'is there a path on the data-flow graph "
          "from source to sink?' CodeQL answers semantically across functions. "
          "The guided questions push the LLM to anchor on exactly this.")

# Control-flow graph mini-diagram
_s = content_slide("Control flow: which paths can execute")
flow_node(_s, Inches(5.3), Inches(1.6), Inches(2.9), Inches(0.85),
          "read len from input", "", SLATE, head_size=14)
flow_node(_s, Inches(5.3), Inches(2.85), Inches(2.9), Inches(0.95),
          "if (len > MAX) ?", "branch", INDIGO, head_size=14)
flow_node(_s, Inches(1.9), Inches(4.4), Inches(3.1), Inches(0.95),
          "return error", "safe path", EMERALD, head_size=14)
flow_node(_s, Inches(8.4), Inches(4.4), Inches(3.1), Inches(0.95),
          "memcpy(dst, src, len)", "sink — overflow if len > MAX", ROSE, head_size=14)
connect(_s, Inches(6.75), Inches(2.45), Inches(6.75), Inches(2.85), SLATE, width=2.5)
connect(_s, Inches(5.6), Inches(3.55), Inches(3.45), Inches(4.4), EMERALD, width=2.5)
connect(_s, Inches(7.9), Inches(3.55), Inches(9.95), Inches(4.4), ROSE, width=2.5)
textbox(_s, Inches(3.3), Inches(3.75), Inches(2.2), Inches(0.4),
        [[("len ≤ MAX", 12, EMERALD, True, True)]])
textbox(_s, Inches(8.1), Inches(3.75), Inches(2.4), Inches(0.4),
        [[("len > MAX", 12, ROSE, True, True)]])
textbox(_s, Inches(0.7), Inches(5.9), Inches(12), Inches(1.0),
        [[("The CFG is the set of paths control CAN take. The bug only exists on "
           "ONE branch — and only if that branch is reachable with a bad len. "
           "Proving reachability is what SAST struggles with (→ false positives) "
           "and what the LLM is asked to reason about.", 15, INK, False, False)]])
_notes(_s,
       "Control flow = the branches and loops. Notice the vulnerability lives on "
       "exactly one edge: len > MAX into the memcpy. A pattern tool sees the "
       "memcpy and flags it; it can't easily tell whether the guard above makes "
       "the bad branch unreachable. That reachability question is the heart of "
       "false positives — and the heart of what verification reasons about.")

# Taint / data-flow diagram
_s = content_slide("Data flow: does tainted input reach a sink?")
taint_diagram(_s, top=2.5)
textbox(_s, Inches(0.7), Inches(5.65), Inches(12), Inches(1.2),
        [[("A vulnerability = a data-flow path from a SOURCE to a SINK with no "
           "adequate SANITIZER. ", 15, INK, True, False),
          ("Taint tracking asks exactly that question across the whole program. "
           "Find the source, find the sink, check what's in between.", 15, INK,
           False, False)]])
_notes(_s,
       "This is the one diagram to remember. Source = untrusted input. Sink = "
       "dangerous operation. Sanitizer = the check that makes it safe. The red "
       "bypass is the bug: input reaches the sink unchecked. Every guided "
       "question set is really pushing the LLM to fill in these three boxes for "
       "the specific finding.")

table_slide(
    "CodeQL vs tree-sitter for building the AST",
    ["", "CodeQL", "tree-sitter"],
    [["Model", "Semantic: AST+CFG+DFG+types", "Syntactic AST only"],
     ["Needs build?", "Yes (compiled langs)", "No — parses source"],
     ["Speed", "Slower", "Very fast"],
     ["Depth", "Cross-function, type-aware", "File-local, structural"],
     ["In VulnHunterX", "Default context backend", "Fast fallback; all 8 languages"]],
    col_widths=[2.0, 4.95, 5.0],
    notes="Both appear as context-extraction backends "
          "(--backend {auto,codeql,treesitter}). CodeQL = semantic context but "
          "needs a build; tree-sitter = fast, build-free, any language. 'auto' "
          "picks CodeQL when a DB exists, else tree-sitter.")

card_grid_slide(
    "tree-sitter — the fast, build-free parser",
    [{"head": "Incremental parser", "accent": CYAN,
      "body": "Re-parses only edited regions — built for editor speed."},
     {"head": "Per-language grammars", "accent": INDIGO,
      "body": "One grammar per language; VulnHunterX covers all 8."},
     {"head": "No build needed", "accent": EMERALD,
      "body": "Reads source directly — works even when the build is broken."},
     {"head": "Error-tolerant", "accent": AMBER,
      "body": "Recovers from syntax errors; parses partial / WIP code."},
     {"head": "Syntactic only", "accent": ROSE,
      "body": "AST + structure, but no types, no cross-function data-flow."},
     {"head": "Role in VulnHunterX", "accent": ORANGE,
      "body": "The 'treesitter' context backend — fallback when no CodeQL DB."}],
    cols=3, top=1.9, card_h=1.95,
    notes="tree-sitter is the parser behind syntax highlighting in many editors. "
          "Its superpowers for us: no build, all languages, error-tolerant. Its "
          "limit: purely syntactic — it can give you the enclosing function and "
          "structure, but not types or taint. That's why it's the FALLBACK "
          "backend: great when CodeQL can't build, but CodeQL is richer when it can.")

table_slide(
    "Doesn't Semgrep already build an AST?",
    ["", "Semgrep's AST", "tree-sitter (in VulnHunterX)"],
    [["Whose engine", "Internal to Semgrep", "VulnHunterX's ContextExtractor"],
     ["Purpose", "DETECT findings (pattern match)", "EXTRACT context for the LLM"],
     ["Reusable?", "No — private to the matcher", "Yes — query for functions, callers, structs"],
     ["Scope", "The rule's match site", "Whole codebase → context CSVs, all 8 langs"],
     ["When used", "Only if Semgrep runs & fires", "Any finding (incl. CodeQL-only), build-free"]],
    col_widths=[1.9, 4.6, 5.45],
    caption="Different ASTs, different jobs — detection vs context.",
    notes="Direct answer to the obvious question. Yes, Semgrep parses to an AST — "
          "but that tree is PRIVATE to its pattern matcher and exists only to "
          "DETECT. VulnHunterX needs its OWN queryable AST to EXTRACT context — "
          "the enclosing function, callers, struct fields — that grounds the LLM. "
          "Crucially, findings often come from CodeQL alone (Semgrep may not run "
          "or may not fire), so the tool can't depend on Semgrep's internal tree "
          "at all. CodeQL is the richer semantic backend; tree-sitter is the fast, "
          "build-free fallback covering all 8 languages "
          "(src/vuln_hunter_x/context/treesitter_extractor.py).")

# AST context — flowchart
_s = content_slide("AST context: from repo to LLM prompt")
flow_chain(
    _s,
    [("Repo", "source tree"),
     ("Extraction", "CodeQL / tree-sitter"),
     ("Context CSVs", "functions, callers,\nstructs, free_sites…"),
     ("ContextProvider", "look-up + cache"),
     ("Code slice", "enclosing fn +\nrequested context"),
     ("LLM prompt", "grounded, not\nthe whole repo")],
    top=3.0, left=0.5, node_w=1.95, node_h=1.4, gap=0.18, accent=INDIGO,
    head_size=13, sub_size=9.5)
textbox(_s, Inches(0.7), Inches(5.2), Inches(12), Inches(1.3),
        [[("The model never sees the 100k-line repo. ", 15, NAVY_TX, True, False),
          ("It sees a tight, AST-derived slice that's actually relevant to the "
           "finding — the single most important defence against hallucination.",
           15, INK, False, False)]])
_notes(_s,
       "This flowchart is the bridge from Part 3 into the VulnHunterX sections. "
       "Stage 1 turns the repo into CSVs; at verify time the ContextProvider "
       "assembles a small, relevant slice and only that goes into the prompt. "
       "Grounding the model in real extracted code — instead of its training "
       "memory — is what keeps verdicts honest. Full diagram in "
       "docs/context-extraction-flow.md.")

# AST context — the 11 context types
table_slide(
    "The context types the LLM can pull",
    ["Context type", "Answers", "When"],
    [["function / method", "What does this code do?", "pre-fetched"],
     ["caller / all_callers", "Who calls it? Where does input come from?", "pre-fetched"],
     ["callees", "What sinks does it call?", "pre-fetched"],
     ["struct / class", "What are the field types & layout?", "reactive"],
     ["global / macro", "What's this constant / expansion?", "reactive"],
     ["typedef / enum", "What's the underlying type / values?", "reactive"],
     ["free_sites  (C/C++)", "Where is this pointer freed?", "reactive"],
     ["destructor  (C/C++)", "What cleanup runs on this type?", "reactive"],
     ["field_writes  (C/C++)", "Where is this field written?", "reactive"]],
    col_widths=[3.3, 5.65, 3.0],
    caption="Two phases: cheap function-keyed context is pre-fetched; the rest is "
            "fetched only when the LLM asks (context_needed).",
    notes="11 context types, grouped by phase. Pre-fetched = cheap, keyed on the "
          "function name, sent up front. Reactive = fetched only when the model "
          "requests it in a Needs-More-Data turn. The C/C++ memory-safety trio "
          "(free_sites, destructor, field_writes) is what powers use-after-free "
          "and double-free verification specifically.")

# 4 — LLM verification ------------------------------------------------------- #
section_slide(4, "Part 4", "LLM verification — pros & cons")

bullets_slide(
    "Why add an LLM after the scanner?",
    [("Scanner: “line 58 might be a heap overflow.”", 0, GREY, True),
     ("It can't add: “…and it's reachable with attacker-controlled length → real.”", 1, GREY, False),
     ("An LLM, given code + context, can reason about reachability & sanitizers.", 0, INK, True),
     ("Goal: noisy candidate list → ranked, explained, triaged report.", 0, EMERALD, True),
     ("It does NOT detect new bugs — it judges the ones SAST produced.", 0, ROSE, True)],
    notes="The LLM does first-pass human-style triage at machine speed and scale. "
          "It is explicitly not the detector — it never invents findings.")

card_grid_slide(
    "What we actually ask the LLM",
    [{"head": "Verdict", "accent": INDIGO,
      "body": "TRUE_POSITIVE · FALSE_POSITIVE · NEEDS_MORE_DATA"},
     {"head": "Confidence", "accent": AMBER,
      "body": "high / medium / low  →  0.85 / 0.6 / 0.3"},
     {"head": "Reasoning", "accent": EMERALD,
      "body": "Anchored to specific lines (no hand-waving)"},
     {"head": "Answers", "accent": CYAN,
      "body": "Responses to the rule's guided questions"},
     {"head": "context_needed", "accent": ORANGE,
      "body": "Extra code it wants to see next → drives the multi-turn loop"}],
    cols=3, top=1.9, card_h=1.7,
    notes="NEEDS_MORE_DATA is the interesting one — instead of guessing, the "
          "model asks for the caller or struct, the engine fetches it and "
          "re-asks. That loop separates this from a one-shot 'is this a bug?'.")

bullets_slide(
    "Pros of LLM verification",
    [("Cuts false-positive noise — the headline benefit.", 0, EMERALD, True),
     ("Explains each verdict in plain language → faster triage & learning.", 0, INK, True),
     ("Fills context gaps SAST leaves (intent, framework conventions).", 0, INK, True),
     ("Language-flexible — one reasoning engine across C, Python, Java, JS, Go…", 0, INK, True)],
    notes="The explanation is underrated. 'CWE-416 at line 67' teaches nothing; "
          "'the pointer freed at 59 is dereferenced at 67 when size1 % 3 == 0' "
          "teaches a junior dev to spot use-after-free. The tool teaches while "
          "it triages.")

bullets_slide(
    "Cons & risks of LLM verification",
    [("Hallucination — can assert a sanitizer or flow that isn't there.", 0, ROSE, True),
     ("Non-determinism — same finding, different runs, maybe different verdict.", 0, ROSE, True),
     ("Cost & latency — tokens and time per finding, especially multi-turn.", 0, ROSE, True),
     ("Prompt injection — the code under review is attacker-influenced text.", 0, ROSE, True),
     ("Overconfidence — fluent prose ≠ correct reasoning.", 0, ROSE, True)],
    notes="Be honest — engineers respect it. None of these are solved perfectly; "
          "they're mitigated (next slide). The tool is an assistant, not an "
          "oracle — the closing line of the whole lesson.")

table_slide(
    "How VulnHunterX mitigates the cons",
    ["Risk", "Mitigation"],
    [["Hallucination", "Grounding — CSV/AST context; questions demand line citations"],
     ["Weak / uncited claims", "Confidence downgrades — vague reasoning auto-demoted to low"],
     ["Premature verdicts", "min_iterations on high-stakes CWEs (e.g. CWE-416, CWE-200)"],
     ["Non-determinism", "Temperature 0.2, structured JSON, tool reconciliation"],
     ["Single-tool error", "Cross-tool reconciliation (CodeQL vs Semgrep)"]],
    col_widths=[3.1, 8.85],
    notes="Recurring theme: grounding and structure beat free-form prompting. "
          "Guided questions force evidence; the context provider supplies real "
          "code; calibrators downgrade hand-waving; high-stakes findings can't be "
          "dismissed on turn one. That's the engineering behind a reliable pipeline.")

# 5 — VulnHunterX architecture ----------------------------------------------- #
section_slide(5, "Part 5", "VulnHunterX (no fuzzing): architecture & stages")

code_slide(
    "The big picture",
    ["SARIF files",
     "  └─► SarifParser → [Finding]",
     "        └─► VerificationEngine",
     "              ├── QuestionsLoader   (rule-specific guided questions)",
     "              ├── ContextProvider   (CSV / AST look-ups)",
     "              └── LLMClient (LiteLLM → OpenAI · Anthropic · Ollama)",
     "                    └─► Verdict (TP | FP | NEEDS_MORE_DATA) + confidence"],
    notes="Everything assembles. SAST tools produce SARIF; the parser normalises; "
          "the engine orchestrates questions + context + LLM; output is a verdict "
          "per finding. LiteLLM is the provider abstraction — same path for "
          "OpenAI, Anthropic, or Ollama (the workshop uses Ollama Cloud).")

table_slide(
    "The pipeline: Stages 1–4 (static + LLM)",
    ["Stage", "Command", "In → Out"],
    [["1", "prepare", "repo URL/path → source + CodeQL DB + context CSVs"],
     ["2", "analyze", "DB/source → SARIF (CodeQL / Semgrep)"],
     ["3", "verify", "SARIF + context → verdicts (TP/FP/NMD + confidence)"],
     ["4", "report", "verdicts → Markdown report (EN / VI)"]],
    col_widths=[1.1, 2.0, 8.85],
    caption="Stages 5–8 add ASan build + fuzzing for C/C++ — out of scope today.",
    notes="Four stages to a report. 'scan' runs all four; the individual commands "
          "let you re-run one stage (e.g. re-verify with a different model without "
          "rebuilding the DB). Stage 1 creates the AST context CSVs from Part 3.")

_s = content_slide("The pipeline, as a flow")
flow_chain(
    _s,
    [("1 · prepare", "clone + CodeQL DB\n+ context CSVs"),
     ("2 · analyze", "CodeQL / Semgrep\n→ SARIF"),
     ("3 · verify", "LLM triage →\nverdicts"),
     ("4 · report", "Markdown report\nEN / VI")],
    top=2.45, left=0.7, node_w=2.7, node_h=1.5, gap=0.62, accent=ORANGE,
    head_size=16, sub_size=11)
textbox(_s, Inches(0.7), Inches(4.5), Inches(12), Inches(0.5),
        [[("artifacts under  output/<lang>/<repo>/ :", 13.5, MUTE, True, True)]])
flow_chain(
    _s,
    [("database/  +  context/", ""),
     ("*.sarif", ""),
     ("verification_results/", ""),
     ("report.md", "")],
    top=5.05, left=0.7, node_w=2.7, node_h=0.7, gap=0.62, accent=SLATE,
    head_size=12.5)
_notes(_s,
       "Same four stages, now as data flow. Call out that each arrow is a "
       "persisted artifact on disk — so you can stop and resume, or re-run just "
       "verify against the same SARIF with a different model. 'scan' is just "
       "this whole chain in one command.")

table_slide(
    "Rule profiles: dialing coverage up",
    ["Profile", "CodeQL suite", "Semgrep packs", "Custom"],
    [["standard*", "security-extended", "auto + per-lang pack", "—"],
     ["extended", "security-extended", "+ security-audit, secrets", "—"],
     ["maximum", "security-and-quality", "+ owasp-top-ten", "—"],
     ["extended-registry", "security-and-quality", "8 universal + per-lang", "—"],
     ["full", "security-and-quality", "extended-registry packs", "+ in-repo custom"]],
    col_widths=[2.7, 3.4, 4.05, 1.8],
    caption="* standard is the default.",
    notes="--profile full is the one to remember for offline/self-hosted use — it "
          "loads the in-repo custom CodeQL queries and Semgrep rules, which work "
          "without semgrep.dev. The workshop uses full for exactly this reason.")

table_slide(
    "What the suites & packs actually contain",
    ["Suite / pack", "What it is"],
    [["CodeQL · security-extended", "High-confidence security queries — fewer FPs, the safe default"],
     ["CodeQL · security-and-quality", "security-extended + maintainability/quality queries — broader, noisier"],
     ["Semgrep · auto", "Registry auto-detects languages and applies matching rulesets"],
     ["Semgrep · p/security-audit", "Broad cross-language security audit ruleset"],
     ["Semgrep · p/secrets", "Hardcoded credentials, API keys, tokens"],
     ["Semgrep · p/owasp-top-ten", "Rules grouped by the OWASP Top-10 categories"],
     ["Semgrep · per-language", "e.g. p/gosec (Go), p/django/p/flask (Py), p/nodejs (JS)"]],
    col_widths=[4.4, 7.55],
    caption="A profile is just a preset bundle of one CodeQL suite + a set of Semgrep packs.",
    notes="This unpacks the previous table. Two CodeQL suites: security-extended "
          "(curated, high-precision security queries) and security-and-quality "
          "(adds maintainability/quality lint — more coverage but more noise, which "
          "is why higher profiles pair it with the LLM triage). Semgrep packs: "
          "'auto' lets the registry pick rulesets per language; p/security-audit is "
          "the broad sweep; p/secrets catches hardcoded keys; p/owasp-top-ten maps "
          "to the OWASP categories; and per-language packs add framework-aware "
          "rules. 'full' layers our in-repo custom rules on top of all of that.")

bullets_slide(
    "Custom rules (the full profile)",
    [(f"Custom CodeQL — config/codeql-custom/<lang>/src/*.ql  ({CODEQL_TOTAL} total).", 0, INK, True),
     ("e.g. cpp/use-after-free, cpp/double-free; @id matches a question key exactly.", 1, GREY, False),
     (f"Custom Semgrep — config/semgrep-custom/<lang>.yaml  ({SEMGREP_TOTAL} total).", 0, INK, True),
     ("e.g. vulnhunterx.cpp.weak-hash; each sets metadata.cwe for routing.", 1, GREY, False),
     ("Counts are computed from config/ at build time — never drift.", 0, EMERALD, True)],
    notes="Custom rules close gaps the built-in suites miss, especially "
          "framework-aware and C/C++ memory-safety classes. They're wired to the "
          "guided-question system. Run scripts/audit_rule_coverage.py to see "
          "coverage by rule × CWE × tool.")

bullets_slide(
    "Guided questions",
    [(f"Per-language YAML banks: config/prompts/<lang>_questions.yaml ({QUESTIONS_TOTAL} entries).", 0, INK, True),
     ("Each rule → evidence-anchored questions the LLM must answer first.", 0, INK, True),
     ("UAF e.g.: “Where is the pointer allocated? Where freed? Quote the flagged line.”", 1, GREY, False),
     ("Also sets additional_context to prefetch, min_iterations, snippet window.", 0, INK, True)],
    notes="The secret sauce against hallucination. Instead of 'is this a bug?', "
          "the model walks the evidence: anchor on the flagged line, trace "
          "allocation → free → use. A model that must cite lines can't hand-wave.")

table_slide(
    "Anatomy of a guided-question entry",
    ["Key", "Purpose"],
    [["short_description", "One-line summary of the vulnerability class"],
     ["questions", "Ordered, evidence-anchored questions the LLM must answer first"],
     ["context_hint", "Human-readable trace plan (allocation → free → use)"],
     ["additional_context", "Context types to prefetch / allow requesting"],
     ["min_iterations", "Floor on rounds before TP/FP (high-stakes CWEs)"],
     ["snippet_window_lines", "How much code around the flagged line to show"]],
    col_widths=[3.3, 8.65],
    caption="config/prompts/<lang>_questions.yaml — one entry per rule.",
    notes="Each rule maps to a YAML entry shaped like this. The questions are the "
          "core, but note min_iterations (can't snap-judge a use-after-free) and "
          "additional_context (which of the 11 context types this class needs). "
          "These keys are read by questions/loader.py.")

bullets_slide(
    "Real example — py/sql-injection (verbatim)",
    [("Quote the EXACT sink (cursor.execute / session.execute, line N) and name "
      "the variable passed to it. Do not paraphrase.", 0, INK, False),
     ("List EVERY assignment to that variable on each path reaching the sink, with "
      "line numbers; identify the LAST one that executes.", 0, INK, False),
     ("For each, does the value derive from user input "
      "(request.args/form/json/cookies) or a constant? Cite the data-flow chain.", 0, INK, False),
     ("What specific defence sanitises it? Name it — parameterised "
      "execute(sql, params), Django ORM .filter(**kwargs)… vague 'sanitisation' "
      "is not acceptable.", 0, INK, False),
     ("If you can't show the tainted value reaching the sink AND the absence of "
      "every defence, verdict FALSE POSITIVE.", 0, ROSE, True)],
    notes="The actual questions from python_questions.yaml (py/sql-injection), "
          "lightly trimmed — chosen to show the same evidence-first philosophy on a "
          "different language and class (taint, not memory safety). It forces a "
          "concrete sink quote, a per-path data-flow trace, a NAMED sanitiser (not "
          "'it's escaped'), and a default-to-FP rule when evidence is missing. "
          "min_iterations=2. This is the CWE-89 class that scored F1=100% in the "
          "OWASP-Python benchmark later.")

bullets_slide(
    "Real example — cpp/overflow-buffer (verbatim)",
    [("Where is the DESTINATION buffer declared and what is its exact size?", 0, INK, False),
     ("Where is the SOURCE buffer/data declared and what is its exact size?", 0, INK, False),
     ("Do these sizes ever CHANGE via assignment, reallocation, or parameters?", 0, INK, False),
     ("Are there any CHECKS on source size, destination size, or copy length "
      "BEFORE the copy?", 0, INK, False),
     ("→ The pattern: pin down sizes, track changes, look for the missing bound.", 0, EMERALD, True)],
    notes="A second class to show the pattern generalises: for overflows the "
          "questions force the model to pin the destination size, the source "
          "size, whether either changes, and whether a bound exists before the "
          "copy. Same philosophy — make the model produce evidence, not vibes.")

bullets_slide(
    "How a finding finds its questions (3-tier routing)",
    [("1. Exact match on SARIF ruleId → <lang>/<name>  (custom CodeQL hits here).", 0, INK, True),
     ("2. Normalized / prefix / language-prefix fallbacks.", 0, INK, True),
     ("3. CWE-based fallback via cwe_question_map  (custom Semgrep routes here).", 0, INK, True),
     ("else → generic question template.", 1, GREY, False)],
    notes="CodeQL IDs look like cpp/use-after-free (slash) → match a key directly. "
          "Semgrep IDs look like vulnhunterx.cpp.weak-hash (dots) → never match by "
          "name, so route by CWE tag. Every finding gets some interrogation.")

_s = content_slide("Routing, traced: two findings, two paths")
textbox(_s, Inches(0.7), Inches(1.45), Inches(12), Inches(0.4),
        [[("CodeQL finding — exact ruleId match (Tier 1):", 14, INDIGO, True, False)]])
flow_chain(
    _s,
    [("ruleId:\ncpp/use-after-free", ""),
     ("exact key\nmatch", ""),
     ("cpp/use-after-free\nquestions", "")],
    top=1.95, left=0.7, node_w=3.3, node_h=1.05, gap=0.7, accent=INDIGO,
    head_size=13, sub_size=10)
textbox(_s, Inches(0.7), Inches(3.45), Inches(12), Inches(0.4),
        [[("Semgrep finding — dotted id, no key → CWE fallback (Tier 3):", 14, CYAN, True, False)]])
flow_chain(
    _s,
    [("id:\nvulnhunterx.cpp.weak-hash", ""),
     ("no key\n(dotted id)", ""),
     ("metadata.cwe:\nCWE-327", ""),
     ("cwe_question_map\n→ weak-hash questions", "")],
    top=3.95, left=0.7, node_w=2.75, node_h=1.05, gap=0.32, accent=CYAN,
    head_size=12, sub_size=9.5)
textbox(_s, Inches(0.7), Inches(5.5), Inches(12), Inches(1.0),
        [[("Both findings end up with the right evidence-anchored questions — one "
           "by name, one by CWE. No finding goes uninterrogated.", 15, INK, False, False)]])
_notes(_s,
       "Trace both paths concretely. The CodeQL rule's @id IS the question key, so "
       "Tier 1 hits immediately. The Semgrep rule's dotted id can never match a "
       "slash-style key, so it falls through to Tier 3 and routes by its "
       "metadata.cwe (CWE-327 for weak-hash) via cwe_question_map. Same "
       "destination, different door.")

bullets_slide(
    "Context extraction & the ContextProvider",
    [("Pre-fetched upfront (keyed on function name):", 0, NAVY_TX, True),
     ("enclosing function, caller, callees, all_callers", 1, GREY, False),
     ("Reactive — only when context_needed asks:", 0, NAVY_TX, True),
     ("struct, global, macro, typedef, enum + (C/C++) free_sites, destructor, field_writes", 1, GREY, False),
     ("All reads are cached and path-traversal-safe; CSVs come from Stage 1.", 0, EMERALD, True)],
    notes="Two-phase model: cheap context goes in immediately; expensive or "
          "finding-specific context is fetched lazily when the model asks. The 11 "
          "context types map to the evidence the guided questions demand.")

_s = content_slide("Multi-turn verification loop")
# nodes
flow_node(_s, Inches(0.85), Inches(2.5), Inches(3.0), Inches(1.1),
          "LLM analyzes", "finding + context", INDIGO, head_size=15, sub_size=11)
# decision diamond
dia = shape(_s, MSO_SHAPE.DIAMOND, Inches(4.55), Inches(2.25), Inches(2.7),
            Inches(1.6), fill=CARD, line=AMBER, line_w=1.6)
text_in(dia, [[("verdict?", 14, NAVY_TX, True, False)]])
flow_node(_s, Inches(8.5), Inches(2.5), Inches(3.4), Inches(1.1),
          "Return Verdict", "TP / FP + confidence", EMERALD, head_size=15, sub_size=11)
flow_node(_s, Inches(4.5), Inches(4.85), Inches(2.8), Inches(1.15),
          "Fetch context", "free_sites · caller ·\nstruct …", AMBER,
          head_size=14, sub_size=10.5)
# arrows
connect(_s, Inches(3.85), Inches(3.05), Inches(4.55), Inches(3.05), SLATE, width=2.5)
connect(_s, Inches(7.25), Inches(3.05), Inches(8.5), Inches(3.05), EMERALD, width=2.8)
textbox(_s, Inches(7.2), Inches(2.55), Inches(1.4), Inches(0.5),
        [[("TP / FP", 11.5, EMERALD, True, True)],
         [("(min_iters met)", 9.5, MUTE, False, True)]], align=CENTER)
connect(_s, Inches(5.9), Inches(3.85), Inches(5.9), Inches(4.85), AMBER, width=2.5)
textbox(_s, Inches(6.0), Inches(4.0), Inches(2.6), Inches(0.4),
        [[("NEEDS_MORE_DATA", 11, AMBER, True, True)]])
# loop back: fetch -> LLM (elbow)
connect(_s, Inches(4.5), Inches(5.4), Inches(2.35), Inches(5.4), SLATE, width=2.5, elbow=False)
connect(_s, Inches(2.35), Inches(5.4), Inches(2.35), Inches(3.6), SLATE, width=2.5)
textbox(_s, Inches(2.5), Inches(4.9), Inches(2.0), Inches(0.4),
        [[("re-ask (dedup)", 11, GREY, True, True)]])
textbox(_s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.7),
        [[("Guards: ", 14, NAVY_TX, True, False),
          ("dedup against already-fetched context · capped at max_iterations (10) · "
           "min_iterations floor on high-stakes CWEs · force-decision if it stalls.",
           14, INK, False, False)]])
_notes(_s,
       "The heart of the engine, as a loop. The model emits a verdict + a "
       "context_needed list. TP/FP (once min_iterations is met) returns. "
       "NEEDS_MORE_DATA routes down: the provider fetches the requested context "
       "(free_sites, caller, struct…), appends it, and the LLM re-asks. Requests "
       "are deduped so it can't spin forever; max_iterations caps it; a "
       "force-decision turn makes it commit if it stalls.")

bullets_slide(
    "Output: the report",
    [("verification_results/ — one verdict_*.json per finding + summary.json.", 0, INK, True),
     ("report.md (English) and report_vi.md (Vietnamese).", 0, INK, True),
     ("Executive summary + per-finding verdict, confidence, reasoning, location.", 1, GREY, False),
     ("A Coverage-Limitations caveat documents what SAST can't reach.", 0, INK, True)],
    notes="The report is what a developer reads. It leads with confident True "
          "Positives, explains each, and is honest about classes the pipeline "
          "can't cover (e.g. missing-authorization, no syntactic signature). "
          "We'll read one in the workshop.")

# 6 — How to use it ---------------------------------------------------------- #
section_slide(6, "Part 6", "How to use it: CLI + a worked example")

table_slide(
    "The command map",
    ["Command", "What it does"],
    [["check-env", "Verify CodeQL / Semgrep / LLM provider"],
     ["prepare", "Clone + CodeQL DB + context CSVs (Stage 1)"],
     ["analyze", "Run CodeQL / Semgrep → SARIF (Stage 2)"],
     ["verify", "LLM triage of SARIF → verdicts (Stage 3)"],
     ["report", "(Re)generate the Markdown report (Stage 4)"],
     ["scan", "All of the above in one shot"],
     ["interactive", "Guided wizard — friendliest entry point"]],
    col_widths=[2.4, 9.55],
    notes="'scan' for convenience, individual stages for control. 'interactive' "
          "walks newcomers through every option and live-tests the LLM before "
          "running — a great first experience.")

code_slide(
    "Key flags",
    ["--tool      {codeql, semgrep, both}              # which scanners",
     "--profile   {standard, extended, maximum, extended-registry, full}",
     "--provider  {openai, anthropic, ollama, deepseek}      # LLM backend",
     "--model     <name>                  # e.g. gpt-4o, ollama/...:cloud",
     "--limit     N           # cap findings to verify (cost control)",
     "--max-iterations N      # multi-turn rounds (default 10)",
     "-j / --jobs N           # parallel findings (default 4)",
     "--lang      <lang>      # required with --url / --local-path",
     "--dry-run               # preview, no LLM calls"],
    notes="--limit is your friend for cost control — verify 5, not 500, while "
          "learning. --provider/--model swap LLMs without touching code (LiteLLM). "
          "--dry-run previews with zero token spend.")

code_slide(
    "Worked example: one repo, end to end",
    ["# One command — prepare, analyze (CodeQL + Semgrep), verify, report:",
     "vuln-hunter-x scan \\",
     "  --url https://github.com/patricia-gallardo/insecure-coding-examples.git \\",
     "  --lang cpp --tool both --profile full --limit 8",
     "",
     "# Read the triaged report:",
     "cat output/cpp/insecure-coding-examples/verification_results/report.md"],
    notes="This is the workshop in one slide. scan clones, builds a CodeQL DB, "
          "runs CodeQL + Semgrep with the full profile, sends the first 8 "
          "findings through LLM triage, writes a Markdown report. In the workshop "
          "we do it stage-by-stage and use Ollama Cloud — no per-token bill.")

# 7 — Results & limitations -------------------------------------------------- #
section_slide(7, "Part 7", "Results & limitations")

card_grid_slide(
    "Why VulnHunterX — and why it keeps improving",
    [{"head": "Actively maintained", "accent": EMERALD,
      "body": "Rules + guided questions are added continuously — coverage grows release over release."},
     {"head": f"{CODEQL_TOTAL + SEMGREP_TOTAL} custom rules", "accent": INDIGO,
      "body": f"{CODEQL_TOTAL} CodeQL queries + {SEMGREP_TOTAL} Semgrep rules, plus {QUESTIONS_TOTAL} guided-question entries (live count)."},
     {"head": "Explainable & calibrated", "accent": CYAN,
      "body": "Every verdict cites lines + a confidence that tracks accuracy (see benchmarks)."},
     {"head": "8 languages", "accent": AMBER,
      "body": "C, C++, Python, JavaScript, Java, PHP, Go, C#."},
     {"head": "Offline-capable", "accent": EMERALD,
      "body": "--profile full ships in-repo rules; runs with a local/cloud Ollama model."},
     {"head": "Provider-flexible", "accent": INDIGO,
      "body": "OpenAI · Anthropic · DeepSeek · Ollama — swap with one flag."}],
    cols=3, top=1.85, card_h=2.0,
    notes="Open with the strengths before the honest limitations. The headline the "
          "user cares about: the tool is actively maintained — custom rules and "
          "rule-specific guided questions are added continuously, so coverage and "
          "triage quality improve over time (the counts here are read live from "
          "config/). Plus: multi-language, offline-capable via --profile full + "
          "Ollama, provider-flexible, and every verdict is explainable and "
          "confidence-calibrated.")

card_grid_slide(
    "What it does well — and what it misses",
    [{"head": "✓ Syntactic classes", "accent": EMERALD,
      "body": "Unsafe functions, format string, double-free, use-after-free."},
     {"head": "✓ Data-flow injection", "accent": EMERALD,
      "body": "Cross-function taint via CodeQL."},
     {"head": "✓ Less review load", "accent": EMERALD,
      "body": "LLM triage + cross-tool reconciliation cut false positives."},
     {"head": "✗ Detection misses", "accent": ROSE,
      "body": "Integer-overflow chains, divide-by-zero, leaks (no rule fires)."},
     {"head": "✗ Hard tier", "accent": ROSE,
      "body": "C++ lifetime/UB: dangling string_view, removed memset, RAII."},
     {"head": "✗ Out of SAST reach", "accent": ROSE,
      "body": "Missing authz, CSRF-disabled (absence-of-control). Failed build → CodeQL recall collapses."}],
    cols=3, top=1.85, card_h=2.0,
    notes="Numbers come from docs/benchmarks/ground-truth-baselines.md. Two "
          "failure modes: detection misses (no rule fired) vs triage misses "
          "(fired but verified FP/NMD). Cost: LLM verification spends tokens and "
          "time — use --limit and a cheap model while iterating.")

bullets_slide(
    "How we benchmark",
    [("A standalone harness scores VulnHunterX vs raw SAST on ground-truth datasets.", 0, INK, True),
     ("Datasets: OWASP Benchmark (Java/Python), Juliet C/C++ (64K cases), "
      "SecLLMHolmes, RealVuln, OpenVuln/ZeroFalse.", 0, INK, False),
     ("Approaches compared: raw-sast (no LLM) vs vulnhunterx (guided questions + "
      "multi-turn) + ablations.", 0, INK, False),
     ("Metrics: precision · recall · F1 · FP-reduction · NMD rate · tokens/finding "
      "· p95 latency · confidence calibration.", 0, INK, False),
     ("Reproducible: benchmarks/scripts/run_benchmark.py & run_model_matrix.py.", 0, GREY, False)],
    notes="Set up the next two slides. The harness in benchmarks/ runs the same "
          "datasets through raw-sast (treat every finding as TP — the baseline) "
          "and through the full VulnHunterX pipeline, then scores both against "
          "ground truth. It also runs a model matrix to compare LLMs. Everything "
          "is reproducible from the repo. Numbers on the next slides are from "
          "committed runs under benchmarks/results/.")

table_slide(
    "Benchmark: false-positive reduction (OWASP-Python)",
    ["Approach", "Precision", "Recall", "F1"],
    [["raw-sast (baseline, no LLM)", "38.6%", "100.0%", "55.7%"],
     ["vulnhunterx (guided + multi-turn)", "65.9%", "98.0%", "78.8%"]],
    col_widths=[5.95, 2.0, 2.0, 2.0],
    caption="Model: qwen3-coder:480b (Ollama Cloud) · 500 Python findings · "
            "benchmarks/results/.../REPORT.md",
    notes="The core result. Raw SAST flags everything, so recall is 100% but "
          "precision is only 39% — three in five alerts are false alarms. "
          "VulnHunterX lifts precision to 66% (F1 56%→79%) while keeping recall at "
          "98% — i.e. it removes most false positives and keeps almost all real "
          "bugs. Per-CWE, SQL-injection (CWE-89) and command-injection (CWE-78) "
          "both hit F1=100% on this run. Confidence is well-calibrated: High-"
          "confidence verdicts were ~81% accurate vs ~68% for Low.")

table_slide(
    "Benchmark: model matrix (Juliet C/C++)",
    ["Model", "Precision", "Recall", "F1", "FP-reduc.", "Cost"],
    [["raw-sast (baseline)", "50.0%", "100.0%", "66.7%", "0%", "$0"],
     ["DeepSeek", "83.8%", "93.8%", "88.5%", "82.2%", "$0.43"],
     ["GPT-4.1-mini", "56.8%", "100.0%", "72.4%", "23.9%", "$1.02"]],
    col_widths=[3.45, 1.85, 1.65, 1.4, 1.8, 1.8],
    caption="vulnhunterx approach on Juliet C/C++ · benchmarks/results/matrix_.../COMPARISON.md",
    notes="Model choice matters — a lot. On Juliet C/C++, DeepSeek pushed "
          "precision from the 50% raw baseline to 84% (F1 88.5%, 82% of false "
          "positives removed) and was cheaper than GPT-4.1-mini, which was weaker "
          "here (precision 57%). Takeaway: benchmark on YOUR languages before "
          "committing to a model; the cheapest frontier-reasoning model was also "
          "the most accurate on memory-safety C/C++. Cost is real provider-"
          "reported USD; Ollama runs are $0.")

table_slide(
    "Two kinds of miss — know which you're fighting",
    ["", "Detection miss", "Triage miss"],
    [["What", "No rule fired at all", "Rule fired, but verified FP / NMD"],
     ["Cause", "Coverage gap (no query/pattern)", "Weak context or LLM error"],
     ["Example", "dvcp integer-overflow chain, div-by-zero", "Real bug downgraded to Needs-More-Data"],
     ["Remedy", "Add custom rule (--profile full)", "More context, more iterations, better model"],
     ["Measure", "Recall vs ground truth", "Precision of verified verdicts"]],
    col_widths=[1.6, 5.15, 5.2],
    notes="Critical distinction for anyone evaluating the tool. A detection miss "
          "is a SAST coverage gap — the LLM never even saw it. A triage miss is a "
          "verification error — it was found but mislabeled. They have completely "
          "different fixes: write a rule vs. improve context/model. Score them "
          "separately.")

table_slide(
    "Observed coverage on the benchmark repos",
    ["Repo (lang)", "Detected well", "Missed (expected)"],
    [["dvcp (C)", "double-free, UAF, OOB read/write", "int-overflow chain, div-by-zero, leak, exhaustion"],
     ["insecure-coding-examples (C++)", "unsafe fns, format string, UAF, double-free, type-confusion", "C++ lifetime/UB: dangling string_view, removed memset, RAII"],
     ["insecure-cpp-dojo (C++)", "≈1/14 (Semgrep only)", "most — CodeQL build failed → recall collapsed"]],
    col_widths=[3.0, 4.6, 4.35],
    caption="Source: docs/benchmarks/ground-truth-baselines.md (observed-coverage notes).",
    notes="Concrete, honest numbers. The dojo row is the cautionary tale: when "
          "its CMake build defeated CodeQL tracing, recall fell to ~1 of 14 bugs. "
          "That single row justifies the whole 'keep the build green' message and "
          "why we picked reliably-building repos for the workshop.")

bullets_slide(
    "Build health drives memory-safety recall",
    [("CodeQL needs to OBSERVE the build for C/C++/Java/Go/C#.", 0, INK, True),
     ("Build fails → database incomplete → semantic queries find little.", 0, ROSE, True),
     ("Semgrep still runs (no build) — but only syntactic coverage.", 0, INK, True),
     ("Practical: get a clean build first; use --tool both so something always runs.", 0, EMERALD, True),
     ("Verify the DB step in prepare logs before trusting low finding counts.", 1, GREY, False)],
    notes="Expand the limitation into actionable advice. The failure is silent — "
          "you get a report with few findings and might think the code is clean, "
          "when really CodeQL never built a database. Always check the prepare "
          "step succeeded, and use --tool both so Semgrep provides a floor "
          "of coverage even when CodeQL can't build.")

table_slide(
    "Cost, latency & best practice",
    ["Lever", "Why it matters", "Recommendation"],
    [["--limit N", "Tokens + time scale with findings", "Start with 5–10 while learning"],
     ["model choice", "Big models cost more per finding", "Cheap/local first, escalate on hard cases"],
     ["--max-iterations", "Multi-turn = more calls", "Default 10; lower for quick passes"],
     ["-j / --jobs", "Parallelism vs rate limits", "Raise to speed up; back off on 429s"],
     ["re-run verify only", "Don't rebuild the DB", "Iterate on triage cheaply"]],
    col_widths=[2.7, 4.6, 4.65],
    notes="Practical knobs. The headline: LLM triage is not free — it spends "
          "tokens and wall-clock per finding, more with multi-turn. Use --limit "
          "and a cheap model to iterate, escalate only on the hard findings, and "
          "re-run just verify (reusing the SARIF + DB) rather than the whole "
          "pipeline. Closes the loop back to 'assistant, not oracle.'")

# Closing
s = content_slide("Key takeaway")
textbox(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.2),
        [[("VulnHunterX is an assistant, not an oracle.", 32, NAVY_TX, True, False)],
         [("", 10, INK, False, False)],
         [("It makes a human reviewer dramatically faster.", 22, GREY, False, False)],
         [("It does not replace one. Now — let's run it.", 22, GREY, False, False)]],
        align=LEFT)
pill(s, Inches(0.9), Inches(5.1), Inches(4.6), Inches(0.7),
     "→  Hands-on workshop (30 min)", ORANGE, size=15)
_notes(s, "Closing thought, then transition to the workshop: scan a real "
          "vulnerable C++ repo (insecure-coding-examples) with a free Ollama "
          "Cloud model, score it against the published ground truth, and take "
          "dvcp home as homework.")

# --------------------------------------------------------------------------- #
prs.save(OUT_PATH)
print(f"Wrote {OUT_PATH}")
print(f"  slides: {len(prs.slides)}  (target ~120 min)")
print(f"  custom CodeQL queries: {CODEQL_TOTAL}  (c/cpp: {CODEQL_CPP})")
print(f"  custom Semgrep rules:  {SEMGREP_TOTAL}  (c/cpp: {SEMGREP_CPP})")
print(f"  guided-question banks: {len(QUESTIONS)}  ({QUESTIONS_TOTAL} entries)")
