#!/usr/bin/env python3
"""Generate the 60-minute VulnHunterX presentation as a .pptx file.

Run from anywhere:  python docs/presentation/generate_deck.py
Requires:           pip install python-pptx

A modern, diagram-driven deck: vector flowcharts (pipeline, the multi-turn
verification loop), a noise->signal funnel, card grids, a mock verdict card and
stat cards. Rule/coverage counts are computed at build time from the repo's
config/ tree so the deck never drifts. Conceptual content is anchored to
README.md and docs/context-extraction-flow.md.
"""
from __future__ import annotations

import glob
import math
import os
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# --------------------------------------------------------------------------- #
# Paths
# --------------------------------------------------------------------------- #
HERE = os.path.dirname(os.path.abspath(__file__))      # docs/presentation/
REPO = os.path.dirname(os.path.dirname(HERE))          # repo root
CONFIG = os.path.join(REPO, "config")
OUT_PATH = os.path.join(HERE, "VulnHunterX-60min.pptx")

# --------------------------------------------------------------------------- #
# Theme  — modern slate / indigo / cyan with orange brand accent
# --------------------------------------------------------------------------- #
DARK = RGBColor(0x0B, 0x11, 0x20)      # slate-950
DARK2 = RGBColor(0x1B, 0x29, 0x47)     # gradient end
GHOST = RGBColor(0x1C, 0x2C, 0x4A)     # watermark on dark
INK = RGBColor(0x0F, 0x17, 0x2A)       # slate-900 body text
NAVY_TX = RGBColor(0x10, 0x2A, 0x4A)   # heading on light
GREY = RGBColor(0x52, 0x5F, 0x70)      # sub-text
MUTE = RGBColor(0x90, 0x9C, 0xAB)      # footers / captions
BORDER = RGBColor(0xDD, 0xE4, 0xEC)    # card border
BG = RGBColor(0xF6, 0xF8, 0xFB)        # content background
CARD = RGBColor(0xFF, 0xFF, 0xFF)
TRACK = RGBColor(0xE2, 0xE8, 0xF0)     # progress track

ORANGE = RGBColor(0xF9, 0x73, 0x16)    # brand accent
INDIGO = RGBColor(0x63, 0x66, 0xF1)    # flow accent
CYAN = RGBColor(0x06, 0xB6, 0xD4)
EMERALD = RGBColor(0x10, 0xB9, 0x81)   # positive / TP
ROSE = RGBColor(0xF4, 0x3F, 0x5E)      # negative / FP
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
SLATE = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# code block palette
CODE_BG = RGBColor(0x0D, 0x17, 0x26)
CODE_FG = RGBColor(0xCC, 0xD8, 0xE6)
CODE_COMMENT = RGBColor(0x6E, 0x8B, 0xA3)
CODE_CMD = RGBColor(0x4F, 0xC9, 0xA8)
CODE_FLAG = RGBColor(0xF5, 0x9E, 0x0B)

BODY_FONT = "Segoe UI"
MONO_FONT = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

LEFT, CENTER, RIGHT = PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.RIGHT
TOP, MID = MSO_ANCHOR.TOP, MSO_ANCHOR.MIDDLE
RR, RECT, OVAL = MSO_SHAPE.ROUNDED_RECTANGLE, MSO_SHAPE.RECTANGLE, MSO_SHAPE.OVAL

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# --------------------------------------------------------------------------- #
# Build-time counters
# --------------------------------------------------------------------------- #
def count_codeql_queries() -> dict[str, int]:
    out: dict[str, int] = {}
    base = os.path.join(CONFIG, "codeql-custom")
    for lang in sorted(os.listdir(base)):
        d = os.path.join(base, lang)
        if not os.path.isdir(d):
            continue
        n = len(glob.glob(os.path.join(d, "**", "*.ql"), recursive=True))
        if n:
            out[lang] = n
    return out


def count_semgrep_rules() -> dict[str, int]:
    out: dict[str, int] = {}
    for path in sorted(glob.glob(os.path.join(CONFIG, "semgrep-custom", "*.yaml"))):
        lang = os.path.splitext(os.path.basename(path))[0]
        with open(path, encoding="utf-8") as fh:
            n = sum(1 for line in fh if re.match(r"\s*-\s+id:", line))
        out[lang] = n
    return out


def count_guided_questions() -> dict[str, int]:
    out: dict[str, int] = {}
    for path in sorted(glob.glob(os.path.join(CONFIG, "prompts", "*questions*.yaml"))):
        name = os.path.basename(path)
        with open(path, encoding="utf-8") as fh:
            n = sum(1 for line in fh if re.match(r"^[^\s#][^:]*:\s*$", line.rstrip()))
        out[name] = n
    return out


CODEQL = count_codeql_queries()
SEMGREP = count_semgrep_rules()
QUESTIONS = count_guided_questions()
CODEQL_TOTAL = sum(CODEQL.values())
SEMGREP_TOTAL = sum(SEMGREP.values())
QUESTIONS_TOTAL = sum(QUESTIONS.values())
QUESTION_BANKS = len(QUESTIONS)


# --------------------------------------------------------------------------- #
# Low-level helpers
# --------------------------------------------------------------------------- #
def _bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def _set(run, size, color, bold=False, italic=False, font=BODY_FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def shape(slide, mso, l, t, w, h, fill=None, line=None, line_w=1.0):
    sp = slide.shapes.add_shape(mso, int(l), int(t), int(w), int(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def text_in(sp, paras, align=CENTER, anchor=MID, wrap=True,
            ml=0.10, mt=0.05, mr=0.10, mb=0.05):
    """paras: list of paragraphs; each = list of (text,size,color,bold,italic)."""
    tf = sp.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(ml)
    tf.margin_right = Inches(mr)
    tf.margin_top = Inches(mt)
    tf.margin_bottom = Inches(mb)
    first = True
    for para in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(2)
        for txt, size, color, bold, italic in para:
            r = p.add_run()
            r.text = txt
            _set(r, size, color, bold=bold, italic=italic)
    return sp


def textbox(slide, l, t, w, h, paras, align=LEFT, anchor=TOP, wrap=True):
    tb = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
    return text_in(tb, paras, align=align, anchor=anchor, wrap=wrap,
                   ml=0, mt=0, mr=0, mb=0)


def soft_shadow(sp, alpha=22000, blur=80000, dist=38000):
    spPr = sp._element.spPr
    old = spPr.find(qn("a:effectLst"))
    if old is not None:
        spPr.remove(old)
    el = spPr.makeelement(qn("a:effectLst"), {})
    sh = el.makeelement(qn("a:outerShdw"),
                        {"blurRad": str(blur), "dist": str(dist),
                         "dir": "5400000", "rotWithShape": "0"})
    clr = sh.makeelement(qn("a:srgbClr"), {"val": "0F172A"})
    clr.append(clr.makeelement(qn("a:alpha"), {"val": str(alpha)}))
    sh.append(clr)
    el.append(sh)
    spPr.append(el)
    return sp


def fill_alpha(sp, pct):
    """Make an already-solid-filled shape translucent (pct = opacity %)."""
    fill = sp._element.spPr.find(qn("a:solidFill"))
    clr = fill.find(qn("a:srgbClr"))
    clr.append(clr.makeelement(qn("a:alpha"), {"val": str(int(pct * 1000))}))
    return sp


def grad(sp, c1, c2, angle=45):
    sp.fill.gradient()
    gs = sp.fill.gradient_stops
    gs[0].position, gs[0].color.rgb = 0.0, c1
    gs[1].position, gs[1].color.rgb = 1.0, c2
    try:
        sp.fill.gradient_angle = angle
    except Exception:
        pass
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def connect(slide, x1, y1, x2, y2, color, width=2.0, elbow=False, arrow=True):
    typ = MSO_CONNECTOR.ELBOW if elbow else MSO_CONNECTOR.STRAIGHT
    cn = slide.shapes.add_connector(typ, int(x1), int(y1), int(x2), int(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    cn.shadow.inherit = False
    if arrow:
        ln = cn.line._get_or_add_ln()
        ln.append(ln.makeelement(qn("a:tailEnd"),
                                 {"type": "triangle", "w": "med", "len": "med"}))
    return cn


def pill(slide, l, t, w, h, label, fill, color=WHITE, size=11):
    p = shape(slide, RR, l, t, w, h, fill=fill)
    text_in(p, [[(label, size, color, True, False)]])
    return p


# --------------------------------------------------------------------------- #
# Slide templates
# --------------------------------------------------------------------------- #
def title_slide(footer):
    s = prs.slides.add_slide(BLANK)
    grad(shape(s, RECT, 0, 0, SLIDE_W, SLIDE_H, fill=DARK), DARK, DARK2, angle=60)
    # decorative translucent blobs
    fill_alpha(shape(s, OVAL, Inches(9.4), Inches(-1.8), Inches(6.4), Inches(6.4),
                     fill=INDIGO), 14)
    fill_alpha(shape(s, OVAL, Inches(10.8), Inches(3.6), Inches(4.6), Inches(4.6),
                     fill=CYAN), 12)
    # logo mark
    text_in(soft_shadow(shape(s, RR, Inches(0.9), Inches(0.85), Inches(0.95),
                              Inches(0.95), fill=ORANGE)),
            [[("VX", 30, WHITE, True, False)]])
    textbox(s, Inches(2.0), Inches(1.02), Inches(7), Inches(0.7),
            [[("VulnHunterX", 22, WHITE, True, False)]])
    # kicker
    shape(s, RECT, Inches(0.92), Inches(2.7), Inches(0.55), Inches(0.12), fill=ORANGE)
    textbox(s, Inches(1.62), Inches(2.46), Inches(10), Inches(0.5),
            [[("LLM-VERIFIED SAST  ·  HANDS-ON WORKSHOP", 15, AMBER, True, False)]])
    # title + subtitle
    textbox(s, Inches(0.85), Inches(3.0), Inches(11.8), Inches(1.5),
            [[("Find real bugs.", 54, WHITE, True, False)],
             [("Skip the false alarms.", 54, CYAN, True, False)]])
    textbox(s, Inches(0.92), Inches(5.0), Inches(11), Inches(0.6),
            [[("Static analysis finds candidates — an LLM decides which are real.",
               20, MUTE, False, False)]])
    # stat chips
    chips = [(f"{CODEQL_TOTAL}", "custom CodeQL queries"),
             (f"{SEMGREP_TOTAL}", "custom Semgrep rules"),
             ("8", "languages  ·  5 profiles")]
    x = Inches(0.9)
    for big, small in chips:
        c = soft_shadow(shape(s, RR, x, Inches(5.75), Inches(3.55), Inches(0.7),
                              fill=GHOST))
        text_in(c, [[(big + "  ", 18, ORANGE, True, False),
                     (small, 12.5, WHITE, False, False)]], align=CENTER)
        x += Inches(3.75)
    textbox(s, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.5),
            [[(footer, 13, MUTE, False, False)]])
    return s


PART_TOTAL = 4


def section_slide(kicker, title):
    s = prs.slides.add_slide(BLANK)
    grad(shape(s, RECT, 0, 0, SLIDE_W, SLIDE_H, fill=DARK), DARK, DARK2, angle=30)
    idx = int("".join(ch for ch in kicker if ch.isdigit()) or 0)
    textbox(s, Inches(7.4), Inches(1.35), Inches(5.55), Inches(4.6),
            [[(str(idx).zfill(2), 250, GHOST, True, False)]], align=RIGHT,
            anchor=MID, wrap=False)
    shape(s, RECT, Inches(0.95), Inches(2.92), Inches(0.55), Inches(0.12), fill=ORANGE)
    textbox(s, Inches(1.65), Inches(2.68), Inches(8), Inches(0.55),
            [[(kicker.upper(), 18, AMBER, True, False)]])
    textbox(s, Inches(0.9), Inches(3.22), Inches(11), Inches(1.6),
            [[(title, 44, WHITE, True, False)]])
    shape(s, RECT, Inches(0.95), Inches(4.55), Inches(3.2), Inches(0.05), fill=ORANGE)
    # progress dots
    x = Inches(0.95)
    for i in range(1, PART_TOTAL + 1):
        col = ORANGE if i == idx else GHOST
        w = Inches(0.5) if i == idx else Inches(0.22)
        shape(s, RR, x, Inches(6.7), w, Inches(0.16), fill=col)
        x += w + Inches(0.16)
    return s


def content_slide(title, kicker=None):
    s = prs.slides.add_slide(BLANK)
    _bg(s, BG)
    shape(s, RECT, 0, 0, Inches(0.14), SLIDE_H, fill=ORANGE)
    textbox(s, Inches(0.55), Inches(0.34), Inches(10.6), Inches(0.85),
            [[(title, 26, NAVY_TX, True, False)]])
    shape(s, RECT, Inches(0.6), Inches(1.14), Inches(1.5), Inches(0.045), fill=ORANGE)
    textbox(s, Inches(9.7), Inches(0.42), Inches(3.05), Inches(0.4),
            [[("VulnHunterX", 12, MUTE, True, False)]], align=RIGHT)
    textbox(s, Inches(12.0), Inches(7.02), Inches(1.0), Inches(0.4),
            [[(str(len(prs.slides)), 11, MUTE, False, False)]], align=RIGHT)
    return s


def bullets_slide(title, bullets, notes="", top=Inches(1.45), left=Inches(0.7),
                  width=Inches(12.0), height=Inches(5.4)):
    s = content_slide(title)
    tb = s.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for text, level, color, bold in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(9)
        size = 22 - 3 * level
        if text:
            if level == 0 and bold:
                m = p.add_run()
                m.text = "▸  "
                _set(m, size, ORANGE, bold=True)
            elif level >= 1:
                m = p.add_run()
                m.text = "–  "
                _set(m, size, MUTE, bold=True)
        r = p.add_run()
        r.text = text
        _set(r, size, color or INK, bold=bold)
    if notes:
        _notes(s, notes)
    return s


def lay_cards(s, items, cols, top=1.55, card_h=1.12, gap=0.25, left=0.6,
              total_w=12.2, head_size=15.5, body_size=12.5):
    """items: list of dict(head, body, accent, pill, pill_color)."""
    gap_e = Inches(gap)
    card_w = int((Inches(total_w) - gap_e * (cols - 1)) / cols)
    ch = Inches(card_h)
    for i, it in enumerate(items):
        r, c = divmod(i, cols)
        x = int(Inches(left) + c * (card_w + gap_e))
        y = int(Inches(top) + r * (ch + gap_e))
        accent = it.get("accent", ORANGE)
        soft_shadow(shape(s, RR, x, y, card_w, ch, fill=CARD, line=BORDER))
        shape(s, RR, x, y, Inches(0.11), ch, fill=accent)
        paras = [[(it["head"], head_size, NAVY_TX, True, False)]]
        if it.get("body"):
            paras.append([(it["body"], body_size, GREY, False, False)])
        body_box = shape(s, RECT, x + Inches(0.22), y, card_w - Inches(0.34), ch)
        body_box.fill.background()
        body_box.line.fill.background()
        text_in(body_box, paras, align=LEFT, anchor=MID, ml=0.06, mr=0.06)
        if it.get("pill"):
            pw, ph = Inches(1.2), Inches(0.4)
            pill(s, x + card_w - pw - Inches(0.16), y + Inches(0.14), pw, ph,
                 it["pill"], it.get("pill_color", accent), size=11)
    return s


def card_grid_slide(title, items, cols, notes="", **kw):
    s = content_slide(title)
    lay_cards(s, items, cols, **kw)
    if notes:
        _notes(s, notes)
    return s


def _emit_code_line(p, line):
    if not line.strip():
        r = p.add_run()
        r.text = " "
        _set(r, 15.5, CODE_FG, font=MONO_FONT)
        return
    if line.lstrip().startswith("#"):
        r = p.add_run()
        r.text = line
        _set(r, 15.5, CODE_COMMENT, italic=True, font=MONO_FONT)
        return
    code_part, comment_part = line, ""
    idx = line.find("  #")
    if idx != -1:
        code_part, comment_part = line[:idx], line[idx:]
    for tok in re.split(r"(\s+)", code_part):
        if tok == "":
            continue
        r = p.add_run()
        r.text = tok
        if tok.startswith("--") or (len(tok) > 1 and tok[0] == "-" and tok[1].isalpha()):
            _set(r, 15.5, CODE_FLAG, font=MONO_FONT)
        elif tok in ("vuln-hunter-x", "python", "git", "uv", "cp", "cat", "source"):
            _set(r, 15.5, CODE_CMD, bold=True, font=MONO_FONT)
        else:
            _set(r, 15.5, CODE_FG, font=MONO_FONT)
    if comment_part:
        r = p.add_run()
        r.text = comment_part
        _set(r, 15.5, CODE_COMMENT, italic=True, font=MONO_FONT)


def code_slide(title, lines, notes="", caption=""):
    s = content_slide(title)
    top = Inches(1.5)
    if caption:
        textbox(s, Inches(0.7), Inches(1.26), Inches(12), Inches(0.5),
                [[(caption, 15, MUTE, False, True)]])
        top = Inches(1.78)
    box = soft_shadow(shape(s, RR, Inches(0.7), top, Inches(11.95),
                            int(SLIDE_H - top - Inches(0.55)), fill=CODE_BG))
    for i, col in enumerate((ROSE, AMBER, EMERALD)):
        shape(s, OVAL, Inches(0.95 + i * 0.28), top + Inches(0.22),
              Inches(0.16), Inches(0.16), fill=col)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = TOP
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.62)
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(3)
        _emit_code_line(p, line)
    if notes:
        _notes(s, notes)
    return s


def flow_node(s, l, t, w, h, head, sub, accent):
    soft_shadow(shape(s, RR, l, t, w, h, fill=CARD, line=BORDER))
    shape(s, RR, l, t, w, Inches(0.13), fill=accent)
    box = shape(s, RECT, l, t + Inches(0.16), w, h - Inches(0.16))
    box.fill.background()
    box.line.fill.background()
    paras = [[(head, 14.5, NAVY_TX, True, False)]]
    if sub:
        paras.append([(sub, 11, GREY, False, False)])
    text_in(box, paras, align=CENTER, anchor=MID, ml=0.06, mr=0.06)


def block_arrow(s, l, t, w, h, color):
    return shape(s, MSO_SHAPE.RIGHT_ARROW, l, t, w, h, fill=color)


# --------------------------------------------------------------------------- #
# SLIDES
# --------------------------------------------------------------------------- #

# 1 — Title
title_slide("60-minute workshop  ·  for developers  ·  CodeQL + Semgrep + LLM triage")
_notes(prs.slides[-1], """
Welcome. Four parts: what VulnHunterX is (features, architecture, methodology and
benchmark results), how to install and use it, a real worked example on dvpwa, then
homework on dvcp. Audience: developers — no security expertise needed.
""")

# 2 — Agenda (4 parts)
card_grid_slide("Agenda — 4 parts", [
    {"head": "1 · VulnHunterX introduction", "body": "Features · architecture · methodology · results", "accent": INDIGO},
    {"head": "2 · How to install & use", "body": "Toolchain, install, configure, run a scan", "accent": CYAN},
    {"head": "3 · Real example — dvpwa", "body": "Scan a vulnerable Python web app, verify the verdicts", "accent": ORANGE},
    {"head": "4 · Homework — dvcp", "body": "Scan the Damn Vulnerable C Program at home", "accent": EMERALD},
], cols=2, top=1.95, card_h=1.55, gap=0.3, head_size=19, body_size=13.5,
   notes="Four parts: understand it, install it, watch it work on dvpwa, then do it yourself on dvcp.")

# 3 — SECTION
section_slide("Part 1", "VulnHunterX — features, architecture, methodology, results")

# 4 — Problem funnel (diagram)
s = content_slide("From noise to signal")
textbox(s, Inches(0.6), Inches(1.25), Inches(12), Inches(0.5),
        [[("SAST over-approximates by design — the cost is human triage, not the scan.",
           17, GREY, False, True)]])
# funnel: 3 stacked centered boxes, decreasing width, down arrows
cx = Inches(5.35)
specs = [
    (Inches(9.0), "Hundreds of raw SAST findings", "every program point that *might* be vulnerable", ROSE),
    (Inches(6.9), "LLM verification — Vulnhalla guided triage", "evidence-anchored, multi-turn reasoning", INDIGO),
    (Inches(4.5), "A handful of real, exploitable bugs", "what your team should actually fix", EMERALD),
]
y = Inches(2.05)
bh = Inches(1.15)
for i, (w, head, sub, col) in enumerate(specs):
    x = int(cx - w / 2)
    n = soft_shadow(shape(s, RR, x, y, w, bh, fill=col))
    text_in(n, [[(head, 18, WHITE, True, False)],
                [(sub, 12.5, WHITE, False, True)]], align=CENTER)
    if i < len(specs) - 1:
        block_arrow_down = shape(s, MSO_SHAPE.DOWN_ARROW, int(cx - Inches(0.35)),
                                 int(y + bh + Inches(0.04)), Inches(0.7), Inches(0.45),
                                 fill=SLATE)
    y = int(y + bh + Inches(0.55))
# side callout
side = soft_shadow(shape(s, RR, Inches(10.3), Inches(2.05), Inches(2.55), Inches(3.85),
                         fill=DARK))
text_in(side, [[("Real bugs get", 14, WHITE, False, False)],
               [("buried.", 22, AMBER, True, False)],
               [("", 8, WHITE, False, False)],
               [("Teams stop", 14, WHITE, False, False)],
               [("reading the", 14, WHITE, False, False)],
               [("report.", 22, ROSE, True, False)]], align=CENTER)
_notes(s, """
SAST flags everything that might be vulnerable. The dominant production cost is
a human triaging the results — alert fatigue means real bugs get buried.
VulnHunterX inserts an LLM verification stage to cut the noise.
""")

# 5 — one-sentence statement card
s = content_slide("VulnHunterX in one sentence")
big = soft_shadow(shape(s, RR, Inches(0.9), Inches(1.7), Inches(11.5), Inches(2.5),
                        fill=CARD, line=BORDER))
shape(s, RR, Inches(0.9), Inches(1.7), Inches(0.16), Inches(2.5), fill=ORANGE)
text_in(big, [
    [("A Python framework that pairs ", 26, INK, False, False),
     ("SAST engines", 26, INDIGO, True, False),
     (" with ", 26, INK, False, False),
     ("multi-turn LLM verification", 26, ORANGE, True, False)],
    [("to suppress false positives — the ", 26, INK, False, False),
     ("Vulnhalla", 26, CYAN, True, False),
     (" guided-question methodology.", 26, INK, False, False)],
], align=CENTER, ml=0.4, mr=0.3)
for i, (h, sub, col) in enumerate([
    ("Cite evidence", "the LLM must point at real code", INDIGO),
    ("Pull context", "callers, structs, globals on demand", CYAN),
    ("Reason in turns", "not pattern-match one snippet", ORANGE)]):
    x = int(Inches(0.9) + i * Inches(3.95))
    c = soft_shadow(shape(s, RR, x, Inches(4.55), Inches(3.7), Inches(1.5),
                          fill=CARD, line=BORDER))
    text_in(c, [[(h, 17, col, True, False)], [(sub, 12.5, GREY, False, False)]],
            align=CENTER)
_notes(s, """
Vulnhalla is a CyberArk methodology (reference at the end). The distinction vs
"just ask ChatGPT": the LLM answers a rule-specific checklist BEFORE a verdict
and fetches real surrounding code instead of guessing from one snippet.
""")

# 6 — Key feature cards
card_grid_slide("Key features", [
    {"head": "8 languages", "body": "C, C++, Python, JavaScript, PHP, Java, Go, C#", "accent": INDIGO},
    {"head": "3 SAST engines", "body": "CodeQL · Semgrep · OpenGrep  (--tool …)", "accent": CYAN},
    {"head": "5 rule profiles", "body": "standard → extended → maximum → … → full", "accent": AMBER},
    {"head": "LLM providers", "body": "OpenAI · Anthropic · Ollama — via LiteLLM", "accent": ORANGE},
    {"head": "Multi-turn verification", "body": "context expansion → verdict + confidence", "accent": EMERALD},
    {"head": "Flexible inputs", "body": "Git URL · local dir · batch repos.yaml", "accent": SLATE},
    {"head": "Reports", "body": "Markdown EN/VI — summary + per-finding detail", "accent": INDIGO},
    {"head": "Fuzzing (C/C++)", "body": "libFuzzer / Atheris / Jazzer + crash triage", "accent": ROSE},
], cols=2, top=1.55, card_h=1.18, gap=0.28,
   notes="""
C# is the newest language — buildless CodeQL extraction, no dotnet build needed.
Most of the room will have at least one of these 8 languages in their stack.
""")

# 7 — Pipeline flowchart (diagram)
s = content_slide("The pipeline — 5 steps, 1 command")
nodes = [
    ("Source", "git / local / repos.yaml", SLATE),
    ("Static Analysis", "CodeQL · Semgrep · OpenGrep", INDIGO),
    ("SARIF findings", "rule · file · line · CWE", CYAN),
    ("LLM Verification", "guided Qs · multi-turn", ORANGE),
    ("Verdicts", "TP / FP / NMD + confidence", EMERALD),
]
nw, nh = Inches(2.12), Inches(1.55)
gap = Inches(0.36)
x = Inches(0.55)
y = Inches(2.2)
for i, (head, sub, col) in enumerate(nodes):
    flow_node(s, x, y, nw, nh, head, sub, col)
    if i < len(nodes) - 1:
        block_arrow(s, int(x + nw + Inches(0.02)), int(y + nh / 2 - Inches(0.22)),
                    int(gap - Inches(0.04)), Inches(0.44), ORANGE)
    x = int(x + nw + gap)
# command pills under stages
labels = ["prepare", "analyze", "(SARIF)", "verify", "report"]
x = Inches(0.55)
for i, lab in enumerate(labels):
    col = INK if lab.startswith("(") else NAVY_TX
    fillc = TRACK if lab.startswith("(") else DARK
    txc = GREY if lab.startswith("(") else WHITE
    pill(s, int(x + nw / 2 - Inches(0.75)), int(y + nh + Inches(0.25)),
         Inches(1.5), Inches(0.42), lab, fillc, color=txc, size=12)
    x = int(x + nw + gap)
# fuzz sub-row
fz = soft_shadow(shape(s, RR, Inches(0.55), Inches(5.55), Inches(12.2), Inches(1.15),
                       fill=RGBColor(0xEE, 0xF2, 0xF7), line=BORDER))
textbox(s, Inches(0.8), Inches(5.66), Inches(11.6), Inches(0.4),
        [[("+ 4 optional fuzz stages  (C/C++ only)", 13, GREY, True, False)]])
fstages = ["build-sanitized", "extract-fuzz-context", "generate-fuzz-drivers", "fuzz-run"]
x = Inches(0.85)
for i, fs in enumerate(fstages):
    pill(s, x, Inches(6.08), Inches(2.6), Inches(0.46), fs, AMBER, size=11.5)
    if i < len(fstages) - 1:
        textbox(s, int(x + Inches(2.62)), Inches(6.12), Inches(0.3), Inches(0.4),
                [[("→", 16, AMBER, True, False)]], align=CENTER)
    x = int(x + Inches(2.92))
_notes(s, """
README "Pipeline Stages". prepare → analyze → verify → report are the core four;
`scan` runs them in one command. Fuzz stages only apply to C/C++ and only on
Linux/macOS (covered under Windows support).
""")

# 8 — 4-stage stepper (diagram)
s = content_slide("The 4 core stages")
stages = [
    ("1", "prepare", "Clone + CodeQL DB +\ncontext CSVs", INDIGO),
    ("2", "analyze", "Run engines →\nSARIF findings", CYAN),
    ("3", "verify", "LLM triage →\nJSON verdicts", ORANGE),
    ("4", "report", "Markdown report\n(EN / VI)", EMERALD),
]
# connecting track
shape(s, RECT, Inches(1.9), Inches(2.85), Inches(9.5), Inches(0.06), fill=TRACK)
cw = Inches(2.7)
x = Inches(0.7)
for num, cmd, desc, col in stages:
    cx = int(x + cw / 2)
    # number circle
    circ = shape(s, OVAL, int(cx - Inches(0.45)), Inches(2.4), Inches(0.9), Inches(0.9),
                 fill=col)
    soft_shadow(circ)
    text_in(circ, [[(num, 26, WHITE, True, False)]])
    # card below
    cd = soft_shadow(shape(s, RR, int(cx - Inches(1.15)), Inches(3.7), Inches(2.3),
                           Inches(1.85), fill=CARD, line=BORDER))
    paras = [[(cmd, 19, col, True, False)]]
    for ln in desc.split("\n"):
        paras.append([(ln, 13, GREY, False, False)])
    text_in(cd, paras, align=CENTER, anchor=MID)
    x = int(x + cw + Inches(0.43))
_notes(s, """
Run stages individually or all at once with `scan`. Stages 2-3 accept
--local-path to run directly on any directory without a config entry.
""")

# 9 — Multi-turn loop flowchart (diagram)
s = content_slide("The Vulnhalla multi-turn loop")
# nodes
flow_node(s, Inches(0.6), Inches(1.65), Inches(2.2), Inches(1.2),
          "SARIF finding", "rule + code snippet", SLATE)
block_arrow(s, Inches(2.86), Inches(2.05), Inches(0.42), Inches(0.4), ORANGE)
flow_node(s, Inches(3.36), Inches(1.65), Inches(2.2), Inches(1.2),
          "Guided questions", "rule-specific checklist", INDIGO)
block_arrow(s, Inches(5.62), Inches(2.05), Inches(0.42), Inches(0.4), ORANGE)
flow_node(s, Inches(6.12), Inches(1.65), Inches(2.3), Inches(1.2),
          "LLM reasoning", "answer + decide", ORANGE)
# decision diamond
dia = soft_shadow(shape(s, MSO_SHAPE.DIAMOND, Inches(9.0), Inches(1.45), Inches(2.0),
                        Inches(1.6), fill=CARD, line=BORDER))
text_in(dia, [[("Verdict", 14, NAVY_TX, True, False)], [("ready?", 14, NAVY_TX, True, False)]])
# yes -> verdict box
connect(s, Inches(11.0), Inches(2.25), Inches(11.55), Inches(2.25), EMERALD, 2.25)
textbox(s, Inches(11.0), Inches(1.8), Inches(1.6), Inches(0.35),
        [[("yes", 12, EMERALD, True, False)]], align=CENTER)
vbox = soft_shadow(shape(s, RR, Inches(11.5), Inches(1.65), Inches(1.45), Inches(1.2),
                         fill=EMERALD))
text_in(vbox, [[("Verdict", 13, WHITE, True, False)], [("TP / FP", 13, WHITE, True, False)]])
# no -> fetch context (down then back into the box)
connect(s, Inches(10.0), Inches(3.05), Inches(9.3), Inches(4.2), AMBER, 2.25, elbow=True)
textbox(s, Inches(10.15), Inches(3.25), Inches(2.2), Inches(0.35),
        [[("needs more data", 12, AMBER, True, False)]])
fbox = soft_shadow(shape(s, RR, Inches(3.36), Inches(4.2), Inches(6.4), Inches(1.25),
                         fill=DARK))
text_in(fbox, [[("Fetch context", 16, AMBER, True, False)],
               [("callers · all_callers · callees · struct · global · macro ·",
                 12, WHITE, False, False)],
               [("typedef · enum · free_sites · destructor · field_writes",
                 12, WHITE, False, False)]], align=CENTER)
# loop back arrow fetch -> LLM
connect(s, Inches(6.5), Inches(4.2), Inches(7.0), Inches(2.85), CYAN, 2.25, elbow=True)
textbox(s, Inches(4.9), Inches(5.55), Inches(8), Inches(0.4),
        [[("Loop up to ", 14, GREY, False, False),
          ("max_iterations (default 3)", 14, ORANGE, True, False),
          ("  until the LLM commits.", 14, GREY, False, False)]], align=CENTER)
_notes(s, """
Anchor: docs/context-extraction-flow.md. Two-phase context: pre-fetched upfront
(function-name keyed) vs reactive (only when the LLM asks). 11 context types;
the last three are C/C++-specific (use-after-free / TOCTOU). The loop
deduplicates requests and can force a decision if the LLM stalls.
""")

# 10 — Guided questions: the contract (YAML anatomy)
code_slide("Guided questions: the per-language contract", [
    "# config/prompts/python_questions.yaml  — one bank per language",
    "py/sql-injection:",
    "  short_description: User input concatenated into SQL query",
    "  questions:                          # answered IN ORDER, before any verdict",
    "    - Quote the EXACT sink (cursor.execute, line N) and name the variable.",
    "    - List EVERY assignment to that variable on each path, with line numbers.",
    "    - Does the value derive from request.args/form/json, or a safe source?",
    "    - Name the SPECIFIC defense — parameterised execute(sql, params), ORM.",
    "  context_hint: Trace request input -> transformations -> query execution.",
    "  additional_context: [caller, callees]   # fetched on demand mid-loop",
    "  min_iterations: 2                        # force a second reasoning pass",
], notes="""
Each *_questions.yaml is a per-language bank of rule-specific checklists, loaded by
QuestionsLoader (src/vuln_hunter_x/questions/loader.py). The questions are written to
force EXACT sink quoting, full assignment tracing, and concrete defense naming — the
opposite of pattern-matching a snippet. additional_context drives the multi-turn fetch
loop; min_iterations forces a second look on FP-prone classes.
""")

# 11 — Same CWE, every language
s = card_grid_slide("Same CWE, every language — CWE-89 / SQL injection", [
    {"head": "py/sql-injection", "body": "sink: cursor.execute(...)  ·  defense: parameterised execute(sql, params) / ORM .filter()", "accent": INDIGO},
    {"head": "java/sql-injection", "body": "sink: Statement.executeQuery(...)  ·  defense: PreparedStatement with bound params", "accent": CYAN},
    {"head": "js/sql-injection", "body": "sink: db.query(`...${x}`)  ·  defense: parameterised driver call / placeholders", "accent": AMBER},
    {"head": "php/sql-injection", "body": "sink: mysqli_query(...)  ·  defense: prepared statements / PDO bound params", "accent": EMERALD},
], cols=2, top=1.6, card_h=1.6, gap=0.3, head_size=18, body_size=13.5)
strip = soft_shadow(shape(s, RR, Inches(0.6), Inches(5.45), Inches(12.15), Inches(1.4),
                          fill=DARK))
text_in(strip, [
    [("Same ordered checklist shape", 16, AMBER, True, False),
     (" — language-specific sinks, sanitizers & framework defenses.", 14, WHITE, False, False)],
    [(f"{QUESTION_BANKS} banks", 15, CYAN, True, False),
     ("  cpp · cs · go · java · js · php · python  (+ default fallback)", 13.5, WHITE, False, False)],
    [("No SQL bank for C/C++", 13.5, ROSE, True, False),
     ("  — instead C/C++ gets memory-safety banks: use-after-free, integer-overflow.", 13, MUTE, False, True)],
], align=CENTER)
_notes(s, """
A common CWE like SQL injection (CWE-89) gets a bank in every framework language. The
checklists share structure but name the REAL APIs of each language — that is what makes
the LLM cite concrete evidence. C/C++ has no native SQL, so there is no SQL bank there;
the C/C++ banks target memory-safety classes instead (use-after-free, integer-overflow).
""")

# 12 — Routing a finding to its questions
s = content_slide("Routing a finding to its question bank")
textbox(s, Inches(0.6), Inches(1.22), Inches(12), Inches(0.45),
        [[("QuestionsLoader resolves every finding to a bank through ordered fallback tiers.",
           15, GREY, False, True)]])
tiers = [
    ("exact", "ruleId matches a bank key — CodeQL: py/sql-injection", EMERALD),
    ("normalized · prefix · lang_prefix", "tolerate id variants & same-language partials", CYAN),
    ("cwe", "Semgrep has no CodeQL-style id → route by CWE tag via cwe_question_map\n(config/rule_categories.yaml):  CWE-89 → \"sql-injection\" → py/sql-injection", INDIGO),
    ("default → generic", "fall back to default_questions.yaml, then an 8-step generic checklist", SLATE),
]
y = Inches(1.85)
for head, body, col in tiers:
    h = Inches(0.92) if "\n" in body else Inches(0.72)
    cd = soft_shadow(shape(s, RR, Inches(0.6), y, Inches(12.15), h, fill=CARD, line=BORDER))
    shape(s, RR, Inches(0.6), y, Inches(0.13), h, fill=col)
    paras = [[(head + ":   ", 15, col, True, False),
              (body.split("\n")[0], 13.5, GREY, False, False)]]
    for extra in body.split("\n")[1:]:
        paras.append([(extra, 12.5, GREY, False, False)])
    text_in(cd, paras, align=LEFT, anchor=MID, ml=0.3)
    y = int(y + h + Inches(0.16))
ovr = soft_shadow(shape(s, RR, Inches(0.6), y, Inches(12.15), Inches(1.2), fill=DARK))
text_in(ovr, [
    [("Multi-iteration override", 15, AMBER, True, False),
     ("   taint CWEs (22/78/79/89/90/94/611/918…) + access-control (200/264/287/862/863…)",
      13, WHITE, False, False)],
    [("force min_iterations = 2.  ", 13.5, WHITE, False, False),
     ("OWASP-python CWE-22: 57.1% (1 iter) → 95.8% (2 iter).", 14, EMERALD, True, False),
     ("  Taint gated to framework langs; access-control all langs.", 12.5, MUTE, False, True)],
], align=CENTER)
_notes(s, """
This is why a Semgrep finding (no CodeQL id) still lands on a rule-specific checklist —
it routes by its CWE tag. And why path-traversal / SQLi verdicts get a forced second
pass: the _CWE_MIN_ITERATIONS_OVERRIDE table (loader.py) raises min_iterations to 2 for
FP-prone classes. Measured on the OWASP-python benchmark: CWE-22 accuracy jumped from
57.1% to 95.8% with the second iteration. Taint overrides are gated to framework
languages (no benefit on C); access-control applies everywhere.
""")

# 13 — Verdict card mock (image-like)
s = content_slide("What a verdict looks like")
card = soft_shadow(shape(s, RR, Inches(0.6), Inches(1.5), Inches(7.4), Inches(5.0),
                         fill=CARD, line=BORDER))
shape(s, RR, Inches(0.6), Inches(1.5), Inches(7.4), Inches(0.7), fill=DARK)
textbox(s, Inches(0.85), Inches(1.6), Inches(5.5), Inches(0.5),
        [[("Finding #12  ·  cpp/use-after-free", 15, WHITE, True, False)]])
pill(s, Inches(6.25), Inches(1.62), Inches(1.55), Inches(0.46), "TRUE POSITIVE",
     EMERALD, size=11.5)
# confidence bar
textbox(s, Inches(0.9), Inches(2.45), Inches(6), Inches(0.4),
        [[("Confidence", 14, GREY, True, False)]])
shape(s, RR, Inches(0.9), Inches(2.9), Inches(6.7), Inches(0.38), fill=TRACK)
shape(s, RR, Inches(0.9), Inches(2.9), Inches(5.76), Inches(0.38), fill=EMERALD)
textbox(s, Inches(0.95), Inches(2.92), Inches(6.6), Inches(0.35),
        [[("High · 0.86", 13, WHITE, True, False)]])
# reasoning
textbox(s, Inches(0.9), Inches(3.55), Inches(6.8), Inches(0.4),
        [[("Reasoning", 14, GREY, True, False)]])
textbox(s, Inches(0.9), Inches(3.95), Inches(6.9), Inches(2.2),
        [[("Tainted length from recv() reaches memcpy() with no bounds check;",
           13.5, INK, False, False)],
         [("attacker-controlled. Confirmed by fetching caller parse_packet(),",
           13.5, INK, False, False)],
         [("which frees buf before the copy. Reachable on the error path.",
           13.5, INK, False, False)]])
pill(s, Inches(0.9), Inches(5.95), Inches(1.5), Inches(0.42), "2 turns", INDIGO, size=12)
pill(s, Inches(2.55), Inches(5.95), Inches(2.2), Inches(0.42), "CWE-416", AMBER, size=12)
# right: field list
fields = [("verdict", "TP / FP / Needs More Data"),
          ("confidence", "High / Medium / Low + score"),
          ("reasoning", "evidence-anchored, cites code"),
          ("iterations", "turns the LLM took"),
          ("data_flow", "source→sink trace (CodeQL)")]
y = Inches(1.7)
for k, v in fields:
    fc = shape(s, RR, Inches(8.3), y, Inches(4.45), Inches(0.84), fill=CARD, line=BORDER)
    text_in(fc, [[(k, 14, ORANGE, True, False), ("   " + v, 12.5, GREY, False, False)]],
            align=LEFT, ml=0.2)
    y = int(y + Inches(0.97))
_notes(s, """
Verdicts are JSON, rendered into report.md. Confidence is downgraded when the
model reasons in pattern-matching language without citing actual code — a
deliberate guard against confident hallucination.
""")

# 14 — Coverage stat cards (diagram)
s = content_slide("Rules & coverage", )
textbox(s, Inches(0.6), Inches(1.25), Inches(12), Inches(0.45),
        [[("Counted from config/ at build time — not copied from the README.",
           15, MUTE, False, True)]])
stats = [(str(CODEQL_TOTAL), "custom CodeQL queries", INDIGO,
          " · ".join(f"{k} {v}" for k, v in CODEQL.items())),
         (str(SEMGREP_TOTAL), "custom Semgrep rules", CYAN,
          " · ".join(f"{k} {v}" for k, v in SEMGREP.items())),
         (f"~{QUESTIONS_TOTAL}", f"guided-question sets · {QUESTION_BANKS} banks", ORANGE,
          "rule-specific checklists")]
x = Inches(0.6)
for big, label, col, sub in stats:
    cd = soft_shadow(shape(s, RR, x, Inches(1.9), Inches(3.9), Inches(2.5),
                           fill=CARD, line=BORDER))
    shape(s, RR, x, Inches(1.9), Inches(3.9), Inches(0.14), fill=col)
    text_in(cd, [[(big, 52, col, True, False)],
                 [(label, 14.5, NAVY_TX, True, False)],
                 [(sub, 11.5, GREY, False, False)]], align=CENTER)
    x = int(x + Inches(4.05))
# bottom strip
strip = shape(s, RR, Inches(0.6), Inches(4.7), Inches(12.15), Inches(1.55),
              fill=DARK)
soft_shadow(strip)
text_in(strip, [
    [("5 rule profiles", 16, AMBER, True, False),
     ("   standard → extended → maximum → extended-registry → full", 14, WHITE, False, False)],
    [("124-entry CWE routing map", 16, CYAN, True, False),
     ("   routes findings to the right guided questions", 14, WHITE, False, False)],
    [("~5–10× more rules", 16, EMERALD, True, False),
     ("   from --profile standard to --profile full", 14, WHITE, False, False)],
], align=CENTER)
_notes(s, f"""
These numbers are computed by generate_deck.py from the live config/ tree (the
README summary table is stale). CodeQL custom = {CODEQL_TOTAL}, Semgrep custom =
{SEMGREP_TOTAL}, guided-question rule-sets ~= {QUESTIONS_TOTAL} across
{QUESTION_BANKS} banks. Custom rules are only layered on under --profile full.
""")

# 15 — Result: precision vs raw SAST (NEW)
s = content_slide("Result — false positives cut by 91%")
textbox(s, Inches(0.6), Inches(1.25), Inches(12), Inches(0.45),
        [[("OWASP-Python benchmark  ·  DeepSeek  ·  300 findings  ·  raw SAST → VulnHunterX.",
           15, MUTE, False, True)]])
res1 = [("87.3%", "Precision", INDIGO, "up from 37.7% (raw SAST)"),
        ("92.4%", "F1 score", CYAN, "up from 54.7%"),
        ("91.4%", "False positives removed", EMERALD, "of all raw-SAST false alarms")]
x = Inches(0.6)
for big, label, col, sub in res1:
    cd = soft_shadow(shape(s, RR, x, Inches(1.9), Inches(3.9), Inches(2.5),
                           fill=CARD, line=BORDER))
    shape(s, RR, x, Inches(1.9), Inches(3.9), Inches(0.14), fill=col)
    text_in(cd, [[(big, 52, col, True, False)],
                 [(label, 14.5, NAVY_TX, True, False)],
                 [(sub, 11.5, GREY, False, False)]], align=CENTER)
    x = int(x + Inches(4.05))
strip = soft_shadow(shape(s, RR, Inches(0.6), Inches(4.7), Inches(12.15), Inches(1.55), fill=DARK))
text_in(strip, [
    [("Recall held at 98.2%", 16, AMBER, True, False),
     ("   real bugs are kept, not thrown away with the noise", 14, WHITE, False, False)],
    [("True-positive preservation 97.4%", 16, CYAN, True, False),
     ("   only 1.0% of findings end up 'Needs More Data'", 14, WHITE, False, False)],
    [("Raw SAST precision was 37.7%", 16, ROSE, True, False),
     ("   roughly 2 in 3 raw findings were false alarms before triage", 14, WHITE, False, False)],
], align=CENTER)
_notes(s, """
Source: benchmarks/results/matrix_20260604_151302/deepseek/REPORT.md (Summary
Comparison). raw-sast vs vulnhunterx on 300 OWASP-Python findings: precision
37.7%→87.3%, F1 54.7%→92.4%, 91.4% of false positives removed while recall stays
at 98.2%. Headline: far fewer false alarms, almost no real bugs lost.
""")

# 16 — Result: the multi-turn loop pays off (NEW)
s = content_slide("Result — the multi-turn loop pays off")
textbox(s, Inches(0.6), Inches(1.25), Inches(12), Inches(0.45),
        [[("High-confidence accuracy, by how many turns the LLM took (OWASP-Python · DeepSeek).",
           15, MUTE, False, True)]])
c1 = soft_shadow(shape(s, RR, Inches(1.1), Inches(2.0), Inches(4.3), Inches(2.3), fill=CARD, line=BORDER))
shape(s, RR, Inches(1.1), Inches(2.0), Inches(4.3), Inches(0.14), fill=ROSE)
text_in(c1, [[("81.8%", 50, ROSE, True, False)],
             [("stop after 1 turn", 15, NAVY_TX, True, False)],
             [("commit too early, miss the defense", 11.5, GREY, False, False)]], align=CENTER)
block_arrow(s, Inches(5.7), Inches(2.95), Inches(1.9), Inches(0.5), ORANGE)
c2 = soft_shadow(shape(s, RR, Inches(7.9), Inches(2.0), Inches(4.3), Inches(2.3), fill=CARD, line=BORDER))
shape(s, RR, Inches(7.9), Inches(2.0), Inches(4.3), Inches(0.14), fill=EMERALD)
text_in(c2, [[("96.9%", 50, EMERALD, True, False)],
             [("after a 2nd turn", 15, NAVY_TX, True, False)],
             [("expand context, then decide", 11.5, GREY, False, False)]], align=CENTER)
strip = soft_shadow(shape(s, RR, Inches(0.6), Inches(4.7), Inches(12.15), Inches(1.55), fill=DARK))
text_in(strip, [
    [("High-confidence accuracy 94.2%", 16, AMBER, True, False),
     ("   well-calibrated over 275 high-confidence verdicts", 14, WHITE, False, False)],
    [("Mean 2.6 turns", 16, CYAN, True, False),
     ("   the loop expands context only when it needs to (max 5)", 14, WHITE, False, False)],
    [("$0.40 for 300 findings", 16, EMERALD, True, False),
     ("   on DeepSeek; $0.00 with a local Ollama model", 14, WHITE, False, False)],
], align=CENTER)
_notes(s, """
Source: REPORT.md (Iteration×Confidence + Cost tables). Cleanest signal is the
early-terminate bucket: 1-turn/High 81.8% vs 2-turn/High 96.9% — a second pass to
expand context flips many borderline verdicts. High-confidence calibration 94.2%;
mean 2.6 turns; $0.40 / 300 findings on DeepSeek (free on local Ollama).
""")

# 17 — Result: per-CWE + cross-model (NEW)
s = card_grid_slide("Result — per-CWE and across models", [
    {"head": "CWE-89 · SQLi", "body": "F1 100%", "accent": EMERALD},
    {"head": "CWE-79 · XSS", "body": "F1 100%", "accent": EMERALD},
    {"head": "CWE-330 · weak RNG", "body": "F1 100%", "accent": EMERALD},
    {"head": "CWE-643 · XPath", "body": "F1 84.4%", "accent": AMBER},
    {"head": "CWE-611 · XXE", "body": "F1 50% — hard case", "accent": ROSE},
], cols=5, top=1.7, card_h=1.5, gap=0.2, head_size=12.5, body_size=12)
strip = soft_shadow(shape(s, RR, Inches(0.6), Inches(3.5), Inches(12.15), Inches(2.95), fill=DARK))
text_in(strip, [
    [("Same approach, three models", 16, AMBER, True, False),
     ("    OWASP-Java · VulnHunterX", 13.5, MUTE, False, True)],
    [("", 6, WHITE, False, False)],
    [("DeepSeek", 15, CYAN, True, False),
     ("    F1 92.0%  ·  FP-reduction 78.6%  ·  $0.30", 14, WHITE, False, False)],
    [("GPT-4.1-mini", 15, CYAN, True, False),
     ("    F1 85.7%  ·  FP-reduction 58.1%  ·  $0.80", 14, WHITE, False, False)],
    [("Ollama Qwen (local)", 15, CYAN, True, False),
     ("    FP-reduction 100%  ·  $0.00 — runs entirely offline", 14, WHITE, False, False)],
    [("", 6, WHITE, False, False)],
    [("Focused slice: DeepSeek, June 2026 runs — directional, not a universal claim.",
      12, MUTE, False, True)],
], align=CENTER)
_notes(s, """
Sources: REPORT.md (Per-CWE Breakdown) and matrix_20260604_224953/COMPARISON.md
(OWASP-Java model comparison). Easy injection classes hit F1 100%; XXE (CWE-611)
is the honest hard case at 50%. Across models the method holds: stronger models
give higher precision, and a local Ollama model still removes false positives at
$0 — useful when code cannot leave the building.
""")

# 18 — SECTION
section_slide("Part 2", "How to install & use")

# 19 — Tool cards
card_grid_slide("The four tools — what each does", [
    {"head": "CodeQL", "body": "Semantic, database-driven taint analysis — the deepest engine. Needs a DB build.", "accent": INDIGO, "pill": "core", "pill_color": INDIGO},
    {"head": "Semgrep", "body": "Fast pattern-based rules across many languages. No DB needed.", "accent": CYAN, "pill": "breadth", "pill_color": CYAN},
    {"head": "OpenGrep", "body": "Semgrep-compatible fork — same rules, drop-in alternative.", "accent": AMBER, "pill": "alt", "pill_color": AMBER},
    {"head": "tree-sitter", "body": "Syntactic context extraction — pip-installed fallback when there's no CodeQL DB.", "accent": EMERALD, "pill": "fallback", "pill_color": EMERALD},
], cols=2, top=1.6, card_h=1.95, gap=0.3, head_size=20, body_size=13.5,
   notes="""
CodeQL is the backbone. Semgrep/OpenGrep add breadth and run without a DB.
tree-sitter is the safety net — pure Python, ships with the pip install.
`check-env` verifies all of them at once.
""")

# 20 — Install
code_slide("Install VulnHunterX", [
    "# Prerequisites: Python 3.12+, CodeQL CLI 2.15+, an LLM provider key",
    "git clone https://github.com/vinsoc-cyber/VulnHunterX.git && cd VulnHunterX",
    "",
    "uv venv --python python3.12 .venv && source .venv/bin/activate",
    "uv pip install -e \".[dev]\"        # plain  python3.12 -m venv + pip  also works",
    "",
    "cp env.example .env                # add OPENAI / ANTHROPIC / OLLAMA keys",
    "vuln-hunter-x check-env            # verify the whole toolchain is green",
], notes="""
Exact commands from README "Install". uv is recommended but optional. Run
check-env first — it reports which engines and provider keys are available.
""")

# 21 — Config
code_slide("Configuration you need to set", [
    "# .env  — provider keys + optional tool paths",
    "LLM_PROVIDER=openai            # openai | anthropic | ollama",
    "LLM_MODEL=gpt-4o",
    "OPENAI_API_KEY=sk-...          # or ANTHROPIC_API_KEY / OLLAMA_API_BASE",
    "CODEQL_PATH=/path/to/codeql    # only if codeql isn't on your PATH",
    "SEMGREP_PATH=/path/to/semgrep  # only if semgrep isn't on your PATH",
    "",
    "# config/confirm_findings.yaml  — verification tuning",
    "provider: openai   model: gpt-4o   temperature: 0.2   max_iterations: 3",
    "",
    "# config/repos.yaml  — batch targets (name, url, language, build_command)",
], notes="""
Priority: CLI args > env vars > config file > defaults. Local Ollama needs no
key. C/C++ need a build_command; C# uses buildless extraction so it doesn't.
""")

# 22 — three ways
code_slide("Three ways to start a scan", [
    "# 1) Friendliest: a guided wizard (prompts for everything, live-tests your LLM)",
    "vuln-hunter-x interactive",
    "",
    "# 2) One-shot: the whole pipeline in a single command",
    "vuln-hunter-x scan --url https://github.com/org/app.git \\",
    "                   --lang python --profile extended --limit 5",
    "",
    "# 3) Stage by stage (full control)",
    "vuln-hunter-x prepare --repo pyyaml",
    "vuln-hunter-x analyze --repo pyyaml --profile extended",
    "vuln-hunter-x verify  --repo pyyaml --limit 5",
    "vuln-hunter-x report  --repo pyyaml --lang python",
], notes="""
Start the demo with `interactive` — it validates inputs and live-tests the
provider so you don't fail halfway. --limit keeps the LLM cost small live.
""")

# 23 — SECTION
section_slide("Part 3", "Real example — dvpwa")

# 24 — dvpwa intro + contrast (NEW)
s = content_slide("dvpwa — a deliberately vulnerable web app")
textbox(s, Inches(0.6), Inches(1.25), Inches(12), Inches(0.5),
        [[("An intentionally insecure Python ", 16, GREY, False, False),
          ("aiohttp", 16, INK, True, False),
          (" app (github.com/anxolerd/dvpwa) — scanned next to a clean library.",
           16, GREY, False, False)]])
c1 = soft_shadow(shape(s, RR, Inches(0.7), Inches(2.1), Inches(5.9), Inches(2.0), fill=CARD, line=BORDER))
shape(s, RR, Inches(0.7), Inches(2.1), Inches(0.16), Inches(2.0), fill=ROSE)
text_in(c1, [[("dvpwa", 22, NAVY_TX, True, False), ("    vulnerable", 13, ROSE, True, True)],
             [("SQL injection · stored XSS · weak MD5 ·", 13.5, GREY, False, False)],
             [("missing authorization · CSRF disabled", 13.5, GREY, False, False)],
             [("→ VulnHunterX should surface real TPs", 13.5, ROSE, True, False)]],
        align=LEFT, ml=0.3, anchor=MID)
c2 = soft_shadow(shape(s, RR, Inches(6.75), Inches(2.1), Inches(5.9), Inches(2.0), fill=CARD, line=BORDER))
shape(s, RR, Inches(6.75), Inches(2.1), Inches(0.16), Inches(2.0), fill=EMERALD)
text_in(c2, [[("pyyaml", 22, NAVY_TX, True, False), ("    real-world", 13, EMERALD, True, True)],
             [("a mature, widely-used YAML library", 13.5, GREY, False, False)],
             [("with no planted vulnerabilities", 13.5, GREY, False, False)],
             [("→ false positives should be suppressed", 13.5, EMERALD, True, False)]],
        align=LEFT, ml=0.3, anchor=MID)
foot = soft_shadow(shape(s, RR, Inches(0.7), Inches(4.45), Inches(11.95), Inches(1.6), fill=DARK))
text_in(foot, [[("The value proposition in one screen:", 17, AMBER, True, False)],
               [("the vulnerable app lights up with real bugs; the clean library stays quiet.",
                 15, WHITE, False, False)],
               [("Ground truth (realvuln): 22 real vulnerabilities + 4 false-positive traps.",
                 13, MUTE, False, True)]], align=CENTER)
_notes(s, """
dvpwa = Damn Vulnerable Python Web App (aiohttp, anxolerd/dvpwa). We run it beside
pyyaml — a clean real-world library — so the room sees both sides: TPs surfaced on
dvpwa, FPs suppressed on pyyaml. The realvuln ground truth labels 22 genuine
vulnerabilities and 4 deliberate false-positive traps; next slides give the
commands and the answer key.
""")

# 25 — Run the dvpwa example (NEW)
code_slide("Run the dvpwa example", [
    "# one command: clone + analyze + verify + report, for both repos",
    "python examples/pipeline_python.py            # add --dry-run to preview",
    "",
    "# or drive just dvpwa, stage by stage:",
    "vuln-hunter-x prepare --repo dvpwa            # clone + CodeQL DB + context",
    "vuln-hunter-x analyze --repo dvpwa --profile extended",
    "vuln-hunter-x verify  --repo dvpwa --limit 5 --report",
    "",
    "# read the verdicts the LLM produced:",
    "cat output/python/dvpwa/verification_results/report.md",
], notes="""
examples/pipeline_python.py runs pyyaml + dvpwa end to end (REPOS list in the
script). --dry-run prints the commands without calling the LLM. report.md holds
per-finding verdicts (TP / FP / NMD) with evidence-anchored reasoning. Keep
--limit small so a live demo stays cheap.
""")

# 26 — dvpwa baseline / answer key (NEW)
s = card_grid_slide("dvpwa answer key — verify your results", [
    {"head": "SQL injection · CWE-89", "body": "dao/student.py:42 — query built with Python % then execute(q)", "accent": EMERALD, "pill": "TP", "pill_color": EMERALD},
    {"head": "Stored XSS ×5 · CWE-79", "body": "templates render unescaped — root cause autoescape=False (app.py:35)", "accent": EMERALD, "pill": "TP", "pill_color": EMERALD},
    {"head": "Weak password hash · CWE-916", "body": "dao/user.py:41 — unsalted MD5", "accent": EMERALD, "pill": "TP", "pill_color": EMERALD},
    {"head": "Missing authz / CSRF off", "body": "views.py (CWE-862) · csrf middleware commented out, app.py:27 (CWE-352)", "accent": EMERALD, "pill": "TP", "pill_color": EMERALD},
    {"head": "NOT SQL injection", "body": "dao/user.py · review.py · course.py — parameterized %s / %(name)s queries", "accent": ROSE, "pill": "trap", "pill_color": ROSE},
    {"head": "NOT XSS", "body": "student.jinja2:14 — {{ name | e }} is explicitly escaped", "accent": ROSE, "pill": "trap", "pill_color": ROSE},
], cols=2, top=1.5, card_h=1.1, gap=0.2, head_size=15.5, body_size=12.5)
strip = soft_shadow(shape(s, RR, Inches(0.6), Inches(5.35), Inches(12.15), Inches(1.45), fill=DARK))
text_in(strip, [
    [("Ground truth: 22 real vulnerabilities + 4 false-positive traps", 15, AMBER, True, False)],
    [("Score yourself: flagging student.py SQLi = correct; flagging the parameterized DAO "
      "queries or the escaped template = a false positive your run should suppress.",
      13, WHITE, False, False)],
], align=CENTER)
_notes(s, """
Answer key: benchmarks/datasets/realvuln/ground-truth/realvuln-dvpwa/ground-truth.json
(22 is_vulnerable=true entries + 4 is_vulnerable=false traps). The traps are the
teaching moment: three DAO create() methods use psycopg2 placeholders (safe) and
student.jinja2:14 uses the | e filter — a good run must NOT report these. Students
compare their report.md verdicts against this list.
""")

# 27 — SECTION
section_slide("Part 4", "Homework — dvcp")

# 28 — dvcp homework + run (NEW)
code_slide("Homework — scan dvcp (Damn Vulnerable C Program)", [
    "# config/repos.yaml already defines dvcp (C needs a build command):",
    "#   url: github.com/hardik05/Damn_Vulnerable_C_Program",
    "#   build_command: gcc -g -o dvcp imgRead.c",
    "",
    "# run the C pipeline — contrast target is c-ares (a real DNS library):",
    "python examples/pipeline_c.py                 # add --dry-run to preview",
    "",
    "# or drive just dvcp, stage by stage:",
    "vuln-hunter-x prepare --repo dvcp             # clone + build + CodeQL DB",
    "vuln-hunter-x analyze --repo dvcp",
    "vuln-hunter-x verify  --repo dvcp --limit 5 --report",
    "",
    "# stretch (Linux/macOS): fuzz the real library",
    "python examples/pipeline_c.py --fuzz",
], caption="Your task: run it, open the report, list the TPs and FPs — then check against the answer key.",
   notes="""
dvcp = Damn Vulnerable C Program (hardik05), a single imgRead.c. C requires a
build_command for CodeQL DB creation — already set in config/repos.yaml.
examples/pipeline_c.py runs c-ares (real-world) + dvcp. The --fuzz stretch runs a
sanitizer build + libFuzzer on c-ares (Linux/macOS only). Bring your TP/FP list to
compare against the next slide.
""")

# 29 — dvcp baseline / answer key (NEW)
s = card_grid_slide("dvcp answer key — what to expect", [
    {"head": "Integer overflow · CWE-190", "body": "width + height, width * height", "accent": EMERALD, "pill": "TP", "pill_color": EMERALD},
    {"head": "Integer underflow · CWE-191", "body": "width - height + 100", "accent": EMERALD, "pill": "TP", "pill_color": EMERALD},
    {"head": "Heap buffer overflow · CWE-122", "body": "memcpy into a buffer sized from the overflowed value", "accent": EMERALD, "pill": "TP", "pill_color": EMERALD},
    {"head": "Double-free · CWE-415", "body": "free(buff1) again when size1 % 2 == 0", "accent": EMERALD, "pill": "TP", "pill_color": EMERALD},
    {"head": "Use-after-free · CWE-416", "body": "buff1[0]='a' after free when size1 % 3 == 0", "accent": EMERALD, "pill": "TP", "pill_color": EMERALD},
    {"head": "OOB read / write · CWE-125/787", "body": "buff3[size3], buff4[size3] past the bounds", "accent": EMERALD, "pill": "TP", "pill_color": EMERALD},
    {"head": "Divide-by-zero · CWE-369", "body": "width / height when height == 0", "accent": EMERALD, "pill": "TP", "pill_color": EMERALD},
    {"head": "Leak + infinite recursion", "body": "CWE-401: buff4=0 without free · CWE-674: stack_operation() recurses forever", "accent": EMERALD, "pill": "TP", "pill_color": EMERALD},
], cols=2, top=1.4, card_h=0.9, gap=0.16, head_size=14, body_size=11.5)
strip = soft_shadow(shape(s, RR, Inches(0.6), Inches(5.5), Inches(12.15), Inches(1.3), fill=DARK))
text_in(strip, [
    [("All bugs live in ProcessImage()", 15, AMBER, True, False),
     ("  — reachable from attacker-controlled image-header fields (width / height / data).", 13.5, WHITE, False, False)],
    [("Contrast: c-ares (production DNS library) should come back mostly false positives.", 13.5, CYAN, True, False)],
    [("Answer key derived from the upstream imgRead.c source (external repo).", 12, MUTE, False, True)],
], align=CENTER)
_notes(s, """
dvcp's imgRead.c packs ~15 memory-safety bugs into one ProcessImage() function, all
reachable from attacker-controlled image-header fields. This is the C counterpart to
the dvpwa key: students list what their run flagged on dvcp (should be these
memory-safety classes) and confirm c-ares stays mostly quiet. This key is from the
upstream source, not the realvuln (Python-only) ground truth.
""")

# 30 — Recap
s = content_slide("Recap & references")
points = [
    ("What it is", "SAST finds candidates; a multi-turn LLM verifies them with evidence", INDIGO),
    ("Results", "OWASP-Python: precision 37.7% → 87.3%, 91.4% of false positives removed", EMERALD),
    ("Install & use", "prepare → analyze → verify → report  ·  vuln-hunter-x interactive", CYAN),
    ("Practice", "dvpwa example today · dvcp homework — check both against the answer keys", ORANGE),
]
y = Inches(1.55)
for h, b, col in points:
    cd = soft_shadow(shape(s, RR, Inches(0.7), y, Inches(11.9), Inches(0.82),
                           fill=CARD, line=BORDER))
    shape(s, RR, Inches(0.7), y, Inches(0.13), Inches(0.82), fill=col)
    text_in(cd, [[(h + ":   ", 16, col, True, False), (b, 14, GREY, False, False)]],
            align=LEFT, ml=0.3)
    y = int(y + Inches(0.95))
start = soft_shadow(shape(s, RR, Inches(0.7), y, Inches(11.9), Inches(0.82), fill=DARK))
text_in(start, [[("Start here →  ", 17, AMBER, True, False),
                 ("vuln-hunter-x interactive", 17, EMERALD, True, False, )]], align=CENTER)
textbox(s, Inches(0.7), int(y + Inches(0.95)), Inches(12), Inches(0.5),
        [[("References:  Vulnhalla (CyberArk) · CodeQL docs · Semgrep docs · SARIF spec · repo README + config/RULES.md",
           12.5, MUTE, False, True)]], align=CENTER)
_notes(s, """
Close on the one command to remember: vuln-hunter-x interactive. dvpwa was the
worked example; dvcp is the take-home — both have answer keys to self-check.
References are in the repo README and WORKSHOP.md.
""")


# --------------------------------------------------------------------------- #
prs.save(OUT_PATH)
print(f"Wrote {OUT_PATH}  ({len(prs.slides)} slides)")
print(f"Counts -> CodeQL custom: {CODEQL_TOTAL} {CODEQL}")
print(f"          Semgrep custom: {SEMGREP_TOTAL} {SEMGREP}")
print(f"          Guided Qs: ~{QUESTIONS_TOTAL} across {QUESTION_BANKS} banks {QUESTIONS}")
