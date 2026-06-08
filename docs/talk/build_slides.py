#!/usr/bin/env python3
"""Generate the combined DefCon talk deck (presenter-style).

This single script builds ONE deck that fuses:
  * the full DefCon talk (problem → method → benchmark evidence → live demo →
    limits → wrap-up), source of truth: docs/talk/defcon-outline.md
  * the "VulnHunterX in the wild" real-bug results segment (5 manually-confirmed,
    vendor-engaged vulnerabilities), data source:
    https://github.com/tuonglnc/VHX-real-bug-confirmed

The real-world-results segment is woven in right after the live demo, before the
limits/red-team wrap-up.

Output: docs/talk/defcon-vulnhunterx.pptx

Usage:
    python docs/talk/build_slides.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# --- palette (dark "terminal" security-conf look) -------------------------
BG = RGBColor(0x0B, 0x0F, 0x14)      # near-black
PANEL = RGBColor(0x16, 0x1B, 0x22)   # slightly lighter panel
FG = RGBColor(0xE6, 0xED, 0xF3)      # off-white text
MUTED = RGBColor(0x8B, 0x94, 0x9E)   # gray
ACCENT = RGBColor(0xFF, 0x6B, 0x35)  # orange (SAST)
ACCENT2 = RGBColor(0x2D, 0xD4, 0xBF)  # teal (LLM)
GOOD = RGBColor(0x3F, 0xB9, 0x50)    # green
BAD = RGBColor(0xF8, 0x51, 0x49)     # red

BODY_FONT = "Calibri"
MONO_FONT = "Consolas"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
TALK_REPO = "github.com/vinsoc-cyber/VulnHunterX"
BUGS_REPO = "github.com/tuonglnc/VHX-real-bug-confirmed"
REPO = TALK_REPO  # default footer repo


def _set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _no_line(shape):
    shape.line.fill.background()


def _box(slide, l, t, w, h, fill=None):
    sp = slide.shapes.add_shape(1, l, t, w, h)  # 1 = rectangle
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    _no_line(sp)
    sp.shadow.inherit = False
    return sp


def _text(slide, l, t, w, h, runs, *, size=18, color=FG, bold=False,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=BODY_FONT,
          line_spacing=1.1, space_after=6):
    """runs: str OR list of paragraphs; a paragraph is str or list of (text, kw)."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        segs = para if isinstance(para, list) else [(para, {})]
        for seg_text, kw in segs:
            r = p.add_run()
            r.text = seg_text
            f = r.font
            f.size = Pt(kw.get("size", size))
            f.bold = kw.get("bold", bold)
            f.name = kw.get("font", font)
            f.color.rgb = kw.get("color", color)
    return tb


_SLIDE_NO = 1  # title slide is #1; each _chrome() call below is the next slide


def _chrome(slide, kicker=None, *, mins=None, tag=None, repo=None,
            title_slide=False):
    """Slide chrome: accent bar, footer (repo + N/TOTAL), optional kicker/tag.

    Slide numbers auto-increment in call order so slides can be inserted
    mid-deck without renumbering every call.
      kicker  upper-cased section label, top-left (talk + real-bug segments)
      mins    appended to the kicker as a timing hint (talk style)
      tag     right-aligned teal status chip (real-bug deep-dive style)
      repo    footer repo override (real-bug slides use BUGS_REPO)
    """
    global _SLIDE_NO
    _SLIDE_NO += 1
    # left accent bar
    _box(slide, 0, 0, Inches(0.16), EMU_H, ACCENT)
    if title_slide:
        return
    # footer
    _text(slide, Inches(0.5), Inches(7.02), Inches(9), Inches(0.4),
          repo or TALK_REPO, size=11, color=MUTED)
    _text(slide, Inches(11.3), Inches(7.02), Inches(1.6), Inches(0.4),
          f"{_SLIDE_NO} / {TOTAL}", size=11, color=MUTED, align=PP_ALIGN.RIGHT)
    if kicker:
        # narrow the kicker box when a tag shares the top row, to avoid overlap
        kw = Inches(9) if tag else Inches(11)
        label = kicker + (f"   ·   {mins}" if mins else "")
        _text(slide, Inches(0.55), Inches(0.42), kw, Inches(0.4),
              label.upper(), size=12.5, color=ACCENT, bold=True)
    if tag:  # right-aligned status chip text
        _text(slide, Inches(8.2), Inches(0.42), Inches(4.6), Inches(0.4),
              tag, size=12.5, color=ACCENT2, bold=True, align=PP_ALIGN.RIGHT)


def _title(slide, title, top=Inches(0.95), size=33):
    _text(slide, Inches(0.55), top, Inches(12.2), Inches(1.1),
          title, size=size, color=FG, bold=True, line_spacing=1.0)


def _bullets(slide, items, *, top=2.25, size=20, gap=14, left=0.85, width=11.6):
    paras = []
    for it in items:
        txt, lvl = (it, 0) if isinstance(it, str) else it
        bullet = "—  " if lvl == 0 else "      ·  "
        col = FG if lvl == 0 else MUTED
        sz = size if lvl == 0 else size - 3
        paras.append([(bullet, {"color": ACCENT2, "size": sz}),
                      (txt, {"color": col, "size": sz})])
    _text(slide, Inches(left), Inches(top), Inches(width), Inches(4.4),
          paras, space_after=gap, line_spacing=1.08)


def _code(slide, lines, *, top=2.3, left=0.85, width=11.6, height=3.6, size=15):
    panel = _box(slide, Inches(left), Inches(top), Inches(width), Inches(height), PANEL)
    panel.line.color.rgb = RGBColor(0x30, 0x36, 0x3D)
    panel.line.width = Pt(0.75)
    paras = []
    for ln, color in lines:
        paras.append([(ln if ln else " ", {"color": color, "font": MONO_FONT, "size": size})])
    _text(slide, Inches(left + 0.25), Inches(top + 0.18), Inches(width - 0.5),
          Inches(height - 0.36), paras, font=MONO_FONT, size=size,
          line_spacing=1.12, space_after=2)


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def _table(slide, headers, rows, *, top=2.3, left=0.7, width=11.9,
           col_w=None, highlight_col=None, row_colors=None, fs=14):
    nrows, ncols = len(rows) + 1, len(headers)
    height = Inches(0.5 + 0.46 * len(rows))
    gtbl = slide.shapes.add_table(nrows, ncols, Inches(left), Inches(top),
                                  Inches(width), height)
    tbl = gtbl.table
    # strip default banding style
    tblPr = tbl._tbl.tblPr
    tblPr.set("firstRow", "0")
    tblPr.set("bandRow", "0")
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Emu(int(Inches(width) * cw / total))
    # header
    for c, htext in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        cell.margin_top = cell.margin_bottom = Pt(3)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT
        r = p.add_run(); r.text = htext
        r.font.size = Pt(fs); r.font.bold = True
        r.font.color.rgb = RGBColor(0x10, 0x12, 0x16); r.font.name = BODY_FONT
    # body
    for ridx, row in enumerate(rows, start=1):
        base = PANEL if ridx % 2 else RGBColor(0x0F, 0x14, 0x1A)
        if row_colors and row_colors[ridx - 1]:
            base = row_colors[ridx - 1]
        for c, val in enumerate(row):
            cell = tbl.cell(ridx, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = base
            cell.margin_top = cell.margin_bottom = Pt(2)
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = val
            r.font.size = Pt(fs); r.font.name = BODY_FONT
            strong = (c == 0) or (highlight_col and c in highlight_col)
            r.font.bold = bool(strong)
            r.font.color.rgb = ACCENT2 if (highlight_col and c in highlight_col) else FG
    return gtbl


def _impact(slide, label, text, *, top, color=BAD):
    bx = _box(slide, Inches(0.85), Inches(top), Inches(11.7), Inches(0.95), PANEL)
    bx.line.color.rgb = color; bx.line.width = Pt(1.0)
    _text(slide, Inches(1.1), Inches(top), Inches(11.2), Inches(0.95),
          [[(label + "  ", {"color": color, "bold": True, "size": 15}), (text, {"size": 15})]],
          anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)


def _verdict_chip(slide, text, *, top=6.45):
    _text(slide, Inches(0.85), Inches(top), Inches(11.7), Inches(0.5),
          [[("LLM verdict  ", {"color": GOOD, "bold": True, "size": 13}),
            (text, {"size": 13, "color": MUTED})]])


IMG_DIR = Path(__file__).resolve().parent / "img"


def _rasterize_pdf(pdf_rel, png_name, dpi=200):
    """Render page 1 of a repo PDF to docs/talk/img/<png_name>; return path or None."""
    out = IMG_DIR / png_name
    if out.exists():
        return out
    pdf = Path(__file__).resolve().parents[2] / pdf_rel
    try:
        import fitz  # PyMuPDF
    except ImportError:
        print(f"[warn] pymupdf not installed — skipping {png_name} "
              "(install with: uv pip install pymupdf)")
        return None
    if not pdf.exists():
        print(f"[warn] {pdf} not found — skipping {png_name}")
        return None
    IMG_DIR.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(str(pdf))
    pix = doc[0].get_pixmap(matrix=fitz.Matrix(dpi / 72.0, dpi / 72.0), alpha=False)
    pix.save(str(out))
    doc.close()
    return out


def _picture_or_placeholder(slide, img_path, left, top, width):
    """Add a picture scaled to `width` (aspect preserved); placeholder if missing."""
    if img_path and Path(img_path).exists():
        return slide.shapes.add_picture(str(img_path), left, top, width=width)
    ph = _box(slide, left, top, width, Inches(2.0), PANEL)
    ph.line.color.rgb = ACCENT
    _text(slide, left, top + Inches(0.8), width, Inches(0.5),
          "[ pipeline diagram — rebuild with pymupdf installed ]",
          size=16, color=MUTED, align=PP_ALIGN.CENTER)
    return ph


# ===========================================================================
prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]
TOTAL = 31


def new():
    s = prs.slides.add_slide(BLANK)
    _set_bg(s)
    return s


# --- 1 · Title / hook ------------------------------------------------------
s = new()
_box(s, 0, 0, Inches(0.16), EMU_H, ACCENT)
_text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.5),
      [[("DEF CON  ·  ", {"color": ACCENT, "bold": True, "size": 16}),
        ("Main Track / Demo Labs", {"color": MUTED, "size": 16})]])
_text(s, Inches(0.9), Inches(2.05), Inches(11.6), Inches(2.0),
      [[("Picking True Bugs from the ", {"bold": True, "size": 44}),
        ("CodeQL Haystack", {"bold": True, "size": 44, "color": ACCENT})]],
      line_spacing=1.0)
_text(s, Inches(0.9), Inches(3.55), Inches(11.6), Inches(0.8),
      "Teaching small LLMs to triage like a senior analyst", size=24, color=ACCENT2)
_box(s, Inches(0.9), Inches(4.6), Inches(7.4), Inches(1.15), PANEL)
_text(s, Inches(1.15), Inches(4.72), Inches(7.0), Inches(1.0),
      [[("30–80%", {"color": ACCENT, "bold": True, "size": 30}),
        ("  of SAST findings are false positives.", {"size": 20})],
       [("What if a $0 model could do the first pass — well?", {"color": MUTED, "size": 15})]],
      line_spacing=1.05, space_after=4)
_text(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5),
      [[("VinSOC Cyber", {"size": 15}), ("     " + TALK_REPO, {"color": MUTED, "size": 15})]])
_notes(s, """
Open on the cost of triage, not detection. Security teams don't drown in MISSED bugs — they
drown in triaging the ones the scanner flagged, most of which are safe. FP rates run 30–80%.
The question for this talk: what if a model that costs zero dollars could do the first triage pass,
and do it well? This builds on CyberArk's Vulnhalla methodology — name it up front for credibility.
State the thesis now and again at the end: reasoning STRUCTURE beats model SIZE.
""")

# --- 2 · Why this is a DefCon problem -------------------------------------
s = new(); _chrome(s, "The problem", mins="1.5 min")
_title(s, "FP fatigue is a security failure mode")
_bullets(s, [
    "Thousands of SAST findings  →  analyst hours  →  a handful of real bugs",
    "When ~70% of alerts are noise, analysts stop reading the queue",
    ("…and the real bug rides in with the noise", 1),
    "Dual-use: defender triage tooling AND an offense recon accelerator",
], top=2.3, gap=18)
_notes(s, """
Draw the funnel. The failure isn't a missed detection — it's that humans can't keep up with the
volume, so they tune out, and a genuine vulnerability gets ignored alongside the noise. Frame the
dual-use angle for a DefCon room: the same triage speed-up helps defenders clear a backlog and
helps offense prioritize attacker-relevant findings fast.
""")

# --- 3 · Anatomy of a false positive --------------------------------------
s = new(); _chrome(s, "The problem", mins="2 min")
_title(s, "Anatomy of a false positive")
_code(s, [
    ("# SAST flags this as path-traversal (CWE-22)", MUTED),
    ("def download(req):", FG),
    ("    name = req.args['file']", FG),
    ("    safe = secure_filename(name)        # <-- guard upstream", GOOD),
    ("    path = os.path.join(BASE, safe)", FG),
    ("    return open(path).read()            # <-- flagged sink", ACCENT),
], top=2.3, height=2.7)
_text(s, Inches(0.85), Inches(5.3), Inches(11.6), Inches(1.0),
      [[("Vote: ", {"bold": True, "size": 22}),
        ("True Positive", {"color": BAD, "bold": True, "size": 22}),
        ("  or  ", {"size": 22}),
        ("False Positive", {"color": GOOD, "bold": True, "size": 22}),
        (" ?", {"size": 22})]])
_notes(s, """
Walk the snippet. The sink genuinely looks dangerous, but secure_filename() upstream neutralizes
the traversal — so it's a false positive. Get a show of hands. This is exactly the judgment we're
automating, and this is the case the LIVE DEMO will resolve later. Plant the hook now.
""")

# --- 4 · Why naive LLM fails ----------------------------------------------
s = new(); _chrome(s, "The problem", mins="2 min")
_title(s, 'Why "just ask GPT if it\'s a bug" fails')
_bullets(s, [
    "Free-form prompting pattern-matches the shape, not the data flow",
    ("free(p) … *p  →  \"use-after-free!\"  — without checking paths overlap or a guard exists", 1),
    "SecLLMHolmes (IEEE S&P 2024): frontier models cap ~40% on hand-crafted cases",
    "The gap isn't model IQ — it's METHOD",
], top=2.3, gap=18)
_notes(s, """
The naive approach is confidently wrong because it matches a textual pattern instead of reasoning
about reachability and guards. SecLLMHolmes measured frontier models capping around 40% accuracy
on hand-crafted scenarios for precisely this reason. Land the line: the gap is method, not
intelligence — which is what the rest of the talk fixes.
""")

# --- 5 · Pipeline ----------------------------------------------------------
s = new(); _chrome(s, "The method", mins="2 min")
_title(s, "The pipeline — SARIF is the spine")
stages = [("1 prepare", ACCENT2), ("2 analyze", ACCENT2), ("3 verify", ACCENT),
          ("4 report", ACCENT2), ("5–8 fuzz", GOOD)]
x = 0.85; w = 2.25; gap = 0.18
for i, (label, col) in enumerate(stages):
    bx = _box(s, Inches(x), Inches(2.7), Inches(w), Inches(1.2), PANEL)
    bx.line.color.rgb = col; bx.line.width = Pt(1.5)
    _text(s, Inches(x), Inches(2.95), Inches(w), Inches(0.7), label,
          size=18, bold=True, color=col, align=PP_ALIGN.CENTER)
    x += w + gap
_text(s, Inches(0.85), Inches(4.35), Inches(11.6), Inches(1.6),
      [[("Stages 1–2", {"bold": True, "color": ACCENT2}), ("  ordinary SAST (CodeQL / Semgrep / OpenGrep)", {})],
       [("Stage 3", {"bold": True, "color": ACCENT}), ("  the LLM verification this talk is about", {})],
       [("Stages 5–8", {"bold": True, "color": GOOD}), ("  optional: PROVE the bug with a crash", {})]],
      size=19, space_after=12)
_notes(s, """
One diagram for the whole system. SARIF is the contract that ties it together — remember that for
the extensibility slide. Stages 1–2 are ordinary static analysis; stage 3 is the LLM verification
that is the heart of the talk; stages 5–8 optionally confirm a true positive with an actual crash.
Everything routes through one SARIF interface.
""")

# --- 6 · Pipeline diagram (from the paper) --------------------------------
s = new(); _chrome(s, "The method", mins="1.5 min")
_title(s, "Stages 1–4 · detection → verified triage")
_img = _rasterize_pdf("docs/paper/diagram-2s-pipeline.pdf", "pipeline-2stage.png")
pic = _picture_or_placeholder(s, _img, Inches(0.7), Inches(3.0), Inches(11.9))
try:  # vertically center the wide banner in the body area
    top_b, bot_b = Inches(2.2), Inches(6.3)
    pic.top = int(top_b + ((bot_b - top_b) - pic.height) / 2)
except Exception:
    pass
_text(s, Inches(0.7), Inches(6.45), Inches(11.9), Inches(0.6),
      "Dotted line = pre-extracted context CSVs feeding multi-turn verification (Stage 1 → Stage 3).",
      size=14, color=MUTED, align=PP_ALIGN.CENTER)
_notes(s, """
The real pipeline figure from the paper. Repo in, triaged findings out. Stage 1 prepares the
CodeQL DB and extracts context CSVs (tree-sitter fallback when the DB build fails); Stage 2 runs
the SAST engines to SARIF; Stage 3 is the multi-turn LLM verification; Stage 4 aggregates a
bilingual report. Point at the dotted line: those pre-extracted context CSVs are what the LLM
draws on in Stage 3 WITHOUT re-running analysis — that's why multi-turn is cheap.
""")

# --- 7 · Guided questions --------------------------------------------------
s = new(); _chrome(s, "The method", mins="2 min")
_title(s, "Guided questions = encoded analyst expertise")
_code(s, [
    ("py/sql-injection:", ACCENT2),
    ("  questions:", FG),
    ('    - "Quote the EXACT sink statement and name the', FG),
    ('       variable passed to it."', FG),
    ('    - "List EVERY assignment to that variable on each', FG),
    ('       path to the sink, with line numbers."', FG),
    ('    - "Does each value derive from user input or a', FG),
    ('       constant/safe source? Cite the chain."', FG),
], top=2.2, height=3.0, size=14)
_text(s, Inches(0.85), Inches(5.4), Inches(11.8), Inches(1.2),
      [[("P1", {"bold": True, "color": ACCENT}), (" evidence-bound (cite lines)   ", {}),
        ("P2", {"bold": True, "color": ACCENT}), (" atomic (one fact)   ", {}),
        ("P3", {"bold": True, "color": ACCENT}), (" refusal allowed → fetch more", {})]],
      size=17)
_notes(s, """
Instead of "is this a bug?", ask the model the same ordered questions a senior reviewer asks.
Three design rules: P1 every question must be answerable only by citing concrete line numbers;
P2 one fact per question, compound questions get split; P3 "not visible in the provided context"
is a legal answer — that's what triggers fetching more context. The repo ships 348 of these banks
across six languages, routed by rule ID and CWE.
""")

# --- 8 · Second CWE example: C/C++ use-after-free -------------------------
s = new(); _chrome(s, "The method", mins="2 min")
_title(s, "Second example — C/C++ use-after-free (CWE-416)")
_code(s, [
    ("cpp/use-after-free:        # lifetime, not taint", ACCENT2),
    ("  questions:", FG),
    ('    - "ANCHOR: quote the flagged line; classify it —', FG),
    ('       pointer USE  /  free·delete  /  declaration?"', FG),
    ('    - "List EVERY free()/delete reaching the use, with', FG),
    ('       line numbers   → request  free_sites:<ptr>"', FG),
    ('    - "Shortest control-flow path free → use: reachable?', FG),
    ('       NULL-set or re-allocated in between?"', FG),
    ("  additional_context: [ free_sites, destructor, field_writes ]", ACCENT2),
    ("  min_iterations: 3", MUTED),
], top=2.05, height=3.55, size=13.5)
_text(s, Inches(0.85), Inches(5.85), Inches(11.9), Inches(0.9),
      [[("Injection asks  ", {"size": 17}),
        ('"where is the taint?"', {"color": ACCENT, "size": 17, "bold": True}),
        ("      use-after-free asks  ", {"size": 17}),
        ('"who owns the lifetime?"', {"color": ACCENT, "size": 17, "bold": True})]])
_notes(s, """
A deliberately different CWE class to show the method generalizes. SQL injection is a taint /
data-flow question — does attacker input reach a sink unsanitized? Use-after-free is a lifetime /
ownership question — these are the real questions from the cpp/use-after-free bank. Note the moves:
anchor and CLASSIFY the flagged line first (many UAF false positives are just declarations); then
enumerate every free site (it requests free_sites: context); then trace the shortest free→use path
and check for a NULL-set or reallocation defense. min_iterations is 3 — memory-safety needs the
turns. The free_sites / destructor / field_writes context it asks for is exactly what the next
slide shows being extracted.
""")

# --- 9 · Answer-before-verdict --------------------------------------------
s = new(); _chrome(s, "The method", mins="2 min")
_title(s, "Answer-before-verdict — the core trick")
_code(s, [
    ("{", FG),
    ('  "answers":   [ "sink at L88; var=path", ... ],   ', ACCENT2),
    ('  "data_flow": "req.args (L10) → join (L15) → open (L88)",', ACCENT2),
    ("  // ^ evidence is generated FIRST", MUTED),
    ('  "verdict":    "False Positive",', ACCENT),
    ('  "confidence": "High",', ACCENT),
    ('  "reasoning":  "secure_filename neutralizes traversal (L14)"', FG),
    ("}", FG),
], top=2.3, height=3.2, size=14.5)
_text(s, Inches(0.85), Inches(5.75), Inches(11.7), Inches(0.9),
      "Verdict token is conditioned on cited evidence. No citation → confidence downgraded.",
      size=17, color=MUTED)
_notes(s, """
This is the single most important slide. The response schema forces the model to write its cited
answers and the data-flow trace FIRST; the verdict field is generated last, conditioned on that
evidence. Because generation is left-to-right, autoregression now works FOR us. A post-processor
downgrades any High/Medium verdict whose reasoning is pure pattern-language with no file:line
citation. Evidence first, conclusion last.
""")

# --- 10 · Context broker ---------------------------------------------------
s = new(); _chrome(s, "The method", mins="2 min")
_title(s, "Context broker + multi-turn")
_bullets(s, [
    "Need more? Model requests from a FIXED vocabulary:",
    ("caller:  ·  struct:  ·  global:  ·  free_sites:  ·  callees:  ·  destructor: …", 1),
    "Resolved from pre-extracted CSVs (tree-sitter fallback) — never re-runs SAST",
    "So multi-turn costs tokens, not another analysis",
], top=2.2, gap=14)
_box(s, Inches(0.85), Inches(5.25), Inches(5.6), Inches(1.4), PANEL)
_text(s, Inches(1.05), Inches(5.45), Inches(5.2), Inches(1.1),
      [[("2.74", {"color": ACCENT, "bold": True, "size": 34}), ("  mean turns", {"size": 17})]],)
_box(s, Inches(6.75), Inches(5.25), Inches(5.6), Inches(1.4), PANEL)
_text(s, Inches(6.95), Inches(5.45), Inches(5.2), Inches(1.1),
      [[("0.8%", {"color": ACCENT, "bold": True, "size": 34}), ("  forced decisions", {"size": 17})]],)
_notes(s, """
When the model says "not visible here," it asks for more context from a fixed vocabulary — caller,
struct, free_sites, and so on. We resolve those from pre-extracted CSVs (or a tree-sitter fallback)
and re-prompt. Crucially we never re-run the SAST engine, so a multi-turn conversation costs a few
thousand extra tokens, not another full analysis. In practice the mean conversation is 2.74 turns
and only 0.8% of findings ever hit a forced decision — source: the 366-entry OWASP-Python run.
""")

# --- 11 · Where context comes from (context_extractor.py) -----------------
s = new(); _chrome(s, "The method", mins="1.5 min")
_title(s, "Where context comes from · context_extractor.py")
flow = [(["CodeQL", "database"], ACCENT2),
        (["codeql query run", "→ bqrs decode"], ACCENT),
        (["per-type", "*.csv"], ACCENT2),
        (["multi-turn", "broker"], GOOD)]
fx, fw, fgap = 0.85, 2.55, 0.5
for i, (lbl, col) in enumerate(flow):
    bx = _box(s, Inches(fx), Inches(2.25), Inches(fw), Inches(1.05), PANEL)
    bx.line.color.rgb = col; bx.line.width = Pt(1.5)
    _text(s, Inches(fx), Inches(2.25), Inches(fw), Inches(1.05), lbl,
          size=15, bold=True, color=col, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
          line_spacing=1.0, space_after=0)
    if i < len(flow) - 1:
        _text(s, Inches(fx + fw - 0.04), Inches(2.4), Inches(fgap + 0.12), Inches(0.7),
              "→", size=24, color=MUTED, align=PP_ALIGN.CENTER)
    fx += fw + fgap
_text(s, Inches(0.85), Inches(3.85), Inches(11.9), Inches(0.5),
      [[("Context types extracted per language  ", {"size": 17, "bold": True}),
        ("(QUERIES_BY_LANG)", {"size": 15, "color": MUTED, "font": MONO_FONT})]])
b1 = _box(s, Inches(0.85), Inches(4.45), Inches(11.7), Inches(1.05), PANEL)
b1.line.color.rgb = ACCENT; b1.line.width = Pt(1.0)
_text(s, Inches(1.1), Inches(4.45), Inches(11.3), Inches(1.05),
      [[("C / C++  (8)", {"color": ACCENT, "bold": True, "size": 16})],
       [("functions · callers · structs · globals · macros · ", {"size": 15}),
        ("free_sites · destructors · field_writes", {"color": ACCENT2, "bold": True, "size": 15})]],
      anchor=MSO_ANCHOR.MIDDLE, space_after=4, line_spacing=1.05)
b2 = _box(s, Inches(0.85), Inches(5.65), Inches(11.7), Inches(0.8), PANEL)
_text(s, Inches(1.1), Inches(5.65), Inches(11.3), Inches(0.8),
      [[("Python / JS / Java / Go / PHP  (3)", {"color": ACCENT, "bold": True, "size": 16})],
       [("functions · callers · classes", {"size": 15})]],
      anchor=MSO_ANCHOR.MIDDLE, space_after=4, line_spacing=1.05)
_text(s, Inches(0.85), Inches(6.6), Inches(11.9), Inches(0.5),
      "The use-after-free questions request free_sites / destructors / field_writes — exactly these CSVs.",
      size=14, color=MUTED)
_notes(s, """
Pull back the curtain on the context broker. context_extractor.py runs small CodeQL queries
against the database — query run, then bqrs decode — and writes one CSV per context type. C and C++
get eight types, including the memory-safety-specific free_sites, destructors, and field_writes
that the use-after-free questions ask for. The other five languages get three (functions, callers,
classes). When no CodeQL DB exists, a tree-sitter fallback emits the same CSV layout, so the rest
of the pipeline is identical. This is the cheap, pre-computed knowledge the LLM pulls from across
multi-turn — no re-analysis.
""")

# --- 12 · Evidence headline ------------------------------------------------
s = new(); _chrome(s, "The evidence", mins="2.5 min")
_title(s, "Precision roughly doubles; FP load collapses")
_table(s,
       ["Dataset", "raw-SAST  P / F1", "VulnHunterX  P / F1", "FP-reduction"],
       [["OWASP-Python (300)", "37.7 / 54.7", "87.3 / 92.4", "91.4%"],
        ["OWASP-Java (full)", "90.0 / 94.7", "97.7 / 96.6", "80.0%"],
        ["Juliet C/C++ (full)", "50.0 / 66.7", "83.8 / 88.5", "82.2%"],
        ["SecLLMHolmes (228)", "52.3 / 68.7", "82.1 / 84.7", "79.4%"]],
       top=2.5, left=0.85, width=11.6, col_w=[3, 2.3, 2.6, 1.8],
       highlight_col=[2, 3], fs=16)
_text(s, Inches(0.85), Inches(5.7), Inches(11.6), Inches(0.9),
      "Every cell is quoted from a published run file (vulnhunterx, best model per row).",
      size=15, color=MUTED)
_notes(s, """
The headline. Precision roughly doubles and we delete 80–91% of false positives while keeping
>90% of real bugs (TP-preservation stays ~90–96%). Point at the repo: every number here is in a
committed benchmark result file, reproducible with the harness. Numbers are percentages, P/F1.
""")

# --- 13 · Small beats big --------------------------------------------------
s = new(); _chrome(s, "The evidence", mins="2.5 min")
_title(s, "Small beats big · cheap beats expensive")
_table(s,
       ["Model", "OWASP-Python F1", "Cost"],
       [["DeepSeek-v4-flash  ($0 pass-through)", "92.4%", "~$0.40 / 300"],
        ["gpt-4.1-mini", "89.4%", "~$1.10 / 300"],
        ["Qwen3-Coder  (local)", "78.7%", "$0"],
        ["GPT-5  (SecLLMHolmes)", "82.0%", "~$16.75 / 228"]],
       top=2.6, left=1.1, width=11.0, col_w=[4.2, 2.2, 2.2],
       highlight_col=[1],
       row_colors=[RGBColor(0x12, 0x2A, 0x20), None, None, RGBColor(0x2A, 0x14, 0x12)],
       fs=16)
_text(s, Inches(1.1), Inches(5.7), Inches(11.0), Inches(0.9),
      "The $0 model wins. GPT-5 costs ~$17 / 228 findings and doesn't lead.",
      size=18, color=FG, bold=True)
_notes(s, """
The surprising result. The zero-dollar pass-through model has the best F1, and GPT-5 — at roughly
seventeen dollars for 228 findings — doesn't lead. The lever is the protocol, not the parameter
count. Be honest about latency: this is batch-speed, tens of seconds per finding, not inline on
keystroke.
""")

# --- 14 · Ablation honesty -------------------------------------------------
s = new(); _chrome(s, "The evidence", mins="2 min")
_title(s, "Ablation honesty — where it doesn't dominate")
_bullets(s, [
    "Biggest jump: raw-SAST → ANY multi-turn LLM  (~20 F1 points)",
    "Zero-shot is already surprisingly strong",
    "Guided questions add the hard-case tail  (~3–5 pts recall)",
    ("On synthetic Juliet, generic questions even edge out specific ones", 1),
    "Saying where you DON'T win is what earns the room's trust",
], top=2.3, gap=15)
_notes(s, """
Don't oversell. The largest single jump is from raw-SAST to any multi-turn LLM — about 20 F1
points. Zero-shot is strong; the guided questions mainly add the hard-case tail, a few points of
recall, and on synthetic Juliet generic questions even nudge ahead. Telling a DefCon audience the
limits of your own method buys credibility for the numbers that ARE strong.
""")

# --- 15 · Demo setup -------------------------------------------------------
s = new(); _chrome(s, "Live demo", mins="1 min")
_title(s, "Live demo — the setup")
_bullets(s, [
    "Two targets side-by-side (mirrors examples/pipeline_python.py):",
    ("a benign real-world library", 1),
    ("a deliberately vulnerable app", 1),
    "Success criteria, stated out loud:",
    ("the real bug SURVIVES as a TP", 1),
    ("the scary-looking benign finding (slide 3) gets KILLED as an FP", 1),
    ("both with cited reasoning", 1),
], top=2.2, gap=10)
_notes(s, """
Set expectations before touching the keyboard. Two terminals: a benign upstream library and a
deliberately vulnerable app. Say the success criteria aloud so the audience knows what "working"
looks like — the real bug survives, the false positive from slide 3 dies, both with citations.
""")

# --- 16 · Demo run ---------------------------------------------------------
s = new(); _chrome(s, "Live demo", mins="6 min")
_title(s, "Live demo — run  ·  vuln-hunter-x verify")
_bullets(s, [
    "1.  Slide-3 false positive — KILLED, with the cited guard",
    "2.  A real bug — SURVIVES as TP, with a data-flow trace",
    "3.  A live multi-turn context request (caller: / struct:) → revised verdict",
    "4.  Confidence downgrade catching a thin, pattern-matched verdict",
], top=2.4, gap=18)
_text(s, Inches(0.85), Inches(5.5), Inches(11.7), Inches(1.0),
      [[("BACKUP (rehearse): ", {"color": BAD, "bold": True, "size": 17}),
        ("pre-recorded screencast + cached verdict JSON — floor network WILL flake", {"size": 17, "color": MUTED})]])
_notes(s, """
Run verify on the pre-staged findings and narrate four moments in order: the false positive dies
with its cited guard; a real bug survives with a data-flow trace; a live multi-turn context
request and the revised verdict; and the confidence downgrade catching a thin verdict. REHEARSE
the fallback: a pre-recorded screencast plus cached verdict JSON. Conference-floor network failure
is a when, not an if. If running long, the fuzz-payoff slide is the safe cut.
""")

# --- 17 · Fuzz payoff ------------------------------------------------------
s = new(); _chrome(s, "Live demo", mins="2 min")
_title(s, "Payoff: from opinion to crashing input")
_bullets(s, [
    "For a C/C++ true positive, stages 5–8 run automatically:",
    ("generate libFuzzer harness  →  ASan crash  →  triaged crashing input", 1),
    "\"The LLM says use-after-free\"  →  \"here is the input that crashes it\"",
    "That's the difference between a triage opinion and a filed, reproducible bug",
], top=2.3, gap=16)
_text(s, Inches(0.85), Inches(5.4), Inches(11.6), Inches(0.7),
      "(Pre-recorded — fuzzing live on stage is slow.)", size=15, color=MUTED)
_notes(s, """
Close the loop. For a memory-unsafe true positive, the fuzz stages generate a harness, build it
with sanitizers, and produce an actual crashing input. That converts an LLM opinion into a filed,
reproducible bug. Use the recording — fuzzing in real time on stage is too slow. This is the bridge
into the real-world results: the next segment shows five such bugs we actually found and confirmed.
""")

# ===========================================================================
# Section B · real-world results (5 confirmed bugs)
# Data source: github.com/tuonglnc/VHX-real-bug-confirmed
# ===========================================================================

# --- 18 · Section divider --------------------------------------------------
s = new(); _chrome(s, "Real-world results", tag="manually confirmed · PoC + vendor", repo=BUGS_REPO)
_text(s, Inches(0.9), Inches(2.0), Inches(11.6), Inches(1.6),
      [[("5 Confirmed Real Bugs ", {"bold": True, "size": 44}),
        ("in the Wild", {"bold": True, "size": 44, "color": ACCENT})]], line_spacing=1.0)
_text(s, Inches(0.9), Inches(3.5), Inches(11.6), Inches(0.8),
      "flac · libevent · vorbis · gradio · safety — all vendor-engaged", size=23, color=ACCENT2)
_box(s, Inches(0.9), Inches(4.55), Inches(9.2), Inches(1.15), PANEL)
_text(s, Inches(1.15), Inches(4.68), Inches(8.8), Inches(1.0),
      [[("175", {"color": ACCENT, "bold": True, "size": 26}),
        (" findings triaged  →  ", {"size": 18}),
        ("5", {"color": GOOD, "bold": True, "size": 26}),
        (" confirmed real bugs", {"size": 18}),
        ("   ·   model GLM-5.1   ·   ", {"size": 16, "color": MUTED}),
        ("$0.00", {"color": GOOD, "bold": True, "size": 18})]],
      anchor=MSO_ANCHOR.MIDDLE)
_notes(s, """
This segment is the proof: VulnHunterX wasn't just benchmarked on synthetic data — it surfaced and
helped confirm five real vulnerabilities in widely-used open-source libraries, every one of them
reported to and engaged by the vendor. The triage that found them ran on a low-cost model
(GLM-5.1) at zero marginal cost. Each bug here has a manual confirmation write-up and a working PoC
in the public repo.
""")

# --- 19 · At a glance ------------------------------------------------------
s = new(); _chrome(s, "Results", tag="manually confirmed · PoC + vendor", repo=BUGS_REPO)
_title(s, "Five confirmed real-world vulnerabilities", size=32)
_table(s,
       ["Library", "Vulnerability", "CWE", "Location", "Vendor status"],
       [["flac", "TOCTOU race condition", "CWE-367", "file.c:116", "xiph/flac #902"],
        ["libevent", "Compiler-elided memset()", "CWE-14", "sha1.c:202", "advance-security triage"],
        ["vorbis", "Stack overflow — alloca() in loop", "CWE-770", "vorbisfile.c:2290", "xiph/vorbis MR !43"],
        ["gradio", "Full SSRF via user URL", "CWE-918", "image_utils.py:253", "GitHub Security triage"],
        ["safety", "Credential leak — URL sanitization", "CWE-295", "uv/command.py:72", "email triage"]],
       top=2.2, left=0.55, width=12.25, col_w=[1.5, 3.5, 1.4, 2.6, 3.0], highlight_col=[2], fs=14)
_text(s, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.7),
      "Memory-safety, crypto-hygiene, DoS, SSRF, and credential-leak — five distinct bug classes "
      "across C and Python, each with a manual PoC.", size=14, color=MUTED)
_notes(s, """
The money slide. Five distinct vulnerability classes — a TOCTOU race, a compiler dead-store that
defeats secret-wiping, an unbounded stack allocation, a server-side request forgery, and a
credential leak from a sloppy URL check — across two languages and five mature, widely-deployed
projects. Every row has a public PoC and a vendor touchpoint. Walk the table, then the next slide
shows the funnel these came out of.
""")

# --- 20 · Scan economics ---------------------------------------------------
s = new(); _chrome(s, "Results", tag="the funnel", repo=BUGS_REPO)
_title(s, "175 findings in · 5 real bugs out · $0", size=32)
_table(s,
       ["Library", "Findings", "True Pos", "False Pos", "Scan time"],
       [["flac", "94", "12", "82", "54.6 min"],
        ["gradio", "49", "24", "12  (+12 err)", "27.0 min"],
        ["vorbis", "18", "18", "0", "17.1 min"],
        ["libevent", "7", "4", "3", "5.3 min"],
        ["safety", "7", "5", "1", "7.2 min"],
        ["TOTAL", "175", "63", "98", "~1 h 51 m"]],
       top=2.15, left=0.55, width=12.25, col_w=[2.0, 1.8, 1.8, 2.2, 2.0], highlight_col=[2], fs=14)
_text(s, Inches(0.85), Inches(6.05), Inches(11.7), Inches(1.0),
      [[("flac: the real TOCTOU bug was 1 of 94 findings — the LLM filtered 82 false positives to "
         "surface it.   ", {"size": 15}),
        ("Model GLM-5.1 · ~2.7K–9.5K tokens/finding · $0.00.",
         {"size": 15, "color": GOOD, "bold": True})]], line_spacing=1.1)
_notes(s, """
This is why the triage layer matters in the real world. 175 raw findings across the five libraries;
the LLM verdicts narrowed those to a manageable set and, on flac alone, filtered 82 false positives
to leave 12 true positives — one of which was the TOCTOU bug. Note gradio's 12 errors (harness/parse
issues) — be honest about that. The whole campaign ran on GLM-5.1 at zero marginal cost, a few
thousand tokens per finding. Cheap triage, real results.
""")

# --- 21 · flac TOCTOU ------------------------------------------------------
s = new(); _chrome(s, "Deep dive · C", tag="TP · conf 0.72 · 3 turns · xiph/flac #902", repo=BUGS_REPO)
_title(s, "flac — TOCTOU race condition (CWE-367)", size=32)
_code(s, [
    ("// grabbag__file_change_stats(filename, ...)", MUTED),
    ("stat(filename, &st);     // line 100  — CHECK", FG),
    ("...                       // attacker swaps file → symlink", BAD),
    ("flac_chmod(filename, mode);  // line 116 — USE (no fd held)", ACCENT),
], top=2.2, height=1.95, size=15)
_impact(s, "Impact",
        "Local attacker controlling the parent dir wins the race → chmod arbitrary files as the "
        "elevated user (e.g. /etc/shadow). Reachable from 2 production paths. PoC won on attempt 7.",
        top=4.5)
_verdict_chip(s, "True Positive · confidence 0.72 (Medium) · 3 multi-turn iterations — the engine "
                 "pulled extra context before committing.")
_notes(s, """
stat() at line 100 and chmod() at line 116 use the same path string with no file descriptor held
between them — the classic check-then-use TOCTOU window. An attacker who controls the parent
directory swaps the file for a symlink in that window, so chmod lands on the symlink target. When
flac/metaflac runs elevated, that's permission changes on arbitrary files. Note the verdict: Medium
confidence, and it took THREE multi-turn iterations — the engine asked for more context before
committing, exactly the behavior the methodology is built for. Vendor issue xiph/flac #902.
""")

# --- 22 · vorbis alloca ----------------------------------------------------
s = new(); _chrome(s, "Deep dive · C", tag="TP · conf 0.85 · 1 turn · xiph/vorbis MR !43", repo=BUGS_REPO)
_title(s, "vorbis — stack overflow via alloca() in loop (CWE-770)", size=32)
_code(s, [
    ("// ov_crosslap() — lib/vorbisfile.c", MUTED),
    ("for (i = 0; i < vi1->channels; i++)        // channels from file header", FG),
    ("    lappcm[i] = alloca(sizeof(**lappcm) * n1);  // line 2290", ACCENT),
    ("// channels only checked >= 1 (no upper bound); blocksize <= 8192", BAD),
], top=2.2, height=1.95, size=14)
_impact(s, "Impact",
        "channels=255 × blocksize=8192 → ~4.17 MB stack on a crafted .ogg → overflows a 1 MB worker "
        "thread → DoS. Media players / game engines decode audio on small-stack threads.",
        top=4.5)
_verdict_chip(s, "True Positive · confidence 0.85 (High) · 1 iteration · ~2.7K tokens · $0.")
_notes(s, """
ov_crosslap allocates with alloca() inside a loop bounded by the channel count, which is parsed
straight from the Ogg header and only checked to be >= 1 — no upper bound. A crafted file with 255
channels and max blocksize forces ~4 MB of stack, overflowing the small worker-thread stacks that
media players and game engines use for audio decode. High-confidence true positive in a single turn.
Vendor merge request !43 on the xiph vorbis GitLab.
""")

# --- 23 · libevent memset --------------------------------------------------
s = new(); _chrome(s, "Deep dive · C", tag="TP · conf 0.90 · 1 turn · advance-security triage", repo=BUGS_REPO)
_title(s, "libevent — compiler-elided memset() (CWE-14)", size=32)
_code(s, [
    ("// sha1.c — clearing sensitive state before return", MUTED),
    ("memset(block,  0, sizeof(block));     // L202  raw SHA1 input", ACCENT),
    ("memset(context,0, sizeof(*context));  // L287  hash state", ACCENT),
    ("// buffers never read again → Dead Store Elimination removes all 3 at -O1+", BAD),
], top=2.2, height=1.95, size=14)
_impact(s, "Impact",
        "SHA1 input blocks, hash state, and counters persist on the stack after return. A memory-read "
        "primitive (core dump, /proc/pid/mem, swap) recovers secrets. Elevated in nginx / memcached / Tor.",
        top=4.5)
_verdict_chip(s, "True Positive · confidence 0.90 (High) · 1 iteration. PoC: objdump shows zero memset "
                 "in the -O2 binary; stack remnant recovered.")
_notes(s, """
Three memset() calls meant to wipe sensitive SHA1 data are silently removed by the compiler's
dead-store elimination at -O1 and above, because the buffers are never read after the wipe. Raw
input, hash state, and counters survive on the stack. The PoC proves it: objdump shows no memset in
the optimized binary, and the secret block is recovered from the stack. Reachable through the public
SHA1Update API; risk is highest in long-lived servers built on libevent. Triaged by advance-security.
""")

# --- 24 · gradio SSRF ------------------------------------------------------
s = new(); _chrome(s, "Deep dive · Python", tag="TP · conf 0.92 · 1 turn · GitHub Security triage", repo=BUGS_REPO)
_title(s, "gradio — full SSRF via user-controlled URL (CWE-918)", size=32)
_code(s, [
    ("# gradio/image_utils.py — extract_svg_content()", MUTED),
    ("if is_http_url_like(image_file):", FG),
    ("    response = httpx.get(image_file)   # line 253 — no SSRF guard", ACCENT),
    ("# sibling processing_utils.py uses safehttpx — this path was overlooked", BAD),
], top=2.2, height=1.95, size=14)
_impact(s, "Impact",
        "A user URL ending in .svg passed to a Gallery/Image component → server fetches any address, "
        "including internal IPs and cloud-metadata endpoints (169.254.169.254). Confirmed on Gradio 6.9.0.",
        top=4.5)
_verdict_chip(s, "True Positive · confidence 0.92 (High) · 1 iteration. The LLM noted the sibling "
                 "module already uses safehttpx — this one didn't.")
_notes(s, """
extract_svg_content fetches a user-controlled URL with httpx.get and no SSRF protection — no scheme
restriction, no host allowlist. The telling detail the LLM surfaced: the sibling processing_utils.py
already routes through safehttpx, so this is an overlooked path, not a missing capability. Impact is
classic SSRF — internal services and cloud metadata. Highest-confidence verdict in the set, single
turn. Triaged privately via GitHub Security.
""")

# --- 25 · safety credential leak -------------------------------------------
s = new(); _chrome(s, "Deep dive · Python", tag="TP · conf 0.70 · 1 turn · email triage", repo=BUGS_REPO)
_title(s, "safety — credential leak via URL sanitization (CWE-295)", size=32)
_code(s, [
    ("# safety/tool/uv/command.py — before()  (L72)", MUTED),
    ('if index_value.startswith("https://pkgs.safetycli.com"):  # prefix only', ACCENT),
    ("    self.__index_url = index_value", FG),
    ('# bypass: https://pkgs.safetycli.com.attacker.com  → JWT injected into netloc', BAD),
], top=2.2, height=1.95, size=13.5)
_impact(s, "Impact",
        "Bypassable prefix check → live JWT embedded into attacker URL → leaked via Authorization header "
        "(uv request). CI/CD argument injection harvests ~24 h tokens for the Safety Platform.",
        top=4.5)
_verdict_chip(s, "True Positive · confidence 0.70 (Medium). Fix: validate urlsplit().hostname == "
                 "exact host and block userinfo.")
_notes(s, """
The --index-url is validated with startswith() on a domain prefix — which a subdomain like
pkgs.safetycli.com.attacker.com trivially passes. The CLI then embeds the user's JWT into the URL,
and uv ships it in an Authorization header to the attacker's host. The vector is CI/CD argument
injection; tokens are valid ~24 hours. The fix is exact-hostname validation via urlsplit. Medium
confidence — honest about the truncated context. Triaged with the vendor over email.
""")

# --- 26 · Why these are hard for SAST alone --------------------------------
s = new(); _chrome(s, "Why it matters", repo=BUGS_REPO)
_title(s, "Why a triage layer made the difference", size=32)
_text(s, Inches(0.85), Inches(2.2), Inches(11.8), Inches(4.0),
      [[("—  ", {"color": ACCENT2}), ("Needle in a haystack: ", {"bold": True}),
        ("the flac TOCTOU was 1 real bug among 94 findings (82 FPs). The verdict + reasoning "
         "surfaced it.", {})],
       [("—  ", {"color": ACCENT2}), ("Reasoning, not pattern-matching: ", {"bold": True}),
        ("the LLM noted gradio's sibling module already uses safehttpx, and that vorbis channels "
         "lacks an upper bound — facts that separate real from benign.", {})],
       [("—  ", {"color": ACCENT2}), ("Multi-turn earns its keep: ", {"bold": True}),
        ("flac needed 3 iterations of context expansion before committing to a verdict.", {})],
       [("—  ", {"color": ACCENT2}), ("Calibrated confidence: ", {"bold": True}),
        ("0.70–0.92 across the set — Medium where context was truncated (flac, safety), High where "
         "the flow was fully visible (gradio, libevent).", {})],
       [("—  ", {"color": ACCENT2}), ("Evidence to act on: ", {"bold": True}),
        ("each verdict shipped a data-flow trace that fed straight into a working PoC.", {})]],
      size=18, line_spacing=1.08, space_after=14)
_notes(s, """
Tie it back to the method. These bugs are exactly the cases raw SAST buries: a single real finding
inside dozens of false positives. What surfaced them was the LLM's reasoning — noticing that a
sibling module already had the safe path, that an attacker-controlled count had no upper bound — and
the multi-turn context expansion (flac took three turns). The confidence scores are calibrated:
Medium where the snippet was truncated, High where the data flow was fully visible. And every
verdict carried a data-flow trace the researcher turned into a PoC.
""")

# ===========================================================================
# Section C · wrap-up
# ===========================================================================

# --- 27 · Limits -----------------------------------------------------------
s = new(); _chrome(s, "Limits & red-team", mins="2.5 min")
_title(s, "Limits & failure modes (be candid)")
_bullets(s, [
    "Snippet-blind spots — e.g. weak-hash judged from a snippet → low recall on that class",
    "Context-heavy classes like XXE are inherently hard",
    "NMD when context is missing — it declines rather than guesses",
    "Batch latency, not real-time",
    "Prompt-injection: the analyzed code is untrusted input to the model",
], top=2.3, gap=14)
_notes(s, """
Candor is a feature for this audience. Call out the real failure modes: classes that can't be
judged from a snippet (weak-hash recall drops), context-heavy classes like XXE, the Needs-More-Data
outcome when context is missing, batch latency rather than real-time, and prompt-injection risk
because the analyzed source is untrusted input. Naming these makes the strong numbers believable.
""")

# --- 28 · Red-team ---------------------------------------------------------
s = new(); _chrome(s, "Limits & red-team", mins="2 min")
_title(s, "Red-team the tool itself")
_bullets(s, [
    "Can crafted code / comments talk the model out of a real verdict?  In principle, yes.",
    "Mitigations:",
    ("evidence-binding — a verdict must cite real file:line, not prose", 1),
    ("confidence downgrade for unsupported reasoning", 1),
    ("second-opinion re-audit on suspicious high-confidence FPs", 1),
    "Treat verdicts as triage assistance, not ground truth — especially on untrusted code",
], top=2.2, gap=12)
_notes(s, """
Invite the audience to attack it. Yes, an attacker could try to poison triage via crafted code or
comments. The defenses are structural: verdicts must cite real file:line locations, unsupported
reasoning is downgraded, and suspicious high-confidence false positives get a second-opinion
re-audit. But the honest framing is: triage assistance, not an oracle — this is the research
frontier, come break it.
""")

# --- 29 · Extensibility ----------------------------------------------------
s = new(); _chrome(s, "Wrap-up", mins="1.5 min")
_title(s, "Extensibility — bring your own")
_bullets(s, [
    "SARIF is the only contract",
    "59 custom CodeQL queries  +  89 custom Semgrep rules ship in-repo",
    "New rules tag a CWE → route through the same verification core",
    "audit_rule_coverage.py --fail-on-gaps  fails CI if a rule isn't wired to a question",
    "No core changes to add a rule, engine, or language",
], top=2.3, gap=15)
_notes(s, """
Because SARIF is the only contract, you bring your own rules, engine, or language. The repo ships
59 custom CodeQL queries and 89 custom Semgrep rules; a new rule just tags a CWE and routes through
the same verification core, and an audit script fails CI when a rule isn't wired to a guided
question. No core code changes required.
""")

# --- 30 · Takeaways --------------------------------------------------------
s = new(); _chrome(s, "Wrap-up", mins="1.5 min")
_title(s, "Takeaways")
_text(s, Inches(0.85), Inches(2.3), Inches(11.7), Inches(3.1),
      [[("1   ", {"color": ACCENT, "bold": True, "size": 25}),
        ("Reasoning STRUCTURE beats model SIZE", {"size": 22, "bold": True})],
       [("2   ", {"color": ACCENT, "bold": True, "size": 25}),
        ("FP-reduction is the ROI security teams actually feel", {"size": 22, "bold": True})],
       [("3   ", {"color": ACCENT, "bold": True, "size": 25}),
        ("Open source (MIT) + full benchmark harness — reproduce every number", {"size": 22, "bold": True})],
       [("4   ", {"color": ACCENT, "bold": True, "size": 25}),
        ("Validated in the wild — 5 confirmed bugs (flac · libevent · vorbis · gradio · safety), "
         "all vendor-engaged, on a $0 model", {"size": 22, "bold": True})]],
      space_after=16, line_spacing=1.04)
_box(s, Inches(0.85), Inches(5.85), Inches(9.6), Inches(1.1), PANEL)
_text(s, Inches(1.1), Inches(5.85), Inches(9.2), Inches(1.1),
      [[("talk   ", {"color": MUTED, "size": 14}), (TALK_REPO, {"color": ACCENT2, "bold": True, "size": 18})],
       [("bugs  ", {"color": MUTED, "size": 14}), (BUGS_REPO, {"color": ACCENT2, "bold": True, "size": 18})]],
      anchor=MSO_ANCHOR.MIDDLE, space_after=4, line_spacing=1.05)
_notes(s, """
Land the takeaways and restate the thesis: structure beats size; FP-reduction is the metric teams
feel; it's open source with a full benchmark harness so every number on the evidence slides is
reproducible; and it holds OUTSIDE the lab — five real, vendor-engaged bugs surfaced on a $0 model.
Two repos on screen: the tool, and the public real-bug confirmations + PoCs. Invite contributions,
then take questions — people screenshot.
""")

# --- 31 · Appendix ---------------------------------------------------------
s = new(); _chrome(s, "Backup", mins="Q&A only")
_title(s, "Appendix / backup")
_bullets(s, [
    "Per-CWE breakdown tables",
    "Token & cost math (~10K tokens/finding)",
    "Architecture deep-dive — context-extraction flow",
    "Exact run_model_matrix.py reproduction commands",
    "Confidence-calibration plots",
    "Real-bug confirmations + PoCs — github.com/tuonglnc/VHX-real-bug-confirmed",
], top=2.3, gap=14)
_notes(s, """
Backup slides for Q&A only — don't present linearly. Pull these up on demand: per-CWE tables, the
token/cost math, the context-extraction architecture, the exact reproduction commands, the
confidence-calibration data, and the public repo with the full real-bug confirmation write-ups and
PoCs.
""")

# ---------------------------------------------------------------------------
out = Path(__file__).resolve().parent / "defcon-vulnhunterx.pptx"
prs.save(str(out))
print(f"Wrote {out}  ({len(prs.slides._sldIdLst)} slides)")
